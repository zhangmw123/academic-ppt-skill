"""Normalize heterogeneous research materials through one ingestion interface."""

from __future__ import annotations

import csv
import hashlib
import mimetypes
from pathlib import Path

from .models import SourceBlock, SourceDocument


ROLE_VALUES = {
    "task_constraint",
    "primary_evidence",
    "incremental_evidence",
    "historical_expression",
    "visual_asset",
    "template",
}


def file_sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def infer_role(path: Path) -> str:
    name = path.stem.lower()
    if any(token in name for token in ("模板", "template")):
        return "template"
    if any(token in name for token in ("要求", "通知", "规范", "brief")):
        return "task_constraint"
    if path.suffix.lower() in {".xlsx", ".xls", ".csv"}:
        return "incremental_evidence"
    if path.suffix.lower() in {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".svg"}:
        return "visual_asset"
    if any(token in name for token in ("旧", "历史", "previous")) and path.suffix.lower() == ".pptx":
        return "historical_expression"
    return "primary_evidence"


class SourceIngestor:
    """Deep adapter that hides format-specific extraction behind ``ingest``."""

    def ingest(self, path: Path | str, role: str | None = None) -> SourceDocument:
        source_path = Path(path).resolve()
        if not source_path.is_file():
            raise FileNotFoundError(source_path)
        resolved_role = role or infer_role(source_path)
        if resolved_role not in ROLE_VALUES:
            raise ValueError(f"unsupported source role: {resolved_role}")
        digest = file_sha256(source_path)
        suffix = source_path.suffix.lower()
        image_metadata = {}
        if suffix in {".txt", ".md", ".markdown"}:
            blocks = self._text_blocks(source_path)
            format_name = "markdown" if suffix in {".md", ".markdown"} else "text"
        elif suffix == ".pdf":
            blocks = self._pdf_blocks(source_path)
            format_name = "pdf"
        elif suffix == ".csv":
            blocks = self._csv_blocks(source_path)
            format_name = "csv"
        elif suffix == ".xlsx":
            blocks = self._xlsx_blocks(source_path)
            format_name = "xlsx"
        elif suffix == ".docx":
            blocks = self._docx_blocks(source_path)
            format_name = "docx"
        elif suffix == ".pptx":
            blocks = self._pptx_blocks(source_path)
            format_name = "pptx"
        elif suffix in {".bmp", ".gif", ".jpeg", ".jpg", ".png", ".tif", ".tiff", ".webp"}:
            blocks, image_metadata = self._image_blocks(source_path)
            format_name = "image"
        else:
            raise ValueError(f"unsupported source format: {suffix or '<none>'}")
        return SourceDocument(
            source_id=f"SRC_{digest[:12].upper()}",
            path=str(source_path),
            format=format_name,
            role=resolved_role,
            sha256=digest,
            blocks=blocks,
            metadata={"mime_type": mimetypes.guess_type(source_path.name)[0], **image_metadata},
        )

    @staticmethod
    def _text_blocks(path: Path) -> list[SourceBlock]:
        blocks = []
        for line_number, line in enumerate(path.read_text(encoding="utf-8-sig").splitlines(), 1):
            text = line.strip()
            if not text:
                continue
            kind = "heading" if text.startswith("#") else "text"
            blocks.append(SourceBlock(
                block_id=f"B{len(blocks) + 1:04d}",
                kind=kind,
                text=text,
                locator={"line": line_number},
            ))
        return blocks

    @staticmethod
    def _pdf_blocks(path: Path) -> list[SourceBlock]:
        try:
            import fitz
        except ImportError as exc:
            raise RuntimeError("PyMuPDF is required to ingest PDF sources") from exc
        document = fitz.open(path)
        blocks = []
        try:
            for page_number, page in enumerate(document, 1):
                for block_number, raw_block in enumerate(page.get_text("blocks"), 1):
                    text = raw_block[4].strip()
                    if not text:
                        continue
                    blocks.append(SourceBlock(
                        block_id=f"B{len(blocks) + 1:04d}",
                        kind="text",
                        text=text,
                        locator={"page": page_number, "block": block_number},
                        data={"bbox": list(raw_block[:4])},
                    ))
                for image_number, image_info in enumerate(page.get_images(full=True), 1):
                    image = document.extract_image(image_info[0])
                    image_bytes = image["image"]
                    blocks.append(SourceBlock(
                        block_id=f"B{len(blocks) + 1:04d}",
                        kind="picture",
                        locator={"page": page_number, "image": image_number, "xref": image_info[0]},
                        data={
                            "image_size": [image["width"], image["height"]],
                            "image_format": image["ext"],
                            "image_sha256": hashlib.sha256(image_bytes).hexdigest(),
                        },
                    ))
        finally:
            document.close()
        return blocks

    @staticmethod
    def _image_blocks(path: Path) -> tuple[list[SourceBlock], dict]:
        try:
            from PIL import Image
        except ImportError as exc:
            raise RuntimeError("Pillow is required to ingest image sources") from exc
        with Image.open(path) as image:
            metadata = {
                "image_size": [image.width, image.height],
                "image_mode": image.mode,
                "image_format": (image.format or path.suffix.lstrip(".")).lower(),
            }
        return [SourceBlock(
            block_id="B0001",
            kind="picture",
            locator={"image": 1},
            data={**metadata, "image_sha256": file_sha256(path)},
        )], metadata

    @staticmethod
    def _csv_blocks(path: Path) -> list[SourceBlock]:
        with path.open("r", encoding="utf-8-sig", newline="") as handle:
            rows = list(csv.reader(handle))
        return [SourceBlock(
            block_id="B0001",
            kind="table",
            locator={"row_start": 1, "row_end": len(rows)},
            data={"rows": rows},
        )]

    @staticmethod
    def _xlsx_blocks(path: Path) -> list[SourceBlock]:
        try:
            from openpyxl import load_workbook
            from openpyxl.utils import get_column_letter
        except ImportError as exc:
            raise RuntimeError("openpyxl is required to ingest XLSX sources") from exc
        workbook = load_workbook(path, data_only=False, read_only=False)
        blocks = []
        for sheet in workbook.worksheets:
            if sheet.max_row == 1 and sheet.max_column == 1 and sheet["A1"].value is None:
                continue
            rows = []
            formulas = {}
            for row in sheet.iter_rows(min_row=1, max_row=sheet.max_row, max_col=sheet.max_column):
                values = []
                for cell in row:
                    values.append(cell.value)
                    if cell.data_type == "f":
                        formulas[cell.coordinate] = cell.value
                rows.append(values)
            blocks.append(SourceBlock(
                block_id=f"B{len(blocks) + 1:04d}",
                kind="table",
                locator={
                    "sheet": sheet.title,
                    "range": f"A1:{get_column_letter(sheet.max_column)}{sheet.max_row}",
                },
                data={"rows": rows, "formulas": formulas},
            ))
        return blocks

    @staticmethod
    def _docx_blocks(path: Path) -> list[SourceBlock]:
        try:
            from docx import Document
            from docx.table import Table
            from docx.text.paragraph import Paragraph
            from docx.oxml.table import CT_Tbl
            from docx.oxml.text.paragraph import CT_P
        except ImportError as exc:
            raise RuntimeError("python-docx is required to ingest DOCX sources") from exc
        document = Document(path)
        paragraph_index = 0
        table_index = 0
        heading_path: list[str] = []
        blocks = []
        for child in document.element.body.iterchildren():
            if isinstance(child, CT_P):
                paragraph_index += 1
                paragraph = Paragraph(child, document)
                text = paragraph.text.strip()
                xml = paragraph._p.xml
                has_equation = "<m:oMath" in xml
                has_image = "<a:blip" in xml
                if not text and not has_equation and not has_image:
                    continue
                style_name = paragraph.style.name if paragraph.style else ""
                is_heading = style_name.lower().startswith("heading") or style_name.startswith("标题")
                if is_heading and text:
                    digits = "".join(character for character in style_name if character.isdigit())
                    level = max(1, int(digits or 1))
                    heading_path = heading_path[:level - 1]
                    heading_path.append(text)
                blocks.append(SourceBlock(
                    block_id=f"B{len(blocks) + 1:04d}",
                    kind="heading" if is_heading else ("equation" if has_equation and not text else "text"),
                    text=text,
                    locator={"paragraph": paragraph_index, "heading_path": list(heading_path)},
                    data={"contains_equation": has_equation, "contains_image": has_image},
                ))
            elif isinstance(child, CT_Tbl):
                table_index += 1
                table = Table(child, document)
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                blocks.append(SourceBlock(
                    block_id=f"B{len(blocks) + 1:04d}",
                    kind="table",
                    locator={"table": table_index, "heading_path": list(heading_path)},
                    data={"rows": rows},
                ))
        return blocks

    @staticmethod
    def _pptx_blocks(path: Path) -> list[SourceBlock]:
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError as exc:
            raise RuntimeError("python-pptx is required to ingest PPTX sources") from exc
        presentation = Presentation(path)
        blocks = []

        def add_shape(slide_number, shape, parent_group_shape_id=None):
            locator = {"slide": slide_number, "shape_id": shape.shape_id}
            if parent_group_shape_id is not None:
                locator["parent_group_shape_id"] = parent_group_shape_id
            geometry = {
                "left": int(shape.left),
                "top": int(shape.top),
                "width": int(shape.width),
                "height": int(shape.height),
            }
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                blocks.append(SourceBlock(
                    block_id=f"B{len(blocks) + 1:04d}",
                    kind="group",
                    locator=locator,
                    data={"geometry": geometry, "child_count": len(shape.shapes)},
                ))
                for child in shape.shapes:
                    add_shape(slide_number, child, shape.shape_id)
                return
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                blocks.append(SourceBlock(
                    block_id=f"B{len(blocks) + 1:04d}",
                    kind="table",
                    locator=locator,
                    data={"rows": rows, "geometry": geometry},
                ))
                return
            if getattr(shape, "has_chart", False):
                blocks.append(SourceBlock(
                    block_id=f"B{len(blocks) + 1:04d}",
                    kind="chart",
                    locator=locator,
                    data={"chart_type": str(shape.chart.chart_type), "geometry": geometry},
                ))
                return
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                blocks.append(SourceBlock(
                    block_id=f"B{len(blocks) + 1:04d}",
                    kind="picture",
                    locator=locator,
                    data={
                        "geometry": geometry,
                        "image_size": list(image.size),
                        "image_format": image.ext,
                        "image_sha256": hashlib.sha256(image.blob).hexdigest(),
                    },
                ))
                return
            text = shape.text.strip() if getattr(shape, "has_text_frame", False) else ""
            if text:
                blocks.append(SourceBlock(
                    block_id=f"B{len(blocks) + 1:04d}",
                    kind="text",
                    text=text,
                    locator=locator,
                    data={"geometry": geometry, "shape_type": str(shape.shape_type)},
                ))

        for slide_number, slide in enumerate(presentation.slides, 1):
            for shape in slide.shapes:
                add_shape(slide_number, shape)
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                blocks.append(SourceBlock(
                    block_id=f"B{len(blocks) + 1:04d}",
                    kind="speaker_notes",
                    text=notes_text,
                    locator={"slide": slide_number, "notes": True},
                ))
        return blocks
