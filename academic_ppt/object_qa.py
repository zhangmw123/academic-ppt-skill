"""Object-level QA for semantic template ownership and rendered PPTX bindings."""

from __future__ import annotations

import hashlib
import json
import math
import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .template_semantics import EXCLUSIVE_RENDER_MODES, MEDIA_PROVENANCE_FIELDS


INTERNAL_LABEL_RE = re.compile(
    r"核心判断|证据依据|机制解释|边界(?:与|和)行动|读图重点|读图结论|证据、读图与边界"
)


@dataclass(frozen=True)
class ObjectQualityResult:
    passed: bool
    errors: tuple[str, ...]
    observations: tuple[str, ...]
    categories: dict[str, tuple[str, ...]]
    inspected_slide_count: int

    def to_dict(self) -> dict:
        return {
            "passed": self.passed,
            "errors": list(self.errors),
            "observations": list(self.observations),
            "categories": {key: list(values) for key, values in self.categories.items()},
            "inspected_slide_count": self.inspected_slide_count,
        }


@dataclass(frozen=True)
class IdentityDifferenceResult:
    passed: bool
    errors: tuple[str, ...]
    signatures: dict[str, str]


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _normalized_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return " ".join(shape.text.split())


def _shape_kind(shape) -> str:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    if getattr(shape, "has_text_frame", False):
        return "text"
    return str(shape.shape_type)


def _shape_payload(shape) -> str | None:
    text = _normalized_text(shape)
    if text:
        return "text:" + text
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            return "picture:" + hashlib.sha256(shape.image.blob).hexdigest()
        except Exception:
            return "picture:unreadable"
    if getattr(shape, "has_table", False):
        values = [cell.text for row in shape.table.rows for cell in row.cells]
        return "table:" + hashlib.sha256("\n".join(values).encode("utf-8")).hexdigest()
    if getattr(shape, "has_chart", False):
        return "chart"
    return None


def _geometry(shape) -> tuple[int, int, int, int]:
    return int(shape.left), int(shape.top), int(shape.width), int(shape.height)


def _overlap_ratio(first, second) -> float:
    left = max(first.left, second.left)
    top = max(first.top, second.top)
    right = min(first.left + first.width, second.left + second.width)
    bottom = min(first.top + first.height, second.top + second.height)
    if right <= left or bottom <= top:
        return 0.0
    overlap = (right - left) * (bottom - top)
    smaller = min(first.width * first.height, second.width * second.height)
    return overlap / max(smaller, 1)


def _font_sizes(shape) -> list[float]:
    values = []
    if not getattr(shape, "has_text_frame", False):
        return values
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.font.size:
            values.append(float(paragraph.font.size.pt))
        values.extend(float(run.font.size.pt) for run in paragraph.runs if run.font.size)
    return values


def _read_json(value: dict | Path | str | None) -> dict | None:
    if value is None or isinstance(value, dict):
        return value
    return json.loads(Path(value).read_text(encoding="utf-8"))


class ObjectLevelQualityGate:
    """Inspect actual PPTX objects against semantic ownership and render bindings."""

    CATEGORY_NAMES = (
        "duplicate_objects",
        "cross_module_overlaps",
        "orphan_components",
        "empty_media_slots",
        "template_residue",
        "font_bounds",
        "template_identity",
        "media_provenance",
        "exclusive_render_modes",
    )

    def inspect(
        self,
        pptx_path: Path | str,
        semantic_specification: dict | Path | str,
        *,
        render_manifest: dict | Path | str | None = None,
        asset_root: Path | str | None = None,
    ) -> ObjectQualityResult:
        specification = _read_json(semantic_specification) or {}
        manifest = _read_json(render_manifest)
        presentation = Presentation(pptx_path)
        categories: dict[str, list[str]] = {name: [] for name in self.CATEGORY_NAMES}
        observations: list[str] = []
        pages_by_id = {page["page_id"]: page for page in specification.get("pages", ())}
        manifest_pages = {
            int(page["slide_number"]): page for page in manifest.get("pages", ())
        } if manifest else {}
        if manifest:
            for error in manifest.get("errors", ()):
                categories["orphan_components"].append(f"binding manifest: {error}")

        if manifest:
            selected = []
            for slide_number in range(1, len(presentation.slides) + 1):
                manifest_page = manifest_pages.get(slide_number)
                if not manifest_page:
                    categories["orphan_components"].append(
                        f"slide {slide_number}: render manifest page is missing"
                    )
                    selected.append((slide_number, None, None))
                    continue
                template_page_id = manifest_page.get("template_page_id")
                selected.append((slide_number, pages_by_id.get(template_page_id), manifest_page))
                if template_page_id not in pages_by_id:
                    categories["orphan_components"].append(
                        f"slide {slide_number}: unknown template page {template_page_id!r}"
                    )
        else:
            selected = [
                (index, page, None)
                for index, page in enumerate(specification.get("pages", ()), 1)
                if index <= len(presentation.slides)
            ]
            if len(presentation.slides) != len(specification.get("pages", ())):
                categories["orphan_components"].append(
                    f"PPTX has {len(presentation.slides)} slides; semantic specification has {len(specification.get('pages', ()))}"
                )

        for slide_number, page, manifest_page in selected:
            if page is None or slide_number > len(presentation.slides):
                continue
            slide = presentation.slides[slide_number - 1]
            shapes = list(_iter_shapes(slide.shapes))
            shape_by_id = {int(shape.shape_id): shape for shape in shapes}
            native_identity_shape_ids = {
                int(shape_id)
                for region in page.get("identity_regions", ())
                for shape_id in region.get("ownership_group", {}).get("shape_ids", ())
            }
            blank_reconstruction = bool(
                manifest_page and manifest_page.get("base_mode") == "blank_reconstruction"
            )
            additional_identity_shape_ids = {
                int(value) for value in manifest_page.get("additional_identity_shape_ids", ())
            } if manifest_page else set()
            identity_shape_ids = (
                additional_identity_shape_ids
                if blank_reconstruction
                else native_identity_shape_ids | additional_identity_shape_ids
            )
            self._check_duplicates(slide_number, shapes, identity_shape_ids, categories)
            # Native template IDs and newly rendered IDs use independent
            # namespaces after a blank reconstruction. Seed only retained
            # source slides with native ownership; the manifest fills the
            # rendered ownership map below.
            module_map = {} if blank_reconstruction else self._module_shape_map(page)
            regions = {
                region.get("module_id"): region
                for region in manifest_page.get("regions", ())
            } if manifest_page else {}
            if manifest_page:
                self._check_manifest_regions(
                    slide_number,
                    page,
                    regions,
                    shape_by_id,
                    module_map,
                    categories,
                    source_components_present=bool(
                        manifest_page.get("source_components_present", not blank_reconstruction)
                    ),
                )
            else:
                self._check_standard_ownership(slide_number, page, shape_by_id, categories)
            self._check_identity(slide_number, identity_shape_ids, shape_by_id, categories)
            if blank_reconstruction:
                required_features = int(
                    specification.get("template_identity", {}).get("minimum_preserved_features", 4)
                )
                feature_bindings = manifest_page.get("identity_feature_bindings", ())
                if len(feature_bindings) < required_features:
                    categories["template_identity"].append(
                        f"slide {slide_number}: blank reconstruction binds only {len(feature_bindings)} "
                        f"template identity features; require {required_features}"
                    )
                for binding in feature_bindings:
                    missing = sorted(
                        int(value) for value in binding.get("shape_ids", ())
                        if int(value) not in shape_by_id
                    )
                    if missing:
                        categories["template_identity"].append(
                            f"slide {slide_number}: identity feature {binding.get('feature')!r} "
                            f"references missing shapes {missing}"
                        )
                if page.get("page_role") not in {"cover", "ending", "section"}:
                    for shape_id in identity_shape_ids:
                        shape = shape_by_id.get(shape_id)
                        if not shape or shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                            continue
                        area_ratio = (shape.width * shape.height) / max(
                            presentation.slide_width * presentation.slide_height, 1
                        )
                        center_y = (shape.top + shape.height / 2) / presentation.slide_height
                        if area_ratio > 0.12 and center_y > 0.16:
                            categories["template_residue"].append(
                                f"slide {slide_number}: content-sized identity picture {shape_id} "
                                f"occupies {area_ratio:.0%} of the slide"
                            )
            self._check_frame_layers(
                slide_number,
                shapes,
                module_map,
                identity_shape_ids,
                categories,
            )
            self._check_overlaps(
                slide_number,
                shapes,
                module_map,
                identity_shape_ids,
                categories,
                include_media=manifest_page is not None,
            )
            self._check_internal_labels(slide_number, shapes, categories)
            if manifest_page:
                self._check_bindings(
                    slide_number,
                    page,
                    regions,
                    shape_by_id,
                    categories,
                    Path(asset_root).resolve() if asset_root else Path(pptx_path).resolve().parent,
                )
                self._check_unmanifested_content(
                    slide_number, shapes, page, regions, identity_shape_ids, categories
                )

        errors = tuple(message for name in self.CATEGORY_NAMES for message in categories[name])
        if not manifest:
            observations.append(
                "No render manifest supplied; empty media, replacement fonts, source provenance, and sample residue are checked only after binding."
            )
        return ObjectQualityResult(
            passed=not errors,
            errors=errors,
            observations=tuple(observations),
            categories={name: tuple(values) for name, values in categories.items()},
            inspected_slide_count=len(selected),
        )

    @staticmethod
    def _module_shape_map(page: dict) -> dict[int, str]:
        values = {}
        for module in page.get("semantic_modules", ()):
            for shape_id in module.get("ownership_group", {}).get("shape_ids", ()):
                values[int(shape_id)] = module["module_id"]
        return values

    @staticmethod
    def _check_duplicates(slide_number, shapes, identity_shape_ids, categories):
        signatures: dict[tuple, list[int]] = {}
        for shape in shapes:
            if int(shape.shape_id) in identity_shape_ids:
                continue
            payload = _shape_payload(shape)
            if payload is None:
                continue
            signature = (_shape_kind(shape), _geometry(shape), payload)
            signatures.setdefault(signature, []).append(int(shape.shape_id))
        for shape_ids in signatures.values():
            if len(shape_ids) > 1:
                categories["duplicate_objects"].append(
                    f"slide {slide_number}: duplicate content objects {shape_ids}"
                )

    @staticmethod
    def _check_overlaps(
        slide_number, shapes, module_map, identity_shape_ids, categories, *, include_media
    ):
        content_shapes = [
            shape for shape in shapes
            if int(shape.shape_id) not in identity_shape_ids
            and int(shape.shape_id) in module_map
            and _shape_payload(shape) is not None
        ]
        for index, first in enumerate(content_shapes):
            for second in content_shapes[index + 1:]:
                first_module = module_map[int(first.shape_id)]
                second_module = module_map[int(second.shape_id)]
                if first_module == second_module:
                    continue
                if not include_media and (
                    _shape_kind(first) != "text" or _shape_kind(second) != "text"
                ):
                    continue
                ratio = _overlap_ratio(first, second)
                if ratio > 0.35:
                    categories["cross_module_overlaps"].append(
                        f"slide {slide_number}: shapes {first.shape_id}/{second.shape_id} from "
                        f"{first_module}/{second_module} overlap {ratio:.0%}"
                    )

    @staticmethod
    def _check_standard_ownership(slide_number, page, shape_by_id, categories):
        for module in page.get("semantic_modules", ()):
            missing = sorted(
                int(shape_id) for shape_id in module.get("ownership_group", {}).get("shape_ids", ())
                if int(shape_id) not in shape_by_id
            )
            if missing:
                categories["orphan_components"].append(
                    f"slide {slide_number}: {module['module_id']} native shapes are missing {missing}"
                )

    @staticmethod
    def _check_identity(slide_number, identity_shape_ids, shape_by_id, categories):
        missing = sorted(identity_shape_ids - set(shape_by_id))
        if missing:
            categories["template_identity"].append(
                f"slide {slide_number}: template identity shapes are missing {missing}"
            )

    @staticmethod
    def _check_manifest_regions(
        slide_number,
        page,
        regions,
        shape_by_id,
        module_map,
        categories,
        *,
        source_components_present,
    ):
        for module in page.get("semantic_modules", ()):
            module_id = module["module_id"]
            region = regions.get(module_id)
            if not region:
                categories["exclusive_render_modes"].append(
                    f"slide {slide_number}: {module_id} has no selected render mode"
                )
                continue
            mode = region.get("render_mode")
            if mode not in EXCLUSIVE_RENDER_MODES:
                categories["exclusive_render_modes"].append(
                    f"slide {slide_number}: {module_id} has invalid render mode {mode!r}"
                )
                continue
            allowed = set(module.get("render_contract", {}).get("allowed_modes", ()))
            if mode not in allowed:
                categories["exclusive_render_modes"].append(
                    f"slide {slide_number}: {module_id} does not allow {mode}"
                )
            native_shapes = {
                int(value) for value in module.get("ownership_group", {}).get("complete_delete_shape_ids", ())
            }
            present = native_shapes & set(shape_by_id)
            if mode == "full_reconstruction" and source_components_present and present:
                categories["template_residue"].append(
                    f"slide {slide_number}: {module_id} reconstructed over retained native shapes {sorted(present)}"
                )
            if mode == "native_reuse":
                missing = native_shapes - set(shape_by_id)
                if missing:
                    categories["orphan_components"].append(
                        f"slide {slide_number}: {module_id} native reuse is missing shapes {sorted(missing)}"
                    )
            for binding in region.get("slot_bindings", ()):
                for shape_id in binding.get("shape_ids", ()):
                    module_map[int(shape_id)] = module_id
            for shape_id in region.get("owned_shape_ids", ()):
                shape_id = int(shape_id)
                previous = module_map.get(shape_id)
                if previous and previous != module_id:
                    categories["duplicate_objects"].append(
                        f"slide {slide_number}: shape {shape_id} is owned by {previous} and {module_id}"
                    )
                module_map[shape_id] = module_id
            declared_removed = {
                int(value) for value in region.get("removed_native_shape_ids", ())
            }
            if mode == "full_reconstruction" and declared_removed != native_shapes:
                missing_removed = sorted(native_shapes - declared_removed)
                extra_removed = sorted(declared_removed - native_shapes)
                categories["template_residue"].append(
                    f"slide {slide_number}: {module_id} incomplete native removal; "
                    f"missing={missing_removed} extra={extra_removed}"
                )

    @staticmethod
    def _check_internal_labels(slide_number, shapes, categories):
        for shape in shapes:
            text = _normalized_text(shape)
            if text and INTERNAL_LABEL_RE.search(text):
                categories["template_residue"].append(
                    f"slide {slide_number}: internal production label remains in shape {shape.shape_id}: {text[:50]}"
                )

    def _check_bindings(self, slide_number, page, regions, shape_by_id, categories, asset_root):
        for module in page.get("semantic_modules", ()):
            module_id = module["module_id"]
            region = regions.get(module_id)
            if not region:
                continue
            bindings = {binding.get("slot_id"): binding for binding in region.get("slot_bindings", ())}
            explicitly_reflowed = set(region.get("reflowed_or_removed_slots", ()))
            for slot in module.get("child_slots", ()):
                slot_id = slot["slot_id"]
                binding = bindings.get(slot_id)
                if slot.get("role") == "image_or_chart":
                    if binding is None:
                        if slot_id not in explicitly_reflowed:
                            categories["empty_media_slots"].append(
                                f"slide {slide_number}: media slot {slot_id} is empty and was not reflowed"
                            )
                        continue
                elif binding is None and slot.get("required") and slot_id not in explicitly_reflowed:
                    categories["orphan_components"].append(
                        f"slide {slide_number}: required text slot {slot_id} has no binding"
                    )
                if binding is None:
                    continue
                bound_shapes = [shape_by_id.get(int(value)) for value in binding.get("shape_ids", ())]
                bound_shapes = [shape for shape in bound_shapes if shape is not None]
                if slot.get("role") == "image_or_chart":
                    if not bound_shapes:
                        categories["empty_media_slots"].append(
                            f"slide {slide_number}: media slot {slot_id} has no actual bound shape"
                        )
                    self._check_media_source(
                        slide_number, slot_id, binding, bound_shapes, categories, asset_root
                    )
                else:
                    if not bound_shapes:
                        categories["orphan_components"].append(
                            f"slide {slide_number}: text slot {slot_id} has no actual bound shape"
                        )
                    self._check_font_bounds(slide_number, slot, bound_shapes, categories)
                    self._check_sample_residue(slide_number, slot, bound_shapes, categories)

    @staticmethod
    def _check_media_source(slide_number, slot_id, binding, shapes, categories, asset_root):
        source = binding.get("source", {})
        missing = [field for field in MEDIA_PROVENANCE_FIELDS if source.get(field) in {None, ""}]
        if missing:
            categories["media_provenance"].append(
                f"slide {slide_number}: {slot_id} lacks media provenance fields {missing}"
            )
            return
        pdf_page = source.get("pdf_page")
        if not (isinstance(pdf_page, int) and pdf_page > 0) and pdf_page != "not_applicable":
            categories["media_provenance"].append(
                f"slide {slide_number}: {slot_id} has invalid PDF page {pdf_page!r}"
            )
        source_path = Path(str(source["source_path"]))
        resolved = source_path if source_path.is_absolute() else asset_root / source_path
        if not resolved.is_file():
            categories["media_provenance"].append(
                f"slide {slide_number}: {slot_id} source file does not exist: {resolved}"
            )
        if resolved.suffix.casefold() == ".pdf" and not (isinstance(pdf_page, int) and pdf_page > 0):
            categories["media_provenance"].append(
                f"slide {slide_number}: {slot_id} requires a positive PDF page for {resolved}"
            )
        if INTERNAL_LABEL_RE.search(str(source.get("semantic_use", ""))):
            categories["media_provenance"].append(
                f"slide {slide_number}: {slot_id} semantic use exposes an internal production label"
            )
        picture_shapes = [shape for shape in shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
        if not picture_shapes:
            return
        asset_value = source.get("asset_path")
        if not asset_value:
            categories["media_provenance"].append(
                f"slide {slide_number}: {slot_id} picture binding lacks asset_path"
            )
            return
        asset_path = Path(str(asset_value))
        asset = asset_path if asset_path.is_absolute() else asset_root / asset_path
        if not asset.is_file():
            categories["media_provenance"].append(
                f"slide {slide_number}: {slot_id} bound image asset does not exist: {asset}"
            )
            return
        expected_digest = hashlib.sha256(asset.read_bytes()).hexdigest()
        for shape in picture_shapes:
            try:
                actual_digest = hashlib.sha256(shape.image.blob).hexdigest()
            except Exception as exc:
                categories["media_provenance"].append(
                    f"slide {slide_number}: {slot_id} picture shape {shape.shape_id} cannot be read: {exc}"
                )
                continue
            if actual_digest != expected_digest:
                categories["media_provenance"].append(
                    f"slide {slide_number}: {slot_id} picture shape {shape.shape_id} does not match {asset}"
                )

    @staticmethod
    def _check_font_bounds(slide_number, slot, shapes, categories):
        typography = slot.get("typography", {})
        minimum = float(typography.get("min_pt", 0))
        maximum = float(typography.get("max_pt", math.inf))
        for shape in shapes:
            for size in _font_sizes(shape):
                if size < minimum - 0.05 or size > maximum + 0.05:
                    categories["font_bounds"].append(
                        f"slide {slide_number}: {slot['slot_id']} shape {shape.shape_id} uses {size:g} pt; "
                        f"allowed {minimum:g}-{maximum:g} pt"
                    )

    @staticmethod
    def _check_sample_residue(slide_number, slot, shapes, categories):
        fingerprints = {item["sha256"] for item in slot.get("source_sample_fingerprints", ())}
        for shape in shapes:
            text = _normalized_text(shape)
            if text and hashlib.sha256(text.encode("utf-8")).hexdigest() in fingerprints:
                categories["template_residue"].append(
                    f"slide {slide_number}: {slot['slot_id']} still contains source-template sample text"
                )

    @staticmethod
    def _check_unmanifested_content(slide_number, shapes, page, regions, identity_shape_ids, categories):
        declared = set(identity_shape_ids)
        for module in page.get("semantic_modules", ()):
            region = regions.get(module["module_id"], {})
            if region.get("render_mode") == "native_reuse":
                declared.update(int(value) for value in module.get("ownership_group", {}).get("shape_ids", ()))
            declared.update(int(value) for value in region.get("owned_shape_ids", ()))
            for binding in region.get("slot_bindings", ()):
                declared.update(int(value) for value in binding.get("shape_ids", ()))
        orphaned = [
            int(shape.shape_id) for shape in shapes
            if int(shape.shape_id) not in declared and _shape_payload(shape) is not None
        ]
        if orphaned:
            categories["orphan_components"].append(
                f"slide {slide_number}: content-bearing shapes are outside semantic ownership {sorted(orphaned)}"
            )

    @staticmethod
    def _check_frame_layers(slide_number, shapes, module_map, identity_shape_ids, categories):
        minimum = 0.55 * 914400
        frames = [
            shape for shape in shapes
            if int(shape.shape_id) not in identity_shape_ids
            and int(shape.shape_id) in module_map
            and shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            and not _normalized_text(shape)
            and shape.width >= minimum
            and shape.height >= minimum
        ]
        for index, outer in enumerate(frames):
            for inner in frames[index + 1:]:
                if module_map[int(outer.shape_id)] != module_map[int(inner.shape_id)]:
                    continue
                first, second = (
                    (outer, inner)
                    if outer.width * outer.height >= inner.width * inner.height
                    else (inner, outer)
                )
                contains = (
                    first.left <= second.left
                    and first.top <= second.top
                    and first.left + first.width >= second.left + second.width
                    and first.top + first.height >= second.top + second.height
                )
                if contains and first.width * first.height > second.width * second.height * 1.1:
                    categories["duplicate_objects"].append(
                        f"slide {slide_number}: nested frame shapes {first.shape_id}/{second.shape_id} "
                        f"remain in {module_map[int(first.shape_id)]}"
                    )


class TemplateIdentityDifferenceGate:
    """Require standardized templates to retain distinct machine-readable identities."""

    def compare(self, specifications: list[dict | Path | str]) -> IdentityDifferenceResult:
        payloads = [_read_json(value) or {} for value in specifications]
        signatures = {
            payload.get("template", {}).get("id", f"template-{index + 1}"):
            payload.get("template_identity", {}).get("identity_signature", "")
            for index, payload in enumerate(payloads)
        }
        errors = []
        if any(not value for value in signatures.values()):
            errors.append("every template requires an identity signature")
        reverse: dict[str, list[str]] = {}
        for template_id, signature in signatures.items():
            reverse.setdefault(signature, []).append(template_id)
        for signature, template_ids in reverse.items():
            if signature and len(template_ids) > 1:
                errors.append(f"templates share one identity signature: {template_ids}")
        for index, first in enumerate(payloads):
            first_id = first.get("template", {}).get("id", str(index))
            first_features = set(first.get("template_identity", {}).get("identity_features", ()))
            for second in payloads[index + 1:]:
                second_id = second.get("template", {}).get("id", "unknown")
                second_features = set(second.get("template_identity", {}).get("identity_features", ()))
                union = first_features | second_features
                similarity = len(first_features & second_features) / max(len(union), 1)
                if similarity >= 0.8:
                    errors.append(
                        f"{first_id}/{second_id} identity features are too similar ({similarity:.0%})"
                    )
        return IdentityDifferenceResult(not errors, tuple(errors), signatures)
