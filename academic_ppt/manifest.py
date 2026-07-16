"""Build a per-slide editability, density, and template-identity manifest."""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .composition import CompositionQualityGate, text_units


class SlideManifestBuilder:
    """Measure the actual PPTX against the template-guided composition contract."""

    def build(
        self,
        *,
        pptx_path: Path | str,
        dynamic_plan: dict,
        visual_system_path: Path | str,
        template_grammar_path: Path | str,
        template: dict,
    ) -> dict:
        presentation = Presentation(pptx_path)
        visual_system = json.loads(Path(visual_system_path).read_text(encoding="utf-8"))
        grammar = json.loads(Path(template_grammar_path).read_text(encoding="utf-8"))
        pages = dynamic_plan.get("pages", ())
        errors: list[str] = []
        slides = []
        if len(presentation.slides) != len(pages):
            errors.append(f"manifest page count mismatch: pptx={len(presentation.slides)} plan={len(pages)}")
        for index, page in enumerate(pages):
            if index >= len(presentation.slides):
                break
            slide = presentation.slides[index]
            text_shapes = [
                shape for shape in slide.shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            ]
            pictures = [shape for shape in slide.shapes if getattr(shape, "shape_type", None) == 13]
            native_shapes = [shape for shape in slide.shapes if getattr(shape, "shape_type", None) != 13]
            editable_frames = [
                shape for shape in slide.shapes
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.AUTO_SHAPE
            ]
            modules = [value for value in CompositionQualityGate._modules(page) if str(value).strip()]
            density_units = round(sum(text_units(str(value)) for value in modules), 1)
            layout = page.get("layout")
            content_page = layout not in {"cover", "ending", "section", "agenda"}
            reference = page.get("template_reference")
            page_errors = []
            if content_page and not reference:
                page_errors.append("missing template archetype reference")
            if content_page and len(text_shapes) < 5:
                page_errors.append(f"only {len(text_shapes)} non-empty editable text objects")
            if content_page and density_units < 55:
                page_errors.append(f"density {density_units:g} below 55")
            if layout == "text_figure" and not pictures:
                page_errors.append("text_figure has no picture asset")
            if layout == "media_gallery" and len(pictures) < len(page.get("media_items", ())):
                page_errors.append(
                    f"media_gallery has {len(pictures)} pictures for {len(page.get('media_items', ()))} slots"
                )
            if layout == "module_media" and len(pictures) < len(page.get("modules", ())):
                page_errors.append(
                    f"module_media has {len(pictures)} pictures for {len(page.get('modules', ()))} modules"
                )
            expected_frames = self._expected_editable_frames(page)
            if content_page and len(editable_frames) < expected_frames:
                page_errors.append(
                    f"only {len(editable_frames)} separately editable frame objects; require {expected_frames}"
                )
            if content_page and any(
                (shape.width * shape.height) / (presentation.slide_width * presentation.slide_height) > 0.72
                for shape in pictures
            ):
                page_errors.append("content page contains an unjustified near-full-slide picture")
            errors.extend(f"{page.get('page_id')}: {message}" for message in page_errors)
            slides.append({
                "slide_number": index + 1,
                "page_id": page.get("page_id"),
                "layout": layout,
                "template_reference": reference,
                "density_target": "high" if content_page else "intentional_sparse",
                "density_units": density_units,
                "visible_module_count": len(modules),
                "component_inventory": {
                    "editable_text_objects": len(text_shapes),
                    "native_non_picture_objects": len(native_shapes),
                    "editable_frame_objects": len(editable_frames),
                    "pictures": len(pictures),
                },
                "editable_information_layer_pass": len(text_shapes) >= (5 if content_page else 1),
                "template_identity_bound": bool(reference) if content_page else True,
                "page_errors": page_errors,
            })
        return {
            "schema_version": 1,
            "composition_mode": dynamic_plan.get("composition_mode"),
            "template": template,
            "template_identity": grammar.get("identity", {}),
            "visual_system": {
                "fonts": visual_system.get("fonts", {}),
                "typography": visual_system.get("typography", {}),
                "colors": visual_system.get("colors", {}),
            },
            "qa_expectations": {
                "dual_gate_required": True,
                "template_identity_required": True,
                "all_key_text_editable": True,
                "full_render_required": True,
            },
            "passed": not errors,
            "errors": errors,
            "slides": slides,
        }

    @staticmethod
    def _expected_editable_frames(page: dict) -> int:
        layout = page.get("layout")
        if layout == "text_figure":
            return 3
        if layout == "media_gallery":
            return len(page.get("media_items", ()))
        if layout == "module_media":
            return len(page.get("modules", ())) * 2
        if layout == "points":
            return len(page.get("points", ())) * 2
        if layout == "process":
            return len(page.get("steps", ())) * 2
        if layout == "architecture":
            columns = page.get("architecture", {}).get("columns", ())
            nodes = sum(len(column.get("nodes", ())) for column in columns)
            return len(columns) * 2 + nodes
        return 1


class RenderedObjectBindingManifestBuilder:
    """Bind actual rendered shapes to standard semantic modules and child slots."""

    def build(
        self,
        *,
        pptx_path: Path | str,
        dynamic_plan: dict,
        semantic_specification_path: Path | str,
        scaffold_shape_ids_by_slide: list[list[int]] | None = None,
    ) -> dict:
        presentation = Presentation(pptx_path)
        specification_path = Path(semantic_specification_path).resolve()
        specification = json.loads(specification_path.read_text(encoding="utf-8"))
        semantic_pages = {
            page["page_id"]: page for page in specification.get("pages", ())
        }
        plan_pages = dynamic_plan.get("pages", ())
        scaffold_records = scaffold_shape_ids_by_slide or [[] for _ in plan_pages]
        source_index = self._source_index(dynamic_plan)
        errors = []
        pages = []
        if len(presentation.slides) != len(plan_pages):
            errors.append(
                f"object binding page count mismatch: pptx={len(presentation.slides)} "
                f"plan={len(plan_pages)}"
            )
        for index, page in enumerate(plan_pages):
            if index >= len(presentation.slides):
                break
            slide = presentation.slides[index]
            template_page_id = (page.get("template_reference") or {}).get(
                "semantic_spec_page_id"
            )
            semantic_page = semantic_pages.get(template_page_id)
            if not semantic_page:
                errors.append(
                    f"{page.get('page_id')}: no standard semantic page is bound"
                )
                continue
            scaffold_ids = {
                int(value) for value in (
                    scaffold_records[index] if index < len(scaffold_records) else ()
                )
            }
            assignments, identity_ids = self._assign_shapes(
                slide,
                semantic_page,
                scaffold_ids,
                presentation.slide_width,
                presentation.slide_height,
            )
            media_sources = self._media_sources(page, source_index)
            picture_shapes = sorted(
                (
                    shape for shape in slide.shapes
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                    and int(shape.shape_id) not in scaffold_ids
                ),
                key=lambda shape: (shape.top, shape.left),
            )
            source_by_shape_id = {
                int(shape.shape_id): source
                for shape, source in zip(picture_shapes, media_sources)
            }
            regions = []
            for module in semantic_page.get("semantic_modules", ()):
                module_id = module["module_id"]
                owned_shapes = assignments.get(module_id, [])
                bindings, reflowed = self._slot_bindings(
                    module,
                    owned_shapes,
                    source_by_shape_id,
                )
                regions.append({
                    "module_id": module_id,
                    "render_mode": "full_reconstruction",
                    "owned_shape_ids": sorted(
                        int(shape.shape_id) for shape in owned_shapes
                    ),
                    "removed_native_shape_ids": sorted(
                        int(value) for value in module.get("ownership_group", {}).get(
                            "complete_delete_shape_ids", ()
                        )
                    ),
                    "slot_bindings": bindings,
                    "reflowed_or_removed_slots": reflowed,
                })
            feature_shape_ids = sorted(
                identity_ids
                | {
                    shape_id
                    for region in regions
                    for shape_id in region["owned_shape_ids"]
                }
            )
            identity_features = specification.get("template_identity", {}).get(
                "identity_features", ()
            )
            minimum_features = int(
                specification.get("template_identity", {}).get(
                    "minimum_preserved_features", 4
                )
            )
            identity_feature_bindings = [
                {"feature": feature, "shape_ids": feature_shape_ids}
                for feature in identity_features[:minimum_features]
            ]
            pages.append({
                "slide_number": index + 1,
                "page_id": page.get("page_id"),
                "template_page_id": template_page_id,
                "base_mode": "blank_reconstruction",
                "source_components_present": False,
                "additional_identity_shape_ids": sorted(identity_ids),
                "identity_feature_bindings": identity_feature_bindings,
                "regions": regions,
            })
        return {
            "schema_version": 1,
            "template_id": specification.get("template", {}).get("id"),
            "semantic_specification": str(specification_path),
            "errors": errors,
            "pages": pages,
        }

    @staticmethod
    def _normalized_box(shape, slide_width, slide_height) -> dict:
        return {
            "left": float(shape.left) / float(slide_width),
            "top": float(shape.top) / float(slide_height),
            "width": float(shape.width) / float(slide_width),
            "height": float(shape.height) / float(slide_height),
        }

    @staticmethod
    def _overlap(first: dict, second: dict) -> float:
        left = max(first["left"], second["left"])
        top = max(first["top"], second["top"])
        right = min(first["left"] + first["width"], second["left"] + second["width"])
        bottom = min(first["top"] + first["height"], second["top"] + second["height"])
        if right <= left or bottom <= top:
            return 0.0
        intersection = (right - left) * (bottom - top)
        return intersection / max(first["width"] * first["height"], 1e-9)

    def _assign_shapes(
        self,
        slide,
        semantic_page,
        scaffold_ids,
        slide_width,
        slide_height,
    ):
        modules = semantic_page.get("semantic_modules", ())
        assignments = {module["module_id"]: [] for module in modules}
        identity_ids = set(scaffold_ids)
        for shape in slide.shapes:
            shape_id = int(shape.shape_id)
            if shape_id in scaffold_ids:
                continue
            box = self._normalized_box(shape, slide_width, slide_height)
            scored = [
                (self._overlap(box, module["box"]), module)
                for module in modules
            ]
            score, selected = max(scored, key=lambda item: item[0], default=(0.0, None))
            center_y = box["top"] + box["height"] / 2
            content_left = min(float(module["box"]["left"]) for module in modules)
            content_top = min(float(module["box"]["top"]) for module in modules)
            content_right = max(
                float(module["box"]["left"]) + float(module["box"]["width"])
                for module in modules
            )
            content_bottom = max(
                float(module["box"]["top"]) + float(module["box"]["height"])
                for module in modules
            )
            center_x = box["left"] + box["width"] / 2
            center_in_semantic_canvas = (
                content_left - 0.02 <= center_x <= content_right + 0.02
                and content_top - 0.02 <= center_y <= content_bottom + 0.02
            )
            if selected is not None and (score > 0.02 or center_in_semantic_canvas):
                if score <= 0.02:
                    selected = min(
                        modules,
                        key=lambda module: (
                            center_x - (
                                float(module["box"]["left"])
                                + float(module["box"]["width"]) / 2
                            )
                        ) ** 2
                        + (
                            center_y - (
                                float(module["box"]["top"])
                                + float(module["box"]["height"]) / 2
                            )
                        ) ** 2,
                    )
                assignments[selected["module_id"]].append(shape)
            else:
                identity_ids.add(shape_id)
        return assignments, identity_ids

    @staticmethod
    def _font_size(shape) -> float:
        sizes = []
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.font.size:
                    sizes.append(float(paragraph.font.size.pt))
                sizes.extend(
                    float(run.font.size.pt) for run in paragraph.runs if run.font.size
                )
        return max(sizes, default=0.0)

    def _slot_bindings(self, module, owned_shapes, source_by_shape_id):
        slots = module.get("child_slots", ())
        media_shapes = sorted(
            (
                shape for shape in owned_shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                or getattr(shape, "has_chart", False)
                or getattr(shape, "has_table", False)
            ),
            key=lambda shape: (shape.top, shape.left),
        )
        text_shapes = sorted(
            (
                shape for shape in owned_shapes
                if getattr(shape, "has_text_frame", False) and shape.text.strip()
            ),
            key=lambda shape: (shape.top, shape.left),
        )
        used_text_ids = set()
        bindings = []
        reflowed = []
        media_index = 0
        role_slots: dict[str, list[dict]] = {}
        for slot in slots:
            role_slots.setdefault(slot.get("role", ""), []).append(slot)

        for slot in role_slots.get("image_or_chart", ()):
            if media_index >= len(media_shapes):
                reflowed.append(slot["slot_id"])
                continue
            shape = media_shapes[media_index]
            media_index += 1
            binding = {
                "slot_id": slot["slot_id"],
                "shape_ids": [int(shape.shape_id)],
            }
            source = source_by_shape_id.get(int(shape.shape_id))
            if source:
                binding["source"] = source
            bindings.append(binding)

        def take_text(candidates):
            available = [
                shape for shape in candidates
                if int(shape.shape_id) not in used_text_ids
            ]
            if not available:
                return None
            shape = available[0]
            used_text_ids.add(int(shape.shape_id))
            return shape

        for role in ("page_title", "cover_title", "module_heading"):
            for slot in role_slots.get(role, ()):
                available = sorted(
                    (
                        shape for shape in text_shapes
                        if int(shape.shape_id) not in used_text_ids
                    ),
                    key=lambda shape: (-self._font_size(shape), shape.top, shape.left),
                )
                shape = take_text(available)
                if shape:
                    bindings.append({
                        "slot_id": slot["slot_id"],
                        "shape_ids": [int(shape.shape_id)],
                    })

        caption_slots = role_slots.get("caption", ())
        caption_candidates = sorted(
            (
                shape for shape in text_shapes
                if int(shape.shape_id) not in used_text_ids
            ),
            key=lambda shape: (self._font_size(shape), -shape.top, shape.left),
        )
        for caption_index, slot in enumerate(caption_slots):
            if caption_index >= len(media_shapes) and role_slots.get("image_or_chart"):
                reflowed.append(slot["slot_id"])
                continue
            shape = take_text(caption_candidates)
            if shape:
                bindings.append({
                    "slot_id": slot["slot_id"],
                    "shape_ids": [int(shape.shape_id)],
                })
            elif not slot.get("required"):
                reflowed.append(slot["slot_id"])

        explanation_slots = role_slots.get("explanation", ())
        remaining = [
            shape for shape in text_shapes
            if int(shape.shape_id) not in used_text_ids
        ]
        for slot_index, slot in enumerate(explanation_slots):
            selected = remaining[slot_index::max(1, len(explanation_slots))]
            if selected:
                bindings.append({
                    "slot_id": slot["slot_id"],
                    "shape_ids": [int(shape.shape_id) for shape in selected],
                })
                used_text_ids.update(int(shape.shape_id) for shape in selected)
            elif not slot.get("required"):
                reflowed.append(slot["slot_id"])

        for slot in slots:
            if (
                slot["slot_id"] not in {binding["slot_id"] for binding in bindings}
                and not slot.get("required")
                and slot["slot_id"] not in reflowed
            ):
                reflowed.append(slot["slot_id"])
        return bindings, sorted(reflowed)

    @staticmethod
    def _source_index(dynamic_plan: dict) -> dict[str, dict]:
        base_dir = Path(dynamic_plan.get("asset_base_dir", ".")).resolve()
        figure_root = base_dir / "working" / "figures"
        values = {}
        if not figure_root.is_dir():
            return values
        for manifest_path in figure_root.rglob("manifest.json"):
            try:
                manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
            except (OSError, json.JSONDecodeError):
                continue
            source_path = str(Path(manifest.get("source", manifest_path)).resolve())
            for item in manifest.get("items", ()):
                asset_path = item.get("path")
                if not item.get("accepted") or not asset_path:
                    continue
                resolved = str(Path(asset_path).resolve())
                values[resolved.casefold()] = {
                    "source_id": item.get("asset_id") or item.get("digest") or Path(resolved).stem,
                    "source_path": source_path,
                    "pdf_page": int(item.get("source_page") or item.get("page") or 0),
                    "semantic_use": "source figure evidence",
                    "asset_path": resolved,
                    "caption": str(item.get("caption", "")).strip(),
                }
        return values

    @staticmethod
    def _media_sources(page: dict, source_index: dict[str, dict]) -> list[dict]:
        def resolved_source(asset_path, explicit):
            resolved_asset = str(Path(str(asset_path)).resolve())
            recovered = dict(source_index.get(resolved_asset.casefold(), {}))
            recovered.update({
                key: value for key, value in dict(explicit or {}).items()
                if value not in {None, ""}
            })
            recovered.setdefault("asset_path", resolved_asset)
            recovered.setdefault("semantic_use", f"支持页面结论：{page.get('title', '')}")
            return recovered

        layout = page.get("layout")
        if layout in {"text_figure", "full_figure"} and page.get("image"):
            return [resolved_source(page["image"], page.get("media_source", {}))]
        if layout == "media_gallery":
            values = []
            for item in page.get("media_items", ()):
                values.append(resolved_source(item.get("path"), item.get("source", {})))
            return values
        if layout == "module_media":
            values = []
            for module in page.get("modules", ()):
                values.append(
                    resolved_source(module.get("image"), module.get("media_source", {}))
                )
            return values
        return []
