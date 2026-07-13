"""Protect user-edited PPTX pages during explicitly scoped revisions."""

from __future__ import annotations

import hashlib
import json
from dataclasses import asdict, dataclass
from pathlib import Path

from pptx import Presentation


@dataclass(frozen=True)
class SlideSnapshot:
    page_id: str
    shape_hashes: dict[int, str]


@dataclass(frozen=True)
class RevisionPlan:
    target_page_ids: tuple[str, ...]
    protected_page_ids: tuple[str, ...]
    changed_page_ids: tuple[str, ...]
    manual_shape_ids: dict[str, tuple[int, ...]]
    replacement_disclosures: tuple[str, ...]

    def to_dict(self) -> dict:
        return {
            "target_page_ids": list(self.target_page_ids),
            "protected_page_ids": list(self.protected_page_ids),
            "changed_page_ids": list(self.changed_page_ids),
            "manual_shape_ids": {page_id: list(shape_ids) for page_id, shape_ids in self.manual_shape_ids.items()},
            "replacement_disclosures": list(self.replacement_disclosures),
        }


class AuthoritativeEditBaseline:
    """Snapshot a generated deck before targeted user-driven revision work."""

    def __init__(self, source_path: Path, page_ids: tuple[str, ...], snapshots: tuple[SlideSnapshot, ...]):
        self.source_path = source_path
        self.page_ids = page_ids
        self._snapshots = snapshots

    @classmethod
    def adopt(cls, pptx_path: Path | str, page_ids: tuple[str, ...] | list[str]) -> "AuthoritativeEditBaseline":
        path = Path(pptx_path).resolve()
        if not path.is_file():
            raise FileNotFoundError(path)
        resolved_page_ids = tuple(page_ids)
        if not resolved_page_ids or any(not page_id.strip() for page_id in resolved_page_ids):
            raise ValueError("baseline requires non-empty page IDs")
        if len(set(resolved_page_ids)) != len(resolved_page_ids):
            raise ValueError("baseline page IDs must be unique")
        presentation = Presentation(path)
        if len(presentation.slides) != len(resolved_page_ids):
            raise ValueError("baseline page IDs must match PPTX slide count")
        snapshots = tuple(
            SlideSnapshot(page_id, cls._shape_hashes(slide))
            for page_id, slide in zip(resolved_page_ids, presentation.slides)
        )
        return cls(path, resolved_page_ids, snapshots)

    def plan_revision(self, edited_pptx_path: Path | str, *, target_page_ids: tuple[str, ...] | list[str]) -> RevisionPlan:
        path = Path(edited_pptx_path).resolve()
        if not path.is_file():
            raise FileNotFoundError(path)
        targets = tuple(target_page_ids)
        unknown = [page_id for page_id in targets if page_id not in self.page_ids]
        if unknown:
            raise ValueError(f"unknown revision target: {', '.join(unknown)}")
        if not targets:
            raise ValueError("revision requires at least one target page")
        presentation = Presentation(path)
        if len(presentation.slides) != len(self.page_ids):
            raise ValueError("edited PPTX slide count differs from the authoritative baseline")

        changed = []
        manual_shape_ids: dict[str, tuple[int, ...]] = {}
        for snapshot, slide in zip(self._snapshots, presentation.slides):
            current = self._shape_hashes(slide)
            if current != snapshot.shape_hashes:
                changed.append(snapshot.page_id)
            added = tuple(sorted(set(current) - set(snapshot.shape_hashes)))
            if added:
                manual_shape_ids[snapshot.page_id] = added
        changed_set = set(changed)
        disclosures = tuple(
            f"{page_id}: user edits may be replaced by the requested rebuild"
            for page_id in targets
            if page_id in changed_set
        )
        return RevisionPlan(
            target_page_ids=targets,
            protected_page_ids=tuple(page_id for page_id in self.page_ids if page_id not in targets),
            changed_page_ids=tuple(changed),
            manual_shape_ids=manual_shape_ids,
            replacement_disclosures=disclosures,
        )

    def to_dict(self) -> dict:
        return {
            "source_path": str(self.source_path),
            "page_ids": list(self.page_ids),
            "snapshots": [asdict(snapshot) for snapshot in self._snapshots],
        }

    def write(self, audit_dir: Path | str) -> Path:
        destination = Path(audit_dir)
        destination.mkdir(parents=True, exist_ok=True)
        output = destination / "authoritative_edit_baseline.json"
        output.write_text(json.dumps(self.to_dict(), ensure_ascii=False, indent=2), encoding="utf-8")
        return output.resolve()

    @classmethod
    def load(cls, baseline_path: Path | str) -> "AuthoritativeEditBaseline":
        path = Path(baseline_path).resolve()
        if not path.is_file():
            raise FileNotFoundError(path)
        payload = json.loads(path.read_text(encoding="utf-8"))
        page_ids = tuple(payload.get("page_ids", ()))
        snapshots = tuple(
            SlideSnapshot(
                page_id=item["page_id"],
                shape_hashes={int(shape_id): digest for shape_id, digest in item["shape_hashes"].items()},
            )
            for item in payload.get("snapshots", ())
        )
        if tuple(snapshot.page_id for snapshot in snapshots) != page_ids:
            raise ValueError("baseline snapshots do not match page IDs")
        return cls(Path(payload["source_path"]).resolve(), page_ids, snapshots)

    @staticmethod
    def _shape_hashes(slide) -> dict[int, str]:
        return {
            shape.shape_id: hashlib.sha256(shape._element.xml.encode("utf-8")).hexdigest()
            for shape in slide.shapes
        }
