"""Inspect editable PPTX components for native template reuse."""

from __future__ import annotations

import json
import subprocess
import sys
import tempfile
from collections import Counter
from dataclasses import dataclass
from pathlib import Path


def _normalize_selection(value: str) -> str:
    return "".join(value.strip().casefold().split()).replace("模板", "模版")


@dataclass(frozen=True)
class TemplateSelection:
    template_id: str | None
    short_name: str
    path: str
    support_level: str
    requested: str | None = None
    source_path: str | None = None
    substitution_reason: str | None = None
    source_fidelity: str | None = None
    source_limitations: str | None = None
    semantic_spec_path: str | None = None
    standardization_status: str | None = None

    def to_dict(self) -> dict:
        return {
            "id": self.template_id,
            "short_name": self.short_name,
            "path": self.path,
            "support_level": self.support_level,
            "requested": self.requested,
            "source_path": self.source_path,
            "substitution_reason": self.substitution_reason,
            "source_fidelity": self.source_fidelity,
            "source_limitations": self.source_limitations,
            "semantic_spec_path": self.semantic_spec_path,
            "standardization_status": self.standardization_status,
        }


class TemplateCatalog:
    """Resolve bundled identities and explicitly classify arbitrary user templates."""

    def __init__(self, root: Path, templates: list[dict]):
        self.root = root.resolve()
        self.templates = tuple(templates)

    @classmethod
    def load(cls, path: Path | str | None = None) -> "TemplateCatalog":
        catalog_path = Path(path) if path else Path(__file__).resolve().parents[1] / "references" / "template-catalog.json"
        payload = json.loads(catalog_path.read_text(encoding="utf-8"))
        return cls(catalog_path.resolve().parents[1], payload["templates"])

    def select(self, scene: str, selection: Path | str | None = None) -> TemplateSelection:
        if selection is None:
            item = next(
                (value for value in self.templates if scene in value.get("recommended_scenes", ())),
                None,
            )
            if item is None:
                raise ValueError(f"no bundled template recommendation for scene: {scene}")
            return self._bundled(item, requested=None, source_requested=False)

        raw = str(selection)
        wanted = _normalize_selection(raw)
        input_path = Path(raw)
        resolved_input = input_path.resolve() if input_path.is_file() else None
        matches: list[tuple[dict, bool]] = []
        for item in self.templates:
            compiled_path = (self.root / item["path"]).resolve()
            source_path = (self.root / item["source_path"]).resolve() if item.get("source_path") else None
            compiled_names = {
                _normalize_selection(item["id"]),
                _normalize_selection(item["short_name"]),
                _normalize_selection(compiled_path.name),
                _normalize_selection(compiled_path.stem),
                *(_normalize_selection(alias) for alias in item.get("aliases", ())),
            }
            source_names = set()
            if source_path is not None:
                source_names = {
                    _normalize_selection(source_path.name),
                    _normalize_selection(source_path.stem),
                }
            source_requested = wanted in source_names or (resolved_input is not None and resolved_input == source_path)
            compiled_requested = wanted in compiled_names or (resolved_input is not None and resolved_input == compiled_path)
            if source_requested or compiled_requested:
                matches.append((item, source_requested))
        if len(matches) == 1:
            item, source_requested = matches[0]
            return self._bundled(item, requested=raw, source_requested=source_requested)
        if len(matches) > 1:
            raise ValueError(f"template selection is ambiguous: {selection}")
        if resolved_input is not None:
            return TemplateSelection(
                template_id=None,
                short_name=resolved_input.stem,
                path=str(resolved_input),
                support_level="conditional_user",
                requested=raw,
                source_path=str(resolved_input),
            )
        raise ValueError(f"template selection does not match a bundled template or existing PPTX: {selection}")

    def _bundled(self, item: dict, *, requested: str | None, source_requested: bool) -> TemplateSelection:
        compiled_path = (self.root / item["path"]).resolve()
        if not compiled_path.is_file():
            raise FileNotFoundError(compiled_path)
        source_path = (self.root / item["source_path"]).resolve() if item.get("source_path") else None
        return TemplateSelection(
            template_id=item["id"],
            short_name=item["short_name"],
            path=str(compiled_path),
            support_level="bundled_recompiled_source" if source_requested else "bundled_formal",
            requested=requested,
            source_path=str(source_path) if source_path else None,
            substitution_reason=(
                "The requested bundled source package is repaired into the catalog path while preserving "
                "its complete slide structure; this is a package recompile, not a style-family substitute."
                if source_requested else None
            ),
            source_fidelity=item.get("source_fidelity") if source_requested else None,
            source_limitations=item.get("source_limitations") if source_requested else None,
            semantic_spec_path=(
                str((self.root / item["semantic_spec_path"]).resolve())
                if item.get("semantic_spec_path") else None
            ),
            standardization_status=item.get("standardization_status"),
        )


@dataclass(frozen=True)
class TemplateAdmissionResult:
    passed: bool
    slide_count: int
    editable_component_count: int
    grammar_extracted: bool
    runtime_status: str
    errors: tuple[str, ...]

    def to_dict(self) -> dict:
        return {
            "passed": self.passed,
            "slide_count": self.slide_count,
            "editable_component_count": self.editable_component_count,
            "grammar_extracted": self.grammar_extracted,
            "runtime_status": self.runtime_status,
            "errors": list(self.errors),
        }


class TemplateAdmissionGate:
    """Admit an arbitrary PPTX only after editable, grammar, and runtime checks."""

    def inspect(
        self,
        path: Path | str,
        *,
        require_runtime: bool,
        runtime: str = "powerpoint",
    ) -> TemplateAdmissionResult:
        template_path = Path(path).resolve()
        errors: list[str] = []
        slide_count = 0
        editable_count = 0
        grammar_extracted = False
        runtime_status = "not_requested" if not require_runtime else "failed"
        if template_path.suffix.lower() != ".pptx":
            errors.append(f"user template must be a PPTX: {template_path}")
        if not template_path.is_file():
            errors.append(f"user template not found: {template_path}")
        if errors:
            return TemplateAdmissionResult(False, 0, 0, False, runtime_status, tuple(errors))

        try:
            graph = TemplateCapabilityGraph.from_presentation(template_path)
            slide_count = len(graph.slides)
            editable_count = sum(
                1
                for slide in graph.slides
                for component in slide.components
                if component.kind in {"text", "picture", "table", "chart"}
            )
            if slide_count == 0:
                errors.append("user template has no slides")
            if editable_count == 0:
                errors.append("user template has no editable text, picture, table, or chart components")
        except Exception as exc:
            errors.append(f"user template cannot be parsed: {exc}")

        skill_root = Path(__file__).resolve().parents[1]
        with tempfile.TemporaryDirectory(prefix="academic-ppt-user-template-") as temp_dir:
            temp = Path(temp_dir)
            grammar_path = temp / "template_grammar.json"
            grammar_command = [
                sys.executable,
                str(skill_root / "scripts" / "extract_template_grammar.py"),
                str(template_path),
                "--output",
                str(grammar_path),
                "--asset-dir",
                str(temp / "assets"),
            ]
            grammar_run = subprocess.run(
                grammar_command,
                cwd=skill_root,
                check=False,
                capture_output=True,
                text=True,
                encoding="utf-8",
            )
            if grammar_run.returncode == 0 and grammar_path.is_file():
                try:
                    json.loads(grammar_path.read_text(encoding="utf-8"))
                    grammar_extracted = True
                except (OSError, json.JSONDecodeError) as exc:
                    errors.append(f"user template grammar is invalid: {exc}")
            else:
                detail = (grammar_run.stderr or grammar_run.stdout).strip()
                errors.append(f"user template grammar extraction failed: {detail}")

            compatibility_path = temp / "template_compatibility.pptx"
            try:
                from pptx import Presentation
                from scripts.pptx_utils import clone_slide, remove_slide

                presentation = Presentation(template_path)
                original_count = len(presentation.slides)
                indices = sorted({0, original_count // 2, original_count - 1}) if original_count else []
                for index in indices:
                    clone_slide(presentation, presentation.slides[index])
                for _ in range(original_count):
                    remove_slide(presentation, 0)
                presentation.save(compatibility_path)
            except Exception as exc:
                errors.append(f"user template clone/save compatibility failed: {exc}")

            report_path = temp / "render_report.json"
            render_command = [
                sys.executable,
                str(skill_root / "scripts" / "validate_pptx.py"),
                str(compatibility_path),
                "--output",
                str(report_path),
                "--render-check",
                "required" if require_runtime else "off",
                "--render-engine",
                runtime if runtime != "portable" else "auto",
            ]
            render_run = subprocess.run(
                render_command,
                cwd=skill_root,
                check=False,
                capture_output=True,
                text=True,
                encoding="utf-8",
            ) if compatibility_path.is_file() else None
            if report_path.is_file():
                report = json.loads(report_path.read_text(encoding="utf-8"))
                if report.get("failed_error", 1):
                    errors.append("user template failed structural or runtime validation")
                real_check = next(
                    (check for check in report.get("checks", ()) if check.get("name") == "Real application render succeeds"),
                    None,
                )
                if not require_runtime:
                    runtime_status = "not_requested"
                elif real_check and real_check.get("passed"):
                    runtime_status = "passed"
                else:
                    detail = str(real_check.get("detail", "")) if real_check else ""
                    unavailable = any(
                        marker in detail.casefold()
                        for marker in ("not installed", "not registered", "not found", "no preview engine")
                    )
                    runtime_status = "unavailable" if unavailable else "failed"
            else:
                detail = ((render_run.stderr or render_run.stdout).strip() if render_run else "compatibility sample missing")
                errors.append(f"user template validation did not produce a report: {detail}")

        return TemplateAdmissionResult(
            passed=not errors,
            slide_count=slide_count,
            editable_component_count=editable_count,
            grammar_extracted=grammar_extracted,
            runtime_status=runtime_status,
            errors=tuple(errors),
        )


@dataclass(frozen=True)
class TemplateComponent:
    kind: str
    shape_id: int
    geometry: dict[str, int]
    parent_group_shape_id: int | None = None


@dataclass(frozen=True)
class TemplateSlideCapability:
    slide_index: int
    components: tuple[TemplateComponent, ...]

    @property
    def component_counts(self) -> dict[str, int]:
        return dict(sorted(Counter(component.kind for component in self.components).items()))


@dataclass(frozen=True)
class TemplateCapabilityGraph:
    template_path: str
    slides: tuple[TemplateSlideCapability, ...]

    @classmethod
    def from_presentation(cls, path: Path | str) -> "TemplateCapabilityGraph":
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError as exc:
            raise RuntimeError("python-pptx is required to inspect templates") from exc

        template_path = Path(path).resolve()
        if not template_path.is_file():
            raise FileNotFoundError(template_path)
        presentation = Presentation(template_path)
        slides = []
        for slide_index, slide in enumerate(presentation.slides, 1):
            components = []
            for shape in slide.shapes:
                geometry = {
                    "left": int(shape.left),
                    "top": int(shape.top),
                    "width": int(shape.width),
                    "height": int(shape.height),
                }
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    components.append(TemplateComponent(
                        kind="group",
                        shape_id=shape.shape_id,
                        geometry=geometry,
                    ))
                elif getattr(shape, "has_table", False):
                    components.append(TemplateComponent("table", shape.shape_id, geometry))
                elif getattr(shape, "has_chart", False):
                    components.append(TemplateComponent("chart", shape.shape_id, geometry))
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    components.append(TemplateComponent("picture", shape.shape_id, geometry))
                elif getattr(shape, "has_text_frame", False):
                    components.append(TemplateComponent("text", shape.shape_id, geometry))
            slides.append(TemplateSlideCapability(
                slide_index=slide_index,
                components=tuple(components),
            ))
        return cls(template_path=str(template_path), slides=tuple(slides))

    def find_compatible_slides(
        self,
        required_components: dict[str, int],
    ) -> tuple[TemplateSlideCapability, ...]:
        invalid = {kind: count for kind, count in required_components.items() if count < 1}
        if invalid:
            raise ValueError(f"component requirements must be positive: {invalid}")
        return tuple(
            slide for slide in self.slides
            if all(slide.component_counts.get(kind, 0) >= count for kind, count in required_components.items())
        )
