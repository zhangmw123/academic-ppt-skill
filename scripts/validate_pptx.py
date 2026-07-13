"""Validate PPTX package integrity, semantic bindings, residue, overflow, and real rendering."""

from __future__ import annotations

import argparse
import json
import math
import posixpath
import re
import statistics
import sys
import zipfile
from datetime import datetime, timezone
from pathlib import Path, PurePosixPath
from xml.etree import ElementTree as ET

from pptx import Presentation

from export_preview import export_preview, find_powerpoint
from pptx_utils import iter_shapes


PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
DOC_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
SAMPLE_RE = re.compile(
    r"这里(填写|输入)|输入(相关|正文|标题|意义|应用)|小云集市|大云集市|"
    r"xxxx|20xx|click to|insert (text|title)|lorem ipsum",
    re.I,
)


def add_check(report, name: str, level: str, passed: bool, detail: str):
    report["checks"].append({"name": name, "level": level, "passed": passed, "detail": detail})
    if passed:
        report["passed"] += 1
    elif level == "ERROR":
        report["failed_error"] += 1
    else:
        report["failed_warning"] += 1


def rels_path_for(part_name: str) -> str:
    path = PurePosixPath(part_name)
    return str(path.parent / "_rels" / f"{path.name}.rels")


def source_path_for_rels(rels_name: str) -> str | None:
    if rels_name == "_rels/.rels":
        return None
    path = PurePosixPath(rels_name)
    if path.parent.name != "_rels":
        return None
    return str(path.parent.parent / path.name.removesuffix(".rels"))


def package_integrity(path: Path) -> list[str]:
    errors = []
    try:
        with zipfile.ZipFile(path) as archive:
            bad_file = archive.testzip()
            if bad_file:
                errors.append(f"CRC failure: {bad_file}")
            names = set(archive.namelist())
            relation_ids: dict[str, set[str]] = {}
            for name in names:
                if not name.endswith(".rels"):
                    continue
                try:
                    root = ET.fromstring(archive.read(name))
                except ET.ParseError as exc:
                    errors.append(f"Invalid relationships XML {name}: {exc}")
                    continue
                source = source_path_for_rels(name)
                base = posixpath.dirname(source) if source else ""
                ids = set()
                for rel in root.findall(f"{{{PKG_REL_NS}}}Relationship"):
                    rid = rel.get("Id")
                    if rid:
                        ids.add(rid)
                    if rel.get("TargetMode") == "External":
                        continue
                    target = posixpath.normpath(posixpath.join(base, rel.get("Target", "")))
                    if target not in names:
                        errors.append(f"Missing relationship target: {name} {rid} -> {target}")
                if source:
                    relation_ids[source] = ids

            relationship_attrs = {f"{{{DOC_REL_NS}}}id", f"{{{DOC_REL_NS}}}embed", f"{{{DOC_REL_NS}}}link"}
            for name in names:
                if not name.endswith(".xml") or name.endswith(".rels"):
                    continue
                try:
                    root = ET.fromstring(archive.read(name))
                except ET.ParseError as exc:
                    errors.append(f"Invalid XML {name}: {exc}")
                    continue
                referenced = {
                    value for element in root.iter() for attr, value in element.attrib.items()
                    if attr in relationship_attrs and value.startswith("rId")
                }
                missing = referenced - relation_ids.get(name, set())
                if missing:
                    errors.append(f"Unresolved relationship IDs in {name}: {sorted(missing)}")
    except (zipfile.BadZipFile, OSError) as exc:
        errors.append(f"Invalid PPTX package: {exc}")
    return errors


def estimate_overflow(shape) -> float:
    if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
        return 0.0
    width_pt = shape.width / 12700
    height_pt = shape.height / 12700
    required_height = 0.0
    for paragraph in shape.text_frame.paragraphs:
        sizes = []
        if paragraph.font.size:
            sizes.append(paragraph.font.size.pt)
        sizes.extend(run.font.size.pt for run in paragraph.runs if run.font.size)
        font_size = statistics.median(sizes) if sizes else 16.0
        chars_per_line = max(1, width_pt / (font_size * 0.8))
        required_lines = max(1, math.ceil(len(paragraph.text) / chars_per_line))
        required_height += required_lines * font_size * 1.2
    return required_height / max(height_pt, 1)


def validate(pptx_path: Path, plan_path: Path | None, render_mode: str, preview_dir: Path,
             selected_page_ids: set[str] | None = None, render_engine: str = "auto"):
    report = {
        "file": str(pptx_path.resolve()),
        "validated_at": datetime.now(timezone.utc).isoformat(),
        "checks": [], "passed": 0, "failed_error": 0, "failed_warning": 0,
    }

    package_errors = package_integrity(pptx_path)
    add_check(report, "OOXML package and relationships are valid", "ERROR", not package_errors,
              "OK" if not package_errors else "; ".join(package_errors[:12]))
    if package_errors:
        return report

    try:
        prs = Presentation(pptx_path)
    except Exception as exc:
        add_check(report, "python-pptx can parse the file", "ERROR", False, str(exc))
        return report
    add_check(report, "python-pptx can parse the file", "ERROR", True, f"{len(prs.slides)} slides")

    plan = None
    if plan_path:
        plan = json.loads(plan_path.read_text(encoding="utf-8"))
        expected_pages = plan.get("pages", [])
        if selected_page_ids:
            expected_pages = [page for page in expected_pages if page.get("page_id") in selected_page_ids]
        expected = len(expected_pages)
        if expected != len(prs.slides) and expected:
            detail = f"plan has {expected} pages; PPTX has {len(prs.slides)}"
            add_check(report, "Slide count matches layout plan", "ERROR", False, detail)
        else:
            add_check(report, "Slide count matches layout plan", "ERROR", True, f"{len(prs.slides)} slides")

    residue = []
    overflow_errors = []
    overflow_warnings = []
    overlap_errors = []
    out_of_bounds = []
    for slide_index, slide in enumerate(prs.slides, 1):
        text_shapes = []
        for shape in iter_shapes(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if text:
                text_shapes.append(shape)
                if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
                    out_of_bounds.append(f"slide {slide_index}, shape {shape.shape_id}")
            if text and SAMPLE_RE.search(text):
                residue.append(f"slide {slide_index}, shape {shape.shape_id}: {text[:60]}")
            ratio = estimate_overflow(shape)
            if ratio > 2.20:
                overflow_errors.append(f"slide {slide_index}, shape {shape.shape_id}: {ratio:.0%}")
            elif ratio > 1.60:
                overflow_warnings.append(f"slide {slide_index}, shape {shape.shape_id}: {ratio:.0%}")
        for left_index, first in enumerate(text_shapes):
            for second in text_shapes[left_index + 1:]:
                overlap_w = min(first.left + first.width, second.left + second.width) - max(first.left, second.left)
                overlap_h = min(first.top + first.height, second.top + second.height) - max(first.top, second.top)
                if overlap_w <= 0 or overlap_h <= 0:
                    continue
                intersection = overlap_w * overlap_h
                smaller = min(first.width * first.height, second.width * second.height)
                if smaller and intersection / smaller > 0.15:
                    overlap_errors.append(
                        f"slide {slide_index}: text shapes {first.shape_id} and {second.shape_id} overlap {intersection / smaller:.0%}"
                    )
    add_check(report, "No template sample text remains", "ERROR", not residue,
              "OK" if not residue else "; ".join(residue[:10]))
    add_check(report, "No severe estimated text overflow", "ERROR", not overflow_errors,
              "OK" if not overflow_errors else "; ".join(overflow_errors[:10]))
    add_check(report, "Text boxes do not materially overlap", "ERROR", not overlap_errors,
              "OK" if not overlap_errors else "; ".join(overlap_errors[:10]))
    add_check(report, "Text boxes stay inside the slide canvas", "ERROR", not out_of_bounds,
              "OK" if not out_of_bounds else "; ".join(out_of_bounds[:10]))
    if overflow_warnings:
        add_check(report, "No possible text overflow", "WARNING", False, "; ".join(overflow_warnings[:10]))

    if render_mode != "off":
        try:
            export_preview(pptx_path, preview_dir, render_engine, 1600, 900)
        except Exception as exc:
            level = "ERROR" if render_mode == "required" or find_powerpoint() else "WARNING"
            add_check(report, "Real application render succeeds", level, False, str(exc))
        else:
            add_check(report, "Real application render succeeds", "ERROR", True,
                      str((preview_dir / "contact-sheet.jpg").resolve()))
    return report


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx")
    parser.add_argument("--layout-plan")
    parser.add_argument("--output", required=True)
    parser.add_argument("--render-check", choices=["off", "auto", "required"], default="auto")
    parser.add_argument("--render-engine", choices=["auto", "powerpoint", "wps", "libreoffice"], default="auto")
    parser.add_argument("--preview-dir")
    parser.add_argument("--pages", help="Comma-separated page IDs when validating a sample render")
    args = parser.parse_args()
    output = Path(args.output)
    preview_dir = Path(args.preview_dir) if args.preview_dir else output.parent / "rendered-preview"
    selected = {value.strip() for value in args.pages.split(",")} if args.pages else None
    report = validate(Path(args.pptx), Path(args.layout_plan) if args.layout_plan else None,
                      args.render_check, preview_dir, selected, args.render_engine)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps({key: report[key] for key in ["passed", "failed_error", "failed_warning"]}, ensure_ascii=False))
    print(f"Report: {output}")
    raise SystemExit(1 if report["failed_error"] else 0)


if __name__ == "__main__":
    main()
