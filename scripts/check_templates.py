"""Check that every bundled template can be cloned, saved, and optionally rendered."""

from __future__ import annotations

import argparse
import json
import tempfile
from pathlib import Path

from pptx import Presentation

from export_preview import export_preview
from extract_template_grammar import extract as extract_template_grammar
from extract_visual_system import extract as extract_visual_system
from pptx_utils import clone_slide, remove_slide
from render_dynamic import render as render_dynamic
from validate_pptx import package_integrity


def check_template(path: Path, work_dir: Path, render_check: bool, style_learning_check: bool,
                   grammar_check: bool) -> dict:
    result = {
        "template": path.name,
        "path": str(path.resolve()),
        "passed": False,
        "slide_count": 0,
        "sampled_source_slides": [],
        "errors": [],
        "style_learning_passed": None,
        "grammar_passed": None,
        "identity": None,
    }
    try:
        prs = Presentation(path)
        original_count = len(prs.slides)
        result["slide_count"] = original_count
        if original_count == 0:
            raise ValueError("template has no slides")
        indices = sorted({0, original_count // 2, original_count - 1})
        result["sampled_source_slides"] = [index + 1 for index in indices]
        for index in indices:
            clone_slide(prs, prs.slides[index])
        for _ in range(original_count):
            remove_slide(prs, 0)
        output = work_dir / f"{path.stem}_compatibility.pptx"
        prs.save(output)
        result["errors"].extend(package_integrity(output))
        preview_dir = work_dir / f"{path.stem}_preview"
        if not result["errors"] and render_check:
            export_preview(output, preview_dir, "auto", 1280, 720)
        if not result["errors"] and style_learning_check:
            if not render_check:
                raise ValueError("--style-learning-check requires --render-check")
            profile = extract_visual_system(path, preview_dir)
            profile_path = work_dir / f"{path.stem}_visual_system.json"
            profile_path.write_text(json.dumps(profile, ensure_ascii=False, indent=2), encoding="utf-8")
            dynamic_plan = {
                "confirmed": True,
                "sections": ["研究问题", "方法设计", "实验验证"],
                "footer": "Template portability check",
                "pages": [
                    {"page_id": "P01", "layout": "cover", "section": None,
                     "title": "Academic Presentation", "subtitle": "Template style-learning check",
                     "metadata": [path.name], "visual_strategy": "native_cover"},
                    {"page_id": "P02", "layout": "comparison", "section": "方法设计",
                     "title": "Navigation and layout follow content structure",
                     "columns": [
                         {"title": "Evidence", "lead": "Source-grounded", "bullets": ["Traceable claims", "Explicit caveats"]},
                         {"title": "Design", "lead": "Template-informed", "bullets": ["Dynamic sections", "Readable hierarchy"]}
                     ], "visual_strategy": "native_comparison"},
                    {"page_id": "P03", "layout": "ending", "section": "实验验证",
                     "title": "Style learned successfully", "subtitle": "The sample page structure was not copied.",
                     "visual_strategy": "native_conclusion"}
                ]
            }
            if grammar_check:
                grammar_path = work_dir / f"{path.stem}_grammar.json"
                if not grammar_path.exists():
                    extract_template_grammar(path, grammar_path, work_dir / f"{path.stem}_grammar_assets")
                dynamic_plan["template_grammar"] = str(grammar_path)
            plan_path = work_dir / f"{path.stem}_dynamic_plan.json"
            plan_path.write_text(json.dumps(dynamic_plan, ensure_ascii=False, indent=2), encoding="utf-8")
            dynamic_output = work_dir / f"{path.stem}_dynamic.pptx"
            render_dynamic(plan_path, profile_path, dynamic_output)
            dynamic_errors = package_integrity(dynamic_output)
            if dynamic_errors:
                result["errors"].extend(dynamic_errors)
            else:
                export_preview(dynamic_output, work_dir / f"{path.stem}_dynamic_preview", "auto", 1280, 720)
                result["style_learning_passed"] = True
        if not result["errors"] and grammar_check:
            grammar_path = work_dir / f"{path.stem}_grammar.json"
            if not grammar_path.exists():
                extract_template_grammar(path, grammar_path, work_dir / f"{path.stem}_grammar_assets")
            grammar = json.loads(grammar_path.read_text(encoding="utf-8"))
            result["grammar_passed"] = True
            result["identity"] = grammar["identity"]
    except Exception as exc:
        result["errors"].append(str(exc))
    result["passed"] = not result["errors"]
    return result


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--template-dir", default=str(Path(__file__).resolve().parents[1] / "assets" / "templates"))
    parser.add_argument("--render-check", action="store_true", help="Open every sample in PowerPoint or LibreOffice")
    parser.add_argument("--style-learning-check", action="store_true", help="Also learn style and render a dynamic three-page deck")
    parser.add_argument("--grammar-check", action="store_true", help="Extract template grammar and identity tokens")
    parser.add_argument("--output", help="Optional JSON report path")
    args = parser.parse_args()

    templates = sorted(Path(args.template_dir).glob("*.pptx"))
    if not templates:
        raise SystemExit("No PPTX templates found")
    with tempfile.TemporaryDirectory(prefix="academic-ppt-template-check-") as temp:
        results = [check_template(path, Path(temp), args.render_check, args.style_learning_check, args.grammar_check) for path in templates]
    report = {
        "template_count": len(results),
        "passed_count": sum(1 for result in results if result["passed"]),
        "render_check": args.render_check,
        "results": results,
    }
    if args.output:
        output = Path(args.output)
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    for result in results:
        status = "PASS" if result["passed"] else "FAIL"
        print(f"[{status}] {result['template']} slides={result['slide_count']} sampled={result['sampled_source_slides']}")
        for error in result["errors"]:
            print(f"       {error}")
    raise SystemExit(0 if report["passed_count"] == report["template_count"] else 1)


if __name__ == "__main__":
    main()
