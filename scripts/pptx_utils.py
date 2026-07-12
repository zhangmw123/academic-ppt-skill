from __future__ import annotations

import copy
from pathlib import Path
from typing import Iterable

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
SLIDE_LAYOUT_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
NOTES_SLIDE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"


def iter_shapes(shapes) -> Iterable:
    """Yield top-level and grouped shapes in stable document order."""
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def find_shape(slide, shape_id: int):
    for shape in iter_shapes(slide.shapes):
        if shape.shape_id == shape_id:
            return shape
    return None


def clone_slide(prs, source_slide):
    """Clone a slide and remap every relationship referenced by its XML."""
    new_slide = prs.slides.add_slide(source_slide.slide_layout)
    rel_map: dict[str, str] = {}

    for old_rid, rel in source_slide.part.rels.items():
        if rel.reltype == NOTES_SLIDE_REL:
            continue
        if rel.is_external:
            new_rid = new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rid = new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)
        rel_map[old_rid] = new_rid

    target = new_slide._element
    source = source_slide._element
    for child in list(target):
        target.remove(child)
    for child in source:
        target.append(copy.deepcopy(child))
    target.attrib.clear()
    target.attrib.update(source.attrib)

    for element in target.iter():
        for attr_name, attr_value in list(element.attrib.items()):
            if attr_name.startswith(f"{{{REL_NS}}}") and attr_value in rel_map:
                element.set(attr_name, rel_map[attr_value])

    # add_slide() materializes a SlideShapes proxy before the source XML is
    # copied. Drop that stale proxy so the next access reflects the cloned tree.
    new_slide.__dict__.pop("shapes", None)

    return new_slide


def remove_slide(prs, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    rid = slide_id.get(qn("r:id"))
    prs.part.drop_rel(rid)
    del prs.slides._sldIdLst[index]


def emu_to_inches(value: int) -> float:
    return round(value / 914400, 3)


def resolve_path(path_value: str, base_dir: Path) -> Path:
    path = Path(path_value)
    if not path.is_absolute():
        path = base_dir / path
    return path.resolve()
