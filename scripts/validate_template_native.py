"""Enforce native template reuse and explicit fallback disclosure."""

from __future__ import annotations

import argparse
import json
from collections import Counter
from pathlib import Path

from pptx import Presentation

from pptx_utils import iter_shapes


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("plan")
    parser.add_argument("template")
    parser.add_argument("--output")
    parser.add_argument("--max-freeform-ratio", type=float, default=0.2)
    args = parser.parse_args()
    plan_path = Path(args.plan)
    plan = json.loads(plan_path.read_text(encoding="utf-8"))
    prs = Presentation(args.template)
    source_ids = {
        index: {shape.shape_id for shape in iter_shapes(slide.shapes)}
        for index, slide in enumerate(prs.slides, 1)
    }
    errors, warnings, pages = [], [], []
    component_ids = {}
    grammar_value = plan.get("template_grammar")
    if grammar_value:
        grammar_path = Path(grammar_value)
        if not grammar_path.is_absolute() and not grammar_path.exists():
            candidates = [
                plan_path.resolve().parent / grammar_path,
                plan_path.resolve().parent.parent / grammar_path,
            ]
            grammar_path = next((path for path in candidates if path.exists()), candidates[0])
        if grammar_path.exists():
            grammar = json.loads(grammar_path.read_text(encoding="utf-8"))
            for archetype in grammar.get("archetypes", []):
                component_ids[int(archetype.get("slide_index", 0))] = {
                    item.get("component_id") for item in archetype.get("components", [])
                }
    mode_counts = Counter()
    used_source_slides = set()
    for page in plan.get("pages", []):
        page_id = page.get("page_id", "UNKNOWN")
        mode = page.get("render_mode", "template_native")
        mode_counts[mode] += 1
        source_slide = page.get("source_slide_index")
        bindings = page.get("text_bindings", [])
        if mode in {"template_native", "template_adaptive"}:
            if source_slide not in source_ids:
                errors.append(f"{page_id}: native page requires a valid source_slide_index")
                available = set()
            else:
                used_source_slides.add(source_slide)
                available = source_ids[source_slide]
            missing = sorted(
                int(item["shape_id"]) for item in bindings
                if int(item["shape_id"]) not in available
            )
            if missing:
                errors.append(f"{page_id}: text bindings not present on source slide {source_slide}: {missing}")
            if page.get("page_role", page.get("layout")) not in {"ending"} and not bindings:
                errors.append(f"{page_id}: native page has no original text-box bindings")
            unknown_components = sorted(
                component_id for component_id in page.get("remove_component_ids", [])
                if component_id not in component_ids.get(source_slide, set())
            )
            if unknown_components:
                errors.append(f"{page_id}: unknown removal components on slide {source_slide}: {unknown_components}")
            for addition in page.get("additions", []):
                if not addition.get("fallback_reason"):
                    errors.append(f"{page_id}: every added object requires fallback_reason")
                if not addition.get("binding_ref"):
                    errors.append(f"{page_id}: every added object requires binding_ref")
            addition_refs = {
                item.get("binding_ref") for item in page.get("additions", []) if item.get("binding_ref")
            }
            for image_index, image_binding in enumerate(page.get("image_bindings", [])):
                if image_binding.get("replace_shape_id") is None:
                    binding_ref = f"image_bindings:{image_index}"
                    if binding_ref not in addition_refs:
                        errors.append(
                            f"{page_id}: coordinate-based image {binding_ref} is a new object; "
                            "bind it to an existing picture shape or register an addition with fallback_reason"
                        )
            for mask in page.get("placeholder_masks", []):
                errors.append(
                    f"{page_id}: placeholder masks are forbidden; remove the complete ownership group "
                    "or reuse the complete native component"
                )
        elif mode == "freeform":
            if not page.get("fallback_reason"):
                errors.append(f"{page_id}: freeform page requires fallback_reason")
        else:
            errors.append(f"{page_id}: unknown render_mode={mode}")
        pages.append({
            "page_id": page_id, "render_mode": mode, "source_slide_index": source_slide,
            "text_binding_count": len(bindings), "fallback_reason": page.get("fallback_reason"),
            "removed_components": page.get("remove_component_ids", []),
            "addition_count": len(page.get("additions", [])),
        })
    total = max(1, len(pages))
    freeform_ratio = mode_counts["freeform"] / total
    if freeform_ratio > args.max_freeform_ratio:
        errors.append(f"freeform ratio {freeform_ratio:.1%} exceeds {args.max_freeform_ratio:.1%}")
    if plan.get("template_mode", "template_native") == "template_native" and not used_source_slides:
        errors.append("no original template slide was selected")
    if len(used_source_slides) < min(3, len(prs.slides)) and len(pages) >= 10:
        warnings.append(f"complete deck uses only {len(used_source_slides)} distinct template slides")
    report = {
        "template": str(Path(args.template).resolve()), "pages": pages,
        "mode_counts": dict(mode_counts), "freeform_ratio": freeform_ratio,
        "distinct_source_slides": sorted(used_source_slides),
        "errors": errors, "warnings": warnings, "passed": not errors,
    }
    print(json.dumps({key: report[key] for key in ("mode_counts", "freeform_ratio", "distinct_source_slides", "errors", "warnings", "passed")}, ensure_ascii=False, indent=2))
    if args.output:
        output = Path(args.output)
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    if errors:
        raise SystemExit(1)


if __name__ == "__main__":
    main()
