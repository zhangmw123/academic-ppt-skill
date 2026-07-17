"""Remove legacy navigation backings from an already curated generated PPTX."""

from __future__ import annotations

import argparse
import hashlib
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _box_matches(first: dict[str, float], second: dict[str, float], tolerance: float = 0.012) -> bool:
    return all(abs(first[key] - second[key]) <= tolerance for key in ("left", "top", "width", "height"))


def _navigation_boxes(grammar: dict) -> list[dict[str, float]]:
    boxes = []
    for archetype in grammar.get("archetypes", ()):
        for asset in archetype.get("decorative_assets", ()):
            if asset.get("role") != "navigation":
                continue
            box = {key: float(asset["box"][key]) for key in ("left", "top", "width", "height")}
            if not any(_box_matches(box, existing, tolerance=0.002) for existing in boxes):
                boxes.append(box)
    if not boxes:
        raise ValueError("template grammar contains no navigation assets")
    return boxes


def _normalized_shape_box(shape, presentation) -> dict[str, float]:
    return {
        "left": shape.left / presentation.slide_width,
        "top": shape.top / presentation.slide_height,
        "width": shape.width / presentation.slide_width,
        "height": shape.height / presentation.slide_height,
    }


def repair(
    source_path: Path,
    grammar_path: Path,
    manifest_path: Path,
    output_path: Path,
    output_manifest_path: Path,
    report_path: Path,
) -> dict:
    grammar = json.loads(grammar_path.read_text(encoding="utf-8"))
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    navigation_boxes = _navigation_boxes(grammar)
    presentation = Presentation(source_path)
    manifest_pages = manifest.get("pages", ())
    if len(manifest_pages) != len(presentation.slides):
        raise ValueError("object manifest page count does not match PPTX")

    removals = []
    for slide_number, (slide, manifest_page) in enumerate(
        zip(presentation.slides, manifest_pages),
        1,
    ):
        owned_shape_ids = {
            int(shape_id)
            for region in manifest_page.get("regions", ())
            for shape_id in region.get("owned_shape_ids", ())
        }
        matched = []
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            shape_box = _normalized_shape_box(shape, presentation)
            if any(_box_matches(shape_box, expected) for expected in navigation_boxes):
                if int(shape.shape_id) in owned_shape_ids:
                    raise ValueError(
                        f"slide {slide_number}: navigation backing {shape.shape_id} is also content-owned"
                    )
                matched.append(shape)
        for shape in matched:
            shape_id = int(shape.shape_id)
            box = _normalized_shape_box(shape, presentation)
            shape._element.getparent().remove(shape._element)
            manifest_page["additional_identity_shape_ids"] = [
                int(value)
                for value in manifest_page.get("additional_identity_shape_ids", ())
                if int(value) != shape_id
            ]
            for binding in manifest_page.get("identity_feature_bindings", ()):
                binding["shape_ids"] = [
                    int(value) for value in binding.get("shape_ids", ()) if int(value) != shape_id
                ]
            removals.append({
                "slide_number": slide_number,
                "shape_id": shape_id,
                "name": shape.name,
                "normalized_box": {key: round(value, 4) for key, value in box.items()},
            })
    if not removals:
        raise ValueError("no legacy navigation backings matched the template grammar")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    output_manifest_path.parent.mkdir(parents=True, exist_ok=True)
    output_manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    report = {
        "schema_version": 1,
        "contract": "curated_navigation_backing_repair",
        "source": str(source_path.resolve()),
        "source_sha256": _sha256(source_path),
        "output": str(output_path.resolve()),
        "output_sha256": _sha256(output_path),
        "grammar": str(grammar_path.resolve()),
        "input_manifest": str(manifest_path.resolve()),
        "output_manifest": str(output_manifest_path.resolve()),
        "removed_count": len(removals),
        "removals": removals,
        "passed": True,
    }
    report_path.parent.mkdir(parents=True, exist_ok=True)
    report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    return report


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx")
    parser.add_argument("--grammar", required=True)
    parser.add_argument("--object-manifest", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--output-manifest", required=True)
    parser.add_argument("--report", required=True)
    args = parser.parse_args()
    result = repair(
        Path(args.pptx),
        Path(args.grammar),
        Path(args.object_manifest),
        Path(args.output),
        Path(args.output_manifest),
        Path(args.report),
    )
    print(f"Removed {result['removed_count']} legacy navigation backings")
    print(f"Output: {result['output']}")


if __name__ == "__main__":
    main()
