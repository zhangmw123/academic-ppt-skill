"""Extract native editable slots, layout signatures, and identity from a PPTX template."""

from __future__ import annotations

import argparse
import hashlib
import json
import statistics
from collections import Counter, defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx_utils import iter_shapes


def norm_box(shape, width, height):
    return {
        "left": round(shape.left / width, 4),
        "top": round(shape.top / height, 4),
        "width": round(shape.width / width, 4),
        "height": round(shape.height / height, 4),
    }


def font_size(shape):
    sizes = []
    if not getattr(shape, "has_text_frame", False):
        return 0.0
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.font.size:
            sizes.append(paragraph.font.size.pt)
        sizes.extend(run.font.size.pt for run in paragraph.runs if run.font.size)
    return max(sizes, default=0.0)


def font_family(shape):
    if not getattr(shape, "has_text_frame", False):
        return None
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.name:
                return run.font.name
    return None


def shape_kind(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if getattr(shape, "is_placeholder", False):
        return "placeholder"
    if getattr(shape, "has_text_frame", False):
        return "textbox" if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX else "autoshape"
    return str(shape.shape_type).lower().replace(" ", "_")


def iter_shape_records(shapes, parent_group_shape_id=None):
    for shape in shapes:
        yield shape, parent_group_shape_id
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shape_records(shape.shapes, shape.shape_id)


def cluster_positions(values, tolerance=0.075):
    clusters = []
    for value in sorted(values):
        target = next((cluster for cluster in clusters if abs(statistics.mean(cluster) - value) <= tolerance), None)
        if target is None:
            clusters.append([value])
        else:
            target.append(value)
    return [round(statistics.mean(cluster), 4) for cluster in clusters]


def estimate_capacity(shape, width, height):
    size = font_size(shape) or 14
    box = norm_box(shape, width, height)
    width_inches = box["width"] * width / 914400
    height_inches = box["height"] * height / 914400
    lines = max(1, int(height_inches * 72 / (size * 1.25)))
    chars_per_line = max(4, int(width_inches * 72 / (size * 0.9)))
    return {"estimated_lines": lines, "estimated_cjk_chars": lines * chars_per_line}


def infer_layout_signature(role, text_shapes, pictures, width, height):
    body = [
        shape for shape in text_shapes
        if shape.top / height > 0.18 and shape.height / height > 0.035
        and 0.08 <= shape.width / width <= 0.42
    ]
    centers = cluster_positions([(shape.left + shape.width / 2) / width for shape in body])
    columns = len(centers)
    significant_pictures = [
        shape for shape in pictures
        if 0.035 <= (shape.width / width) * (shape.height / height) <= 0.72 and shape.top / height > 0.1
    ]
    if role in {"cover", "agenda", "section", "ending"}:
        return role, max(1, columns)
    if columns >= 5:
        return "five_columns", columns
    if columns == 4:
        return "four_columns", columns
    if columns == 3:
        return "three_columns", columns
    if columns == 2 and significant_pictures:
        picture_center = statistics.mean((shape.left + shape.width / 2) / width for shape in significant_pictures)
        return ("text_figure_right" if picture_center > 0.5 else "figure_text_right"), columns
    if columns == 2:
        return "two_columns", columns
    if len(significant_pictures) >= 4:
        return "image_grid", max(1, columns)
    if len(significant_pictures) == 1 and significant_pictures[0].width / width > 0.62:
        return "full_figure", max(1, columns)
    return "title_body", max(1, columns)


def text_slots(record, repeated_texts, width, height):
    shapes = record["text_shapes"]
    title = max(
        (shape for shape in shapes if shape.top / height < 0.27),
        key=lambda shape: (font_size(shape), shape.width), default=None,
    )
    body = [shape for shape in shapes if shape is not title and shape.top / height > 0.14]
    column_seeds = [shape for shape in body if 0.08 <= shape.width / width <= 0.42]
    centers = cluster_positions([(shape.left + shape.width / 2) / width for shape in column_seeds])
    slots = []
    for z_order, shape in enumerate(shapes):
        text = shape.text.strip()
        box = norm_box(shape, width, height)
        center = (shape.left + shape.width / 2) / width
        column = min(range(len(centers)), key=lambda index: abs(centers[index] - center)) + 1 if centers else None
        if shape is title:
            slot_role = "page_title"
        elif text in repeated_texts:
            slot_role = "navigation_or_footer"
        elif box["top"] > 0.92 or (text.isdigit() and len(text) <= 3):
            slot_role = "page_number_or_footer"
        elif font_size(shape) >= 16 and box["height"] < 0.16:
            slot_role = f"column_{column}_title" if column else "panel_title"
        else:
            slot_role = f"column_{column}_body" if column else "body"
        slots.append({
            "shape_id": shape.shape_id,
            "name": shape.name,
            "role": slot_role,
            "box": box,
            "font_size_pt": round(font_size(shape), 1) or None,
            "font_family": font_family(shape),
            "sample_text": text,
            "repeated": text in repeated_texts,
            "column": column,
            "z_order": z_order,
            "capacity": estimate_capacity(shape, width, height),
            "editable": True,
        })
    return slots


def classify_slide(index, total, text_shapes, picture_shapes, width, height):
    texts = [shape.text.strip() for shape in text_shapes if shape.text.strip()]
    joined = " ".join(texts).lower()
    max_font = max((font_size(shape) for shape in text_shapes), default=0)
    if index == 0:
        return "cover"
    if index == total - 1 or any(token in joined for token in ("感谢", "thank", "q & a", "致谢")):
        return "ending"
    if any(text.lower() in {"目录", "contents", "content"} for text in texts) and max_font >= 30:
        return "agenda"
    if max_font >= 36 and index <= 1:
        return "cover"
    if len(texts) <= 5 and max_font >= 30:
        return "section"

    significant_pictures = []
    for shape in picture_shapes:
        box = norm_box(shape, width, height)
        area = box["width"] * box["height"]
        if 0.06 <= area <= 0.65 and box["top"] > 0.12:
            significant_pictures.append(box)
    body_text = [shape for shape in text_shapes if shape.top / height > 0.17 and shape.text.strip()]
    centers = [((shape.left + shape.width / 2) / width) for shape in body_text]
    left = sum(center < 0.48 for center in centers)
    right = sum(center > 0.52 for center in centers)
    if significant_pictures:
        picture_center = statistics.mean(box["left"] + box["width"] / 2 for box in significant_pictures)
        if picture_center > 0.57 and left:
            return "text_figure_right"
        if picture_center < 0.43 and right:
            return "text_figure_left"
        if len(significant_pictures) >= 2:
            return "gallery"
        return "full_figure"
    if left >= 2 and right >= 2:
        return "comparison"
    if len(body_text) >= 8:
        return "grid"
    return "title_body"


def median_box(shapes, width, height, fallback):
    if not shapes:
        return fallback
    boxes = [norm_box(shape, width, height) for shape in shapes]
    return {
        key: round(statistics.median(box[key] for box in boxes), 4)
        for key in ("left", "top", "width", "height")
    }


def extract(template: Path, output: Path, asset_dir: Path):
    prs = Presentation(template)
    width, height = prs.slide_width, prs.slide_height
    slide_records = []
    text_occurrences = defaultdict(list)
    picture_digests = Counter()

    for slide_index, slide in enumerate(prs.slides):
        text_shapes, pictures, shape_records = [], [], []
        for shape, parent_group_shape_id in iter_shape_records(slide.shapes):
            shape_records.append((shape, parent_group_shape_id))
            if getattr(shape, "has_text_frame", False):
                text_shapes.append(shape)
                if shape.text.strip():
                    text_occurrences[shape.text.strip()].append((slide_index, shape))
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures.append(shape)
                picture_digests[hashlib.sha256(shape.image.blob).hexdigest()] += 1
        role = classify_slide(slide_index, len(prs.slides), text_shapes, pictures, width, height)
        slide_records.append({
            "slide_index": slide_index + 1,
            "role": role,
            "text_shapes": text_shapes,
            "pictures": pictures,
            "shape_records": shape_records,
            "shape_count": len(list(iter_shapes(slide.shapes))),
        })

    repeated = [
        (text, entries) for text, entries in text_occurrences.items()
        if len({slide_index for slide_index, _ in entries}) >= max(3, len(prs.slides) // 2)
    ]
    repeated_texts = {text for text, _ in repeated}
    repeated_shapes = [shape for _, entries in repeated for _, shape in entries]
    top_nav = [shape for shape in repeated_shapes if shape.top / height < 0.16]
    side_nav = [shape for shape in repeated_shapes if shape.left / width < 0.25 and shape.top / height > 0.15]
    if len(side_nav) >= 3 and len(side_nav) > len(top_nav):
        navigation = "sidebar"
    elif len(top_nav) >= 3:
        navigation = "top"
    else:
        navigation = "minimal"

    navigation_shape_ids_by_slide = defaultdict(set)
    for text, entries in repeated:
        for slide_index, shape in entries:
            in_navigation_region = (
                navigation == "top" and shape.top / height < 0.16
            ) or (
                navigation == "sidebar" and shape.left / width < 0.25 and shape.top / height > 0.12
            )
            if in_navigation_region:
                navigation_shape_ids_by_slide[slide_index].add(shape.shape_id)

    # Large or repeated navigation labels must not turn a content page into an
    # agenda/section page. Reclassify after the navigation region is known.
    for slide_index, record in enumerate(slide_records):
        content_text_shapes = [
            shape for shape in record["text_shapes"]
            if shape.shape_id not in navigation_shape_ids_by_slide[slide_index]
        ]
        record["role"] = classify_slide(
            slide_index, len(prs.slides), content_text_shapes,
            record["pictures"], width, height,
        )

    content_slides = [record for record in slide_records if record["role"] not in {"cover", "ending", "agenda", "section"}]
    title_candidates = []
    for record in content_slides:
        candidates = [shape for shape in record["text_shapes"] if shape.top / height < 0.24 and font_size(shape) >= 18]
        if candidates:
            title_candidates.append(max(candidates, key=font_size))
    title_box = median_box(title_candidates, width, height, {"left": 0.05, "top": 0.13, "width": 0.9, "height": 0.08})

    role_examples = {}
    archetypes = []
    navigation_examples = []
    asset_dir.mkdir(parents=True, exist_ok=True)
    for record in slide_records:
        role_examples.setdefault(record["role"], record["slide_index"])
        assets = []
        picture_slots = []
        for shape in record["pictures"]:
            box = norm_box(shape, width, height)
            area = box["width"] * box["height"]
            digest = hashlib.sha256(shape.image.blob).hexdigest()
            decorative = (
                picture_digests[digest] >= 2
                or area >= 0.42
                or (box["top"] < 0.13 and area < 0.22)
                or (box["left"] < 0.035 and area < 0.22)
                or (box["top"] + box["height"] > 0.93 and area < 0.18)
            )
            logo_candidate = area <= 0.045 and box["top"] <= 0.2 and (box["left"] <= 0.2 or box["left"] + box["width"] >= 0.8)
            picture_slots.append({
                "shape_id": shape.shape_id,
                "name": shape.name,
                "role": "logo" if logo_candidate else ("decoration" if decorative else "content_image"),
                "box": box,
                "area": round(area, 4),
                "replaceable": not decorative or logo_candidate,
                "logo_candidate": logo_candidate,
            })
            if not decorative:
                continue
            try:
                extension = shape.image.ext or "png"
            except ValueError:
                extension = "bin"
            filename = f"slide_{record['slide_index']:02d}_shape_{shape.shape_id}_{digest[:8]}.{extension}"
            target = asset_dir / filename
            if not target.exists():
                target.write_bytes(shape.image.blob)
            try:
                stored_path = str(target.resolve().relative_to(output.parent.resolve()))
            except ValueError:
                stored_path = str(target.resolve())
            assets.append({"path": stored_path, "box": box, "digest": digest, "area": round(area, 4)})
        for slot in picture_slots:
            inner = slot["box"]
            containers = []
            for candidate in picture_slots:
                if candidate["shape_id"] == slot["shape_id"]:
                    continue
                outer = candidate["box"]
                contains = (
                    outer["left"] <= inner["left"] + 0.004
                    and outer["top"] <= inner["top"] + 0.004
                    and outer["left"] + outer["width"] >= inner["left"] + inner["width"] - 0.004
                    and outer["top"] + outer["height"] >= inner["top"] + inner["height"] - 0.004
                    and candidate["area"] > slot["area"] * 1.08
                )
                if contains:
                    containers.append(candidate)
            container = min(containers, key=lambda item: item["area"], default=None)
            slot["nested_within_shape_id"] = container["shape_id"] if container else None
            slot["likely_placeholder_layer"] = bool(
                container and slot["role"] == "content_image" and not slot["logo_candidate"]
            )
        components = []
        component_boxes = []
        for shape, parent_group_shape_id in record["shape_records"]:
            box = norm_box(shape, width, height)
            component_boxes.append((shape, parent_group_shape_id, box))
        for shape, parent_group_shape_id, box in component_boxes:
            area = box["width"] * box["height"]
            containers = []
            for candidate, _, outer in component_boxes:
                if candidate.shape_id == shape.shape_id:
                    continue
                outer_area = outer["width"] * outer["height"]
                if outer_area <= area * 1.08:
                    continue
                if (
                    outer["left"] <= box["left"] + 0.004
                    and outer["top"] <= box["top"] + 0.004
                    and outer["left"] + outer["width"] >= box["left"] + box["width"] - 0.004
                    and outer["top"] + outer["height"] >= box["top"] + box["height"] - 0.004
                ):
                    containers.append((candidate, outer_area))
            geometric_parent = min(containers, key=lambda item: item[1], default=(None, None))[0]
            kind = shape_kind(shape)
            has_text = bool(getattr(shape, "has_text_frame", False) and shape.text.strip())
            role = "navigation" if shape.shape_id in navigation_shape_ids_by_slide[record["slide_index"] - 1] else kind
            if kind == "picture" and geometric_parent is not None:
                role = "nested_picture_placeholder"
            components.append({
                "component_id": f"S{record['slide_index']:02d}_C{shape.shape_id}",
                "shape_ids": [shape.shape_id],
                "shape_kind": kind,
                "component_role": role,
                "box": box,
                "parent_group_shape_id": parent_group_shape_id,
                "nested_within_shape_id": geometric_parent.shape_id if geometric_parent else None,
                "text_owner": kind if has_text else None,
                "sample_text": shape.text.strip() if has_text else "",
                "editable": kind in {"textbox", "autoshape", "placeholder", "picture"},
                "replace_policy": (
                    "edit_text_frame" if has_text else
                    "replace_child_picture_or_remove_child" if role == "nested_picture_placeholder" else
                    "replace_in_place" if kind == "picture" else
                    "preserve"
                ),
                "preserve_geometry": True,
            })
        component_by_shape_id = {
            item["shape_ids"][0]: item for item in components if item.get("shape_ids")
        }
        for item in components:
            parent_shape_id = item.get("parent_group_shape_id") or item.get("nested_within_shape_id")
            parent = component_by_shape_id.get(parent_shape_id)
            item["parent_component_id"] = parent["component_id"] if parent else None
            item["child_component_ids"] = []
        for item in components:
            parent = next(
                (candidate for candidate in components if candidate["component_id"] == item["parent_component_id"]),
                None,
            )
            if parent:
                parent["child_component_ids"].append(item["component_id"])
        signature, column_count = infer_layout_signature(
            record["role"], record["text_shapes"], record["pictures"], width, height
        )
        slots = text_slots(record, repeated_texts, width, height)
        archetypes.append({
            "slide_index": record["slide_index"],
            "role": record["role"],
            "layout_signature": signature,
            "column_count": column_count,
            "shape_count": record["shape_count"],
            "text_slots": slots,
            "picture_slots": picture_slots,
            "nested_picture_slots": [slot for slot in picture_slots if slot["nested_within_shape_id"]],
            "components": components,
            "editable_text_count": len(slots),
            "logo_candidates": [slot for slot in picture_slots if slot["logo_candidate"]],
            "decorative_assets": assets,
        })

        nav_components = [item for item in components if item["component_role"] == "navigation"]
        if nav_components:
            members = []
            for item in sorted(nav_components, key=lambda value: (value["box"]["top"], value["box"]["left"])):
                backing_id = item["nested_within_shape_id"]
                if backing_id is None and item["shape_kind"] in {"autoshape", "placeholder"}:
                    backing_id = item["shape_ids"][0]
                members.append({
                    "position": len(members) + 1,
                    "label_shape_id": item["shape_ids"][0],
                    "backing_shape_ids": [backing_id] if backing_id is not None else [],
                    "text_owner": item["text_owner"],
                    "sample_text": item["sample_text"],
                    "box": item["box"],
                })
            navigation_examples.append({
                "slide_index": record["slide_index"],
                "members": members,
            })

    cover_record = next((record for record in slide_records if record["role"] == "cover"), slide_records[0])
    cover_title = max(cover_record["text_shapes"], key=font_size, default=None)
    ending_record = next((record for record in reversed(slide_records) if record["role"] == "ending"), slide_records[-1])
    ending_title = max(ending_record["text_shapes"], key=font_size, default=None)
    avg_shapes = statistics.mean(record["shape_count"] for record in slide_records)
    picture_count = sum(len(record["pictures"]) for record in slide_records)
    text_count = sum(len(record["text_shapes"]) for record in slide_records)
    panel_mode = "image_scaffold" if picture_count > text_count * 0.35 else "native_shapes"
    composition = "sidebar_canvas" if navigation == "sidebar" else ("top_band_canvas" if navigation == "top" else "minimal_canvas")

    result = {
        "schema_version": 3,
        "source_template": str(template.resolve()),
        "template_name": template.name,
        "slide_size": {"width_inches": round(width / 914400, 3), "height_inches": round(height / 914400, 3)},
        "identity": {
            "navigation": navigation,
            "composition": composition,
            "panel_mode": panel_mode,
            "average_shape_count": round(avg_shapes, 2),
            "picture_to_text_ratio": round(picture_count / max(text_count, 1), 3),
            "signature_features": [
                f"{navigation} navigation",
                f"{composition} composition",
                f"{panel_mode} panel treatment",
                f"{round(avg_shapes)} shapes per source slide",
            ],
        },
        "geometry": {
            "title_box": title_box,
            "cover_title_box": norm_box(cover_title, width, height) if cover_title else None,
            "ending_title_box": norm_box(ending_title, width, height) if ending_title else None,
            "content_region": {"left": 0.25 if navigation == "sidebar" else 0.05, "top": 0.19, "width": 0.71 if navigation == "sidebar" else 0.9, "height": 0.7},
        },
        "role_examples": role_examples,
        "navigation_contract": {
            "orientation": navigation,
            "label_detection": "repetition_plus_geometry",
            "members_example": navigation_examples[0]["members"] if navigation_examples else [],
            "member_slide_index": navigation_examples[0]["slide_index"] if navigation_examples else None,
            "members_by_slide": navigation_examples,
            "active_state": {"rule": "template_dark_fill_only"},
            "inactive_state": {"rule": "template_light_fill"},
            "font": {
                "family": next((font_family(shape) for shape in repeated_shapes if font_family(shape)), None),
                "size_pt": round(statistics.median([font_size(shape) for shape in repeated_shapes if font_size(shape)]), 1) if any(font_size(shape) for shape in repeated_shapes) else None,
                "uniform_across_pages": True,
            },
            "clear_protection": True,
            "reset_state_after_clone": True,
        },
        "archetypes": archetypes,
        "native_editing": {
            "default_mode": "template_native",
            "rule": "Clone the complete source slide and bind content to existing shape IDs.",
            "freeform_requires_reason": True,
        },
    }
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Template grammar: {output}")
    print(json.dumps(result["identity"], ensure_ascii=False))


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("template")
    parser.add_argument("--output", required=True)
    parser.add_argument("--asset-dir", required=True)
    args = parser.parse_args()
    extract(Path(args.template), Path(args.output), Path(args.asset_dir))


if __name__ == "__main__":
    main()
