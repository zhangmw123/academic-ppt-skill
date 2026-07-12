"""Verify that two rendered decks retain materially different template identities."""

from __future__ import annotations

import argparse
import hashlib
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx_utils import iter_shapes


def jaccard(first, second):
    union = first | second
    return len(first & second) / len(union) if union else 1.0


def fingerprint(path):
    prs = Presentation(path)
    media, media_boxes, shape_boxes = set(), set(), set()
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            box = (
                round(shape.left / prs.slide_width, 2), round(shape.top / prs.slide_height, 2),
                round(shape.width / prs.slide_width, 2), round(shape.height / prs.slide_height, 2),
            )
            shape_boxes.add((str(shape.shape_type), *box))
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                media.add(hashlib.sha256(shape.image.blob).hexdigest())
                media_boxes.add(box)
    return {"media": media, "media_boxes": media_boxes, "shape_boxes": shape_boxes}


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("first")
    parser.add_argument("second")
    parser.add_argument("--minimum", type=float, default=0.18)
    parser.add_argument("--output")
    args = parser.parse_args()
    first, second = fingerprint(Path(args.first)), fingerprint(Path(args.second))
    score = (
        0.5 * (1 - jaccard(first["media"], second["media"]))
        + 0.3 * (1 - jaccard(first["media_boxes"], second["media_boxes"]))
        + 0.2 * (1 - jaccard(first["shape_boxes"], second["shape_boxes"]))
    )
    result = {
        "first": str(Path(args.first).resolve()), "second": str(Path(args.second).resolve()),
        "rendered_identity_distance": round(score, 3), "minimum": args.minimum,
        "passed": score >= args.minimum,
        "first_media_assets": len(first["media"]), "second_media_assets": len(second["media"]),
    }
    print(json.dumps(result, ensure_ascii=False, indent=2))
    if args.output:
        target = Path(args.output)
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    if not result["passed"]:
        raise SystemExit(1)


if __name__ == "__main__":
    main()
