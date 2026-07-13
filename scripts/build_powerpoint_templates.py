"""Generate clean, PowerPoint-readable bundled templates for the 1.0 catalog."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


THEMES = (
    ("T01", "green_research", "0F6B5B", "DDF0E9", "2F9E78", "top"),
    ("T02", "blue_research", "155EEF", "E7F0FF", "3B82F6", "top"),
    ("T03", "blue_defense", "0B4F8A", "E5F4FF", "2F80ED", "top"),
    ("T04", "project_application", "0B6E4F", "E6F3EC", "35A276", "top"),
    ("T05", "red_general", "B42318", "FCE8E6", "D92D20", "top"),
    ("T06", "yunnan_purple", "7F56D9", "F3EEFF", "9E77ED", "sidebar"),
    ("T07", "yunnan_red", "B42318", "FCE8E6", "D92D20", "sidebar"),
    ("T08", "yunnan_blue", "175CD3", "E7F0FF", "2E90FA", "sidebar"),
)


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _box(slide, left, top, width, height, fill: str, *, line: str | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line or fill)
    return shape


def _text(slide, left, top, width, height, text: str, size: float, color: str, *, bold: bool = False,
          align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = "Aptos"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = _rgb(color)
    return shape


def _chrome(slide, primary: str, pale: str, accent: str, navigation: str, active: int | None = None):
    if navigation == "sidebar":
        _box(slide, 0, 0, 1.25, 7.5, primary)
        _text(slide, 0.16, 0.35, 0.9, 0.65, "ACADEMIC\nPPT", 11, "FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        for index, label in enumerate(("研究背景", "研究方法", "实验结果", "总结展望")):
            color = "FFFFFF" if index == active else pale
            _text(slide, 0.16, 1.55 + index * 0.74, 0.9, 0.42, label, 9, color,
                  bold=index == active, align=PP_ALIGN.CENTER)
        _box(slide, 1.25, 0, 12.08, 0.12, accent)
        return
    _box(slide, 0, 0, 13.333, 0.62, primary)
    _text(slide, 0.45, 0.09, 2.8, 0.32, "ACADEMIC PPT", 10, "FFFFFF", bold=True)
    for index, label in enumerate(("研究背景", "研究方法", "实验结果", "总结展望")):
        color = "FFFFFF" if index == active else pale
        _text(slide, 6.0 + index * 1.55, 0.12, 1.38, 0.25, label, 8.5, color,
              bold=index == active, align=PP_ALIGN.CENTER)
    _box(slide, 0, 7.38, 13.333, 0.12, accent)


def _title(slide, navigation: str, text: str, primary: str):
    left = 1.65 if navigation == "sidebar" else 0.65
    _text(slide, left, 0.95, 10.8, 0.7, text, 27, primary, bold=True)
    _box(slide, left, 1.68, 1.0, 0.06, primary)


def build_template(output: Path, primary: str, pale: str, accent: str, navigation: str, placeholder_image: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    cover = presentation.slides.add_slide(blank)
    _box(cover, 0, 0, 13.333, 7.5, primary)
    _box(cover, 0.7, 0.85, 0.13, 5.8, accent)
    _text(cover, 1.2, 1.55, 10.5, 1.3, "学术汇报标题", 36, "FFFFFF", bold=True)
    _text(cover, 1.24, 3.08, 8.6, 0.55, "Evidence-grounded academic presentation", 16, pale)
    _text(cover, 1.24, 5.55, 7.5, 0.35, "汇报人 | 指导教师 | 单位 | 日期", 12, pale)

    text_page = presentation.slides.add_slide(blank)
    _chrome(text_page, primary, pale, accent, navigation, 0)
    _title(text_page, navigation, "关键问题决定研究的价值", primary)
    start = 1.65 if navigation == "sidebar" else 0.65
    _box(text_page, start, 2.1, 5.0, 4.35, "FFFFFF", line=pale)
    _box(text_page, start + 5.35, 2.1, 5.65, 4.35, pale)
    _text(text_page, start + 0.35, 2.45, 4.3, 0.55, "证据与论点", 19, primary, bold=True)
    _text(text_page, start + 0.35, 3.25, 4.25, 2.35, "• 来源证据一\n• 来源证据二\n• 解释与边界", 15, "344054")
    text_page.shapes.add_picture(str(placeholder_image), Inches(start + 5.55), Inches(2.3), Inches(5.25), Inches(3.95))
    _text(text_page, start + 5.75, 3.65, 4.85, 0.65, "图表 / 图像证据位", 18, primary, bold=True, align=PP_ALIGN.CENTER)

    comparison = presentation.slides.add_slide(blank)
    _chrome(comparison, primary, pale, accent, navigation, 1)
    _title(comparison, navigation, "方法设计围绕可检验的机制", primary)
    start = 1.65 if navigation == "sidebar" else 0.65
    for index, title in enumerate(("输入", "核心方法", "输出验证")):
        left = start + index * 3.65
        _box(comparison, left, 2.2, 3.15, 3.65, "FFFFFF", line=pale)
        _box(comparison, left, 2.2, 3.15, 0.16, accent)
        _text(comparison, left + 0.28, 2.6, 2.5, 0.48, title, 18, primary, bold=True)
        _text(comparison, left + 0.28, 3.35, 2.45, 1.7, "关键步骤\n可追溯证据\n实施边界", 14, "344054")

    result = presentation.slides.add_slide(blank)
    _chrome(result, primary, pale, accent, navigation, 2)
    _title(result, navigation, "结果支持核心结论，并保留适用边界", primary)
    start = 1.65 if navigation == "sidebar" else 0.65
    _box(result, start, 2.15, 7.1, 4.35, "FFFFFF", line=pale)
    _box(result, start + 7.45, 2.15, 3.55, 4.35, pale)
    _text(result, start + 0.38, 2.55, 6.2, 0.45, "主结果图 / 重建图表", 18, primary, bold=True)
    for index, height in enumerate((1.2, 2.2, 1.65, 3.05)):
        _box(result, start + 0.65 + index * 1.25, 5.5 - height, 0.72, height, accent)
    _text(result, start + 7.8, 2.65, 2.8, 0.5, "解读", 18, primary, bold=True)
    _text(result, start + 7.8, 3.45, 2.65, 1.9, "结果含义\n比较口径\n适用边界", 14, "344054")

    conclusion = presentation.slides.add_slide(blank)
    _chrome(conclusion, primary, pale, accent, navigation, 3)
    _title(conclusion, navigation, "结论、贡献与下一步行动", primary)
    start = 1.65 if navigation == "sidebar" else 0.65
    _box(conclusion, start, 2.25, 10.85, 3.6, pale)
    _text(conclusion, start + 0.55, 2.85, 9.7, 0.55, "一句话核心结论", 25, primary, bold=True, align=PP_ALIGN.CENTER)
    _text(conclusion, start + 1.1, 4.05, 8.6, 0.85, "贡献、边界和可执行的下一步在此清晰收束。", 16, "344054", align=PP_ALIGN.CENTER)

    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)


def main() -> None:
    destination = Path(__file__).resolve().parents[1] / "assets" / "powerpoint_templates"
    destination.mkdir(parents=True, exist_ok=True)
    placeholder_image = destination / "_visual_placeholder.png"
    Image.new("RGB", (1050, 790), "#E4E7EC").save(placeholder_image)
    for template_id, filename, primary, pale, accent, navigation in THEMES:
        output = destination / f"{template_id}_{filename}.pptx"
        build_template(output, primary, pale, accent, navigation, placeholder_image)
        print(output)


if __name__ == "__main__":
    main()
