"""Render a deck by cloning explicit template slides and binding explicit shape IDs."""

from __future__ import annotations

import argparse
import copy
import io
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from pptx_utils import clone_slide, find_shape, iter_shapes, remove_slide, resolve_path


def apply_font_policy(shape, policy: dict | None) -> None:
    if not policy or not getattr(shape, "has_text_frame", False):
        return
    zh_font = policy.get("zh") or policy.get("body_zh")
    latin_font = policy.get("latin")
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            has_cjk = any("\u4e00" <= character <= "\u9fff" for character in run.text)
            primary_font = zh_font if has_cjk and zh_font else latin_font
            properties = run._r.get_or_add_rPr()
            for tag, font_name in (("a:ea", zh_font), ("a:latin", primary_font)):
                if not font_name:
                    continue
                element = properties.find(qn(tag))
                if element is None:
                    element = OxmlElement(tag)
                    properties.append(element)
                element.set("typeface", font_name)


def set_text(shape, content: str, font_policy: dict | None = None) -> None:
    if not getattr(shape, "has_text_frame", False):
        raise ValueError(f"shape {shape.shape_id} has no text frame")
    tf = shape.text_frame
    styles = []
    for paragraph in tf.paragraphs:
        run_properties = None
        if paragraph.runs and paragraph.runs[0]._r.rPr is not None:
            run_properties = copy.deepcopy(paragraph.runs[0]._r.rPr)
        styles.append({
            "paragraph_properties": copy.deepcopy(paragraph._p.pPr) if paragraph._p.pPr is not None else None,
            "run_properties": run_properties,
            "end_properties": copy.deepcopy(paragraph._p.endParaRPr) if paragraph._p.endParaRPr is not None else None,
        })
    tf.clear()
    for index, line in enumerate(content.splitlines() or [""]):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = line
        style = styles[min(index, len(styles) - 1)] if styles else None
        if not style:
            continue
        if paragraph._p.pPr is not None:
            paragraph._p.remove(paragraph._p.pPr)
        if style["paragraph_properties"] is not None:
            paragraph._p.insert(0, copy.deepcopy(style["paragraph_properties"]))
        if paragraph.runs and style["run_properties"] is not None:
            run = paragraph.runs[0]._r
            if run.rPr is not None:
                run.remove(run.rPr)
            run.insert(0, copy.deepcopy(style["run_properties"]))
        if paragraph._p.endParaRPr is not None:
            paragraph._p.remove(paragraph._p.endParaRPr)
        if style["end_properties"] is not None:
            paragraph._p.append(copy.deepcopy(style["end_properties"]))
    apply_font_policy(shape, font_policy)


def add_fitted_picture(slide, image_path: Path, box: dict, fit: str = "contain"):
    left, top = Inches(box["left"]), Inches(box["top"])
    width, height = Inches(box["width"]), Inches(box["height"])
    with Image.open(image_path) as image:
        image_ratio = image.width / image.height
    box_ratio = width / height

    if fit == "cover":
        picture = slide.shapes.add_picture(str(image_path), left, top, width, height)
        if image_ratio > box_ratio:
            visible = box_ratio / image_ratio
            picture.crop_left = picture.crop_right = (1 - visible) / 2
        else:
            visible = image_ratio / box_ratio
            picture.crop_top = picture.crop_bottom = (1 - visible) / 2
        return picture

    if image_ratio > box_ratio:
        final_w, final_h = width, int(width / image_ratio)
    else:
        final_h, final_w = height, int(height * image_ratio)
    return slide.shapes.add_picture(
        str(image_path), left + (width - final_w) // 2, top + (height - final_h) // 2,
        final_w, final_h,
    )


def place_picture_binding(slide, image_path: Path, binding: dict, errors: list[str], page_id: str):
    box = binding.get("box")
    replace_id = binding.get("replace_shape_id")
    target = None
    parent = None
    target_index = None
    if replace_id is not None:
        target = find_shape(slide, int(replace_id))
        if target is None:
            errors.append(f"[{page_id}] image replacement shape not found: {replace_id}")
            return
        box = {
            "left": target.left / 914400, "top": target.top / 914400,
            "width": target.width / 914400, "height": target.height / 914400,
        }
        parent = target._element.getparent()
        target_index = parent.index(target._element)
    if not box:
        errors.append(f"[{page_id}] image binding requires box or replace_shape_id")
        return
    picture = add_fitted_picture(slide, image_path, box, binding.get("fit", "contain"))
    if target is not None and parent is not None and target_index is not None:
        # Preserve original component identity for later semantic bindings and edits.
        for non_visual_properties in picture._element.iterdescendants(qn("p:cNvPr")):
            non_visual_properties.set("id", str(target.shape_id))
            break
        parent.remove(target._element)
        picture_parent = picture._element.getparent()
        picture_parent.remove(picture._element)
        parent.insert(target_index, picture._element)


def apply_shape_adjustments(slide, adjustments: list[dict], errors: list[str], page_id: str):
    for adjustment in adjustments:
        shape_id = int(adjustment["shape_id"])
        shape = find_shape(slide, shape_id)
        if shape is None:
            errors.append(f"[{page_id}] adjustment shape not found: {shape_id}")
            continue
        for key in ("left", "top", "width", "height"):
            if key in adjustment:
                setattr(shape, key, Inches(float(adjustment[key])))
        if adjustment.get("font_size_pt") and getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                paragraph.font.size = Pt(float(adjustment["font_size_pt"]))
                for run in paragraph.runs:
                    run.font.size = Pt(float(adjustment["font_size_pt"]))


def picture_color(shape) -> tuple[int, int, int]:
    try:
        with Image.open(io.BytesIO(shape.image.blob)) as image:
            sample = image.convert("RGB")
            sample.thumbnail((64, 64))
            colors = sample.getcolors(maxcolors=4096) or []
            useful = [(count, color) for count, color in colors if sum(color) < 735]
            return max(useful or colors, default=(1, (230, 235, 232)))[1]
    except Exception:
        return (230, 235, 232)


def navigation_shape_color(shape) -> tuple[int, int, int]:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return picture_color(shape)
    try:
        rgb = shape.fill.fore_color.rgb
        if rgb is not None:
            return tuple(rgb)
    except Exception:
        pass
    return (230, 235, 232)


def navigation_members(contract: dict | None, source_index: int) -> list[dict]:
    if not contract:
        return []
    entries = contract.get("members_by_slide", contract.get("members_by_source_slide", []))
    entry = next((item for item in entries if int(item.get("slide_index", 0)) == source_index), None)
    if entry:
        return entry.get("members", [])
    if int(contract.get("member_slide_index", 0) or 0) == source_index:
        return contract.get("members_example", [])
    return []


def apply_native_navigation_contract(slide, sections: list[str], active: str | None,
                                     contract: dict | None, source_index: int,
                                     font_policy: dict | None, errors: list[str], page_id: str) -> int:
    members = navigation_members(contract, source_index)
    if not members:
        return 0
    if len(members) < len(sections):
        # Let the caller rebuild a complete navigation strip when a source page
        # carries only a subset of the deck's confirmed sections.
        return 0

    sampled = []
    for member in members:
        backing_ids = member.get("backing_shape_ids", [])
        backing = find_shape(slide, int(backing_ids[0])) if backing_ids else None
        label = find_shape(slide, int(member["label_shape_id"]))
        if backing or label:
            sampled.append(navigation_shape_color(backing or label))
    active_color = min(sampled, key=sum, default=(31, 78, 65))
    inactive_color = max(sampled, key=sum, default=(225, 235, 230))
    active_index = sections.index(active) if active in sections else -1
    font = (contract or {}).get("font", {})

    for index, member in enumerate(members):
        related_ids = set(member.get("backing_shape_ids", [])) | {member.get("label_shape_id")}
        related_ids.discard(None)
        if index >= len(sections):
            for shape_id in related_ids:
                shape = find_shape(slide, int(shape_id))
                if shape is not None:
                    shape._element.getparent().remove(shape._element)
            continue

        label = find_shape(slide, int(member["label_shape_id"]))
        if label is None or not getattr(label, "has_text_frame", False):
            errors.append(f"[{page_id}] navigation label shape not found/editable: {member.get('label_shape_id')}")
            continue
        set_text(label, sections[index], font_policy)
        for paragraph in label.text_frame.paragraphs:
            text_color = RGBColor(255, 255, 255) if index == active_index else RGBColor(*active_color)
            if font.get("size_pt"):
                paragraph.font.size = Pt(float(font["size_pt"]))
            paragraph.font.bold = bool(font.get("bold", False))
            paragraph.font.color.rgb = text_color
            for run in paragraph.runs:
                if font.get("size_pt"):
                    run.font.size = Pt(float(font["size_pt"]))
                run.font.bold = bool(font.get("bold", False))
                run.font.color.rgb = text_color

        backing_ids = member.get("backing_shape_ids", [])
        backing = find_shape(slide, int(backing_ids[0])) if backing_ids else label
        if backing is not None and backing.shape_type != MSO_SHAPE_TYPE.PICTURE:
            color = active_color if index == active_index else inactive_color
            try:
                backing.fill.solid()
                backing.fill.fore_color.rgb = RGBColor(*color)
                backing.line.color.rgb = RGBColor(*color)
            except (AttributeError, ValueError):
                errors.append(f"[{page_id}] navigation backing cannot be recolored: {backing.shape_id}")
    return len(members)


def replace_raster_navigation(slide, sections: list[str], active: str | None,
                              slide_width: int, slide_height: int, font_policy: dict | None,
                              navigation_style: dict | None = None):
    if not sections or len(sections) > 6:
        return 0
    # The active tab may be a raster picture while inactive tabs are native
    # rectangles. Select the tallest backing object at each tab position instead
    # of assuming that every tab has the same OOXML shape type.
    candidates = [
        shape for shape in iter_shapes(slide.shapes)
        if 0.03 <= shape.top / slide_height <= 0.13
        and 0.08 <= shape.width / slide_width <= 0.16
        and shape.height / slide_height >= 0.055
        and shape.left / slide_width < 0.83
        and not (getattr(shape, "has_text_frame", False) and shape.text.strip())
    ]
    candidates = sorted(candidates, key=lambda shape: shape.left)
    if len(candidates) < len(sections):
        return 0
    candidates = candidates[:len(sections)]
    # Raster navigation templates often keep only the active tab label as a
    # separate text object. Remove those labels before rebuilding all tabs.
    for text_shape in list(iter_shapes(slide.shapes)):
        if (getattr(text_shape, "has_text_frame", False)
                and text_shape.text.strip()
                and text_shape.top / slide_height < 0.14
                and text_shape.left / slide_width < 0.83):
            text_shape._element.getparent().remove(text_shape._element)
    sampled = [navigation_shape_color(shape) for shape in candidates]
    active_color = min(sampled, key=lambda color: sum(color))
    inactive_color = max(sampled, key=lambda color: sum(color))
    active_index = sections.index(active) if active in sections else -1
    navigation_style = navigation_style or {}
    font_size_pt = float(navigation_style.get("font_size_pt", 8.8))
    for index, (source, label) in enumerate(zip(candidates, sections)):
        parent = source._element.getparent()
        z_index = parent.index(source._element)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, source.left, source.top, source.width, source.height)
        color = active_color if index == active_index else inactive_color
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        shape.line.color.rgb = RGBColor(*color)
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = Inches(0.03)
        frame.margin_top = frame.margin_bottom = Inches(0.01)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.text = label
        paragraph.alignment = PP_ALIGN.CENTER
        # Navigation uses one stable type size. The active section is expressed
        # by the tab fill, not by changing typography between pages.
        paragraph.font.size = Pt(font_size_pt)
        paragraph.font.bold = bool(navigation_style.get("active_bold", False)) and index == active_index
        paragraph.font.color.rgb = RGBColor(255, 255, 255) if index == active_index else RGBColor(*active_color)
        apply_font_policy(shape, font_policy)
        parent.remove(source._element)
        shape_parent = shape._element.getparent()
        shape_parent.remove(shape._element)
        parent.insert(z_index, shape._element)
    return len(candidates)


def render(plan_path: Path, template_path: Path, output_path: Path, page_ids: set[str] | None,
           allow_unconfirmed: bool, include_notes: bool = True):
    plan = json.loads(plan_path.read_text(encoding="utf-8"))
    if not allow_unconfirmed and plan.get("confirmed") is not True:
        raise ValueError("layout_plan.json is not confirmed; use --allow-unconfirmed only for an approved sample render")

    prs = Presentation(template_path)
    original_count = len(prs.slides)
    base_dir = plan_path.resolve().parent
    while base_dir.name.lower() in {"ppt_output", "sample", "final"}:
        base_dir = base_dir.parent
    pages = plan.get("pages", [])
    if page_ids:
        pages = [page for page in pages if page.get("page_id") in page_ids]
    if not pages:
        raise ValueError("No pages selected")

    errors = []
    font_policy = plan.get("font_policy") if isinstance(plan.get("font_policy"), dict) else None
    global_replacements = {
        str(key).strip(): str(value) for key, value in plan.get("global_text_replacements", {}).items()
    }
    navigation_contract = plan.get("navigation_contract", {})
    component_index = {}
    grammar_value = plan.get("template_grammar")
    if grammar_value:
        grammar_path = resolve_path(grammar_value, base_dir)
        if grammar_path.exists():
            grammar = json.loads(grammar_path.read_text(encoding="utf-8"))
            grammar_navigation = grammar.get("navigation_contract", {})
            merged_navigation = dict(grammar_navigation)
            merged_navigation.update(navigation_contract)
            merged_font = dict(grammar_navigation.get("font", {}))
            merged_font.update(navigation_contract.get("font", {}))
            if merged_font:
                merged_navigation["font"] = merged_font
            navigation_contract = merged_navigation
            for archetype in grammar.get("archetypes", []):
                source_slide = int(archetype.get("slide_index", 0))
                for component in archetype.get("components", []):
                    component_index[(source_slide, component.get("component_id"))] = component.get("shape_ids", [])
    for page in pages:
        page_id = page.get("page_id", "UNKNOWN")
        if page.get("render_mode", "template_native") not in {"template_native", "template_adaptive"}:
            errors.append(f"[{page_id}] render_pptx only accepts template_native/template_adaptive pages")
            continue
        source_index = page.get("source_slide_index")
        if not isinstance(source_index, int) or not 1 <= source_index <= original_count:
            errors.append(f"[{page_id}] invalid source_slide_index: {source_index}")
            continue
        slide = clone_slide(prs, prs.slides[source_index - 1])

        remove_ids = list(page.get("remove_shape_ids", []))
        for component_id in page.get("remove_component_ids", []):
            component_shape_ids = component_index.get((source_index, component_id))
            if component_shape_ids is None:
                errors.append(f"[{page_id}] removal component not found: {component_id}")
                continue
            remove_ids.extend(component_shape_ids)
        for remove_id in dict.fromkeys(remove_ids):
            target = find_shape(slide, int(remove_id))
            if target is None:
                errors.append(f"[{page_id}] removal shape not found: {remove_id}")
                continue
            target._element.getparent().remove(target._element)

        bindings = {int(binding["shape_id"]): binding for binding in page.get("text_bindings", [])}
        preserve = {int(value) for value in page.get("preserve_shape_ids", [])}
        nav_members = navigation_members(navigation_contract, source_index)
        for member in nav_members:
            preserve.add(int(member["label_shape_id"]))
            preserve.update(int(value) for value in member.get("backing_shape_ids", []))
        globally_bound = set()
        if global_replacements:
            for shape in iter_shapes(slide.shapes):
                if not getattr(shape, "has_text_frame", False):
                    continue
                replacement = global_replacements.get(shape.text.strip())
                if replacement is not None:
                    set_text(shape, replacement, font_policy)
                    globally_bound.add(shape.shape_id)
        available_ids = {shape.shape_id for shape in iter_shapes(slide.shapes)}
        missing = sorted(set(bindings) - available_ids)
        if missing:
            errors.append(f"[{page_id}] missing text shape IDs on source slide {source_index}: {missing}")
            continue

        if page.get("clear_unbound_text", True):
            for shape in iter_shapes(slide.shapes):
                if (getattr(shape, "has_text_frame", False) and shape.shape_id not in bindings
                        and shape.shape_id not in preserve and shape.shape_id not in globally_bound
                        and shape.text.strip()):
                    set_text(shape, "", font_policy)
        for shape_id, binding in bindings.items():
            set_text(find_shape(slide, shape_id), str(binding.get("content", "")), font_policy)

        apply_shape_adjustments(slide, page.get("adjustments", []), errors, page_id)
        if page.get("placeholder_masks"):
            errors.append(
                f"[{page_id}] placeholder masks are forbidden; remove the complete ownership group "
                "or reuse the complete native component"
            )
        navigation_enabled = page.get("navigation_enabled", True)
        native_navigation_count = apply_native_navigation_contract(
            slide, plan.get("sections", []), page.get("section"), navigation_contract,
            source_index, font_policy, errors, page_id,
        ) if navigation_enabled else 0
        if (navigation_enabled and native_navigation_count == 0
                and page.get("replace_raster_navigation", plan.get("replace_raster_navigation", False))):
            replace_raster_navigation(
                slide, plan.get("sections", []), page.get("section"),
                prs.slide_width, prs.slide_height, font_policy,
                plan.get("navigation_style"),
            )

        for image_binding in page.get("image_bindings", []):
            image_path = resolve_path(image_binding["path"], base_dir)
            if not image_path.exists():
                errors.append(f"[{page_id}] image not found: {image_path}")
                continue
            place_picture_binding(slide, image_path, image_binding, errors, page_id)

        logo_path_value = page.get("logo_path", plan.get("logo_path"))
        logo_binding = page.get("logo_binding")
        if logo_path_value and logo_binding:
            logo_path = resolve_path(logo_path_value, base_dir)
            if not logo_path.exists():
                errors.append(f"[{page_id}] logo not found: {logo_path}")
            else:
                place_picture_binding(slide, logo_path, logo_binding, errors, page_id)

        # Template scaffolds often place a picture panel above editable text slots.
        # Keep all authored text visible without changing its geometry or formatting.
        for shape_id in bindings:
            shape = find_shape(slide, shape_id)
            if shape is None:
                continue
            parent = shape._element.getparent()
            parent.remove(shape._element)
            parent.append(shape._element)

        notes = page.get("speaker_notes", "").strip()
        if notes and include_notes:
            slide.notes_slide.notes_text_frame.text = notes

    if errors:
        raise ValueError("\n".join(errors))

    for _ in range(original_count):
        remove_slide(prs, 0)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"Rendered {len(pages)} slides to {output_path}")


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--layout-plan", "-l", required=True)
    parser.add_argument("--template", "-t", required=True)
    parser.add_argument("--output", "-o", required=True)
    parser.add_argument("--pages", help="Comma-separated page IDs for a sample")
    parser.add_argument("--allow-unconfirmed", action="store_true")
    parser.add_argument("--no-notes", action="store_true", help="Skip speaker notes (diagnostic/fallback)")
    args = parser.parse_args()
    page_ids = {value.strip() for value in args.pages.split(",")} if args.pages else None
    render(Path(args.layout_plan), Path(args.template), Path(args.output), page_ids,
           args.allow_unconfirmed, not args.no_notes)


if __name__ == "__main__":
    main()
