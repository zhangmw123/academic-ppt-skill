"""Learn reusable visual tokens from a rendered PPT template without copying its page slots."""

from __future__ import annotations

import argparse
import json
import re
from collections import Counter
from pathlib import Path

from PIL import Image
from pptx import Presentation

from pptx_utils import iter_shapes


def hex_color(rgb) -> str:
    return "#%02X%02X%02X" % tuple(rgb)


def mix(color: tuple[int, int, int], white_ratio: float) -> tuple[int, int, int]:
    return tuple(round(channel * (1 - white_ratio) + 255 * white_ratio) for channel in color)


def palette_from_preview(preview_dir: Path) -> list[tuple[int, int, int]]:
    images = sorted([*preview_dir.glob("*.PNG"), *preview_dir.glob("*.png")])
    if not images:
        raise ValueError(f"no rendered slide PNGs in {preview_dir}")
    image = Image.open(images[min(1, len(images) - 1)]).convert("RGB")
    image.thumbnail((320, 180), Image.Resampling.LANCZOS)
    quantized = image.quantize(colors=16, method=Image.Quantize.MEDIANCUT).convert("RGB")
    colors = Counter(quantized.getdata())
    return [color for color, _ in colors.most_common(16)]


def palette_from_template(template: Path) -> list[tuple[int, int, int]]:
    """Fallback for templates that python-pptx can parse but PowerPoint cannot preview."""
    prs = Presentation(template)
    colors = Counter()
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            for value in re.findall(r'<a:srgbClr[^>]*val="([0-9A-Fa-f]{6})"', shape._element.xml):
                color = tuple(int(value[index:index + 2], 16) for index in (0, 2, 4))
                colors[color] += 1
    if not colors:
        return [(31, 91, 77), (248, 250, 248), (255, 255, 255)]
    return [color for color, _ in colors.most_common(24)]


def choose_accent(colors: list[tuple[int, int, int]]) -> tuple[int, int, int]:
    candidates = []
    for color in colors:
        maximum, minimum = max(color), min(color)
        saturation = maximum - minimum
        luminance = sum(color) / 3
        if 35 <= luminance <= 190 and saturation >= 20:
            candidates.append((saturation + (190 - luminance) * 0.4, color))
    return max(candidates, default=(0, (31, 91, 77)))[1]


def learn_fonts(template: Path) -> tuple[str, str]:
    prs = Presentation(template)
    east_asian = Counter()
    latin = Counter()
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            xml = shape._element.xml
            east_asian.update(re.findall(r'<a:ea[^>]*typeface="([^"]+)"', xml))
            latin.update(re.findall(r'<a:latin[^>]*typeface="([^"]+)"', xml))
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.font.name:
                    latin[paragraph.font.name] += 1
                for run in paragraph.runs:
                    if run.font.name:
                        latin[run.font.name] += 1
    zh_font = east_asian.most_common(1)[0][0] if east_asian else "Microsoft YaHei"
    if zh_font.startswith("+") or zh_font.lower() in {"arial", "times new roman"}:
        zh_font = "Microsoft YaHei"
    en_font = latin.most_common(1)[0][0] if latin else "Times New Roman"
    return zh_font, en_font


def extract(template: Path, preview_dir: Path) -> dict:
    try:
        colors = palette_from_preview(preview_dir)
        palette_source = "rendered_preview"
    except ValueError:
        colors = palette_from_template(template)
        palette_source = "template_ooxml_fallback"
    accent = choose_accent(colors)
    near_white = next((color for color in colors if sum(color) / 3 > 235), (248, 250, 248))
    zh_font, en_font = learn_fonts(template)
    return {
        "schema_version": 1,
        "source_template": str(template.resolve()),
        "mode": "style_learning",
        "palette_source": palette_source,
        "colors": {
            "background": hex_color(near_white),
            "surface": "#FFFFFF",
            "primary": hex_color(accent),
            "primary_soft": hex_color(mix(accent, 0.82)),
            "primary_pale": hex_color(mix(accent, 0.92)),
            "text": "#1F2933",
            "muted": "#667085",
            "line": "#D9E2DD"
        },
        "fonts": {"title": zh_font, "body": zh_font, "latin": en_font, "fallback": "Microsoft YaHei"},
        "typography": {
            "cover_title": 30, "cover_subtitle": 17, "page_title": 20,
            "section_nav": 10, "panel_title": 15, "body": 13,
            "caption": 9, "footer": 8
        },
        "geometry": {
            "slide_width": 13.333, "slide_height": 7.5,
            "margin_x": 0.62, "nav_y": 0.28, "nav_height": 0.42,
            "title_y": 0.92, "content_y": 1.62, "content_bottom": 6.92
        },
        "learned_principles": [
            "light background with a template-derived accent",
            "thin rules and flat panels rather than heavy shadows",
            "compact top navigation with one active section",
            "wide content canvas and small academic footer"
        ]
    }


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("template")
    parser.add_argument("--preview-dir", required=True)
    parser.add_argument("--output", required=True)
    args = parser.parse_args()
    profile = extract(Path(args.template), Path(args.preview_dir))
    output = Path(args.output)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(json.dumps(profile, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Visual system: {output}")


if __name__ == "__main__":
    main()
