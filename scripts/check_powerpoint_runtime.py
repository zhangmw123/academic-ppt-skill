"""Run repeated PowerPoint COM controls against a minimal image deck and a known candidate."""

from __future__ import annotations

import argparse
import hashlib
import json
import subprocess
import sys
import time
from datetime import datetime, timezone
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from export_preview import export_preview, find_powerpoint, list_slide_images


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def powerpoint_process_count() -> int | None:
    if sys.platform != "win32":
        return None
    completed = subprocess.run(
        [
            "powershell",
            "-NoProfile",
            "-Command",
            "@(Get-Process -Name POWERPNT -ErrorAction SilentlyContinue).Count",
        ],
        capture_output=True,
        text=True,
        timeout=15,
    )
    try:
        return int(completed.stdout.strip()) if completed.returncode == 0 else None
    except ValueError:
        return None


def wait_for_process_baseline(baseline: int | None, timeout_seconds: float = 10.0) -> int | None:
    if baseline is None:
        return None
    deadline = time.monotonic() + timeout_seconds
    current = powerpoint_process_count()
    while current is not None and current > baseline and time.monotonic() < deadline:
        time.sleep(0.25)
        current = powerpoint_process_count()
    return current


def build_minimal_image_control(output_path: Path) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    image_path = output_path.with_suffix(".png")
    image = Image.new("RGB", (960, 360), "white")
    draw = ImageDraw.Draw(image)
    colors = ("#0F766E", "#DC2626", "#2563EB")
    for index, color in enumerate(colors):
        left = index * 320
        draw.rectangle((left, 0, left + 319, 359), fill=color)
    draw.rectangle((300, 120, 660, 240), fill="white")
    draw.text((370, 170), "COM IMAGE CONTROL", fill="black")
    image.save(image_path, "PNG")

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1.0), Inches(1.6), Inches(11.333333), Inches(4.25))
    label = slide.shapes.add_textbox(Inches(1.0), Inches(0.55), Inches(11.333333), Inches(0.55))
    paragraph = label.text_frame.paragraphs[0]
    paragraph.text = "PowerPoint COM image-render control"
    paragraph.font.size = Pt(24)
    paragraph.alignment = PP_ALIGN.CENTER
    presentation.save(output_path)
    return output_path


def inspect_images(
    output_dir: Path,
    expected_count: int,
    *,
    expected_size: tuple[int, int],
    require_color_control: bool,
) -> list[dict]:
    paths = list_slide_images(output_dir)
    if len(paths) != expected_count:
        raise RuntimeError(f"expected {expected_count} slide images; rendered {len(paths)}")
    details = []
    for path in paths:
        with Image.open(path) as image:
            rgb = image.convert("RGB")
            if rgb.size != expected_size:
                raise RuntimeError(f"{path.name} has size {rgb.size}; expected {expected_size}")
            sample = rgb.resize((32, 18))
            colors = sample.getcolors(maxcolors=32 * 18)
            color_count = len(colors) if colors is not None else 32 * 18 + 1
            details.append({
                "name": path.name,
                "width": rgb.width,
                "height": rgb.height,
                "bytes": path.stat().st_size,
                "sample_color_count": color_count,
            })
    if require_color_control and details[0]["sample_color_count"] < 4:
        raise RuntimeError("minimal image control rendered without the expected color variation")
    return details


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("candidate", help="Known-good PPTX used as the realistic render control")
    parser.add_argument("--rounds", type=int, default=3)
    parser.add_argument("--output-root", required=True)
    parser.add_argument("--report", required=True)
    parser.add_argument("--width", type=int, default=1600)
    parser.add_argument("--height", type=int, default=900)
    args = parser.parse_args()

    if args.rounds < 1:
        raise SystemExit("--rounds must be at least 1")
    if not find_powerpoint():
        raise SystemExit("Microsoft PowerPoint is not installed")
    candidate = Path(args.candidate).resolve()
    if not candidate.is_file():
        raise SystemExit(f"candidate not found: {candidate}")
    output_root = Path(args.output_root).resolve()
    if output_root.exists() and any(output_root.iterdir()):
        raise SystemExit(f"output root must be empty or absent: {output_root}")
    output_root.mkdir(parents=True, exist_ok=True)
    report_path = Path(args.report).resolve()
    control = build_minimal_image_control(output_root / "control" / "minimal-image-control.pptx")
    candidate_count = len(Presentation(candidate).slides)
    baseline_processes = powerpoint_process_count()
    report = {
        "schema_version": 1,
        "contract": "repeatable_powerpoint_com_control",
        "started_at": datetime.now(timezone.utc).isoformat(),
        "candidate": {
            "path": str(candidate),
            "sha256": sha256(candidate),
            "slide_count": candidate_count,
        },
        "minimal_control": {
            "path": str(control),
            "sha256": sha256(control),
            "slide_count": 1,
        },
        "rounds_requested": args.rounds,
        "baseline_powerpoint_process_count": baseline_processes,
        "runs": [],
        "passed": False,
    }

    failed = False
    try:
        for round_index in range(1, args.rounds + 1):
            for name, source, expected_count, require_color in (
                ("minimal_image", control, 1, True),
                ("known_candidate", candidate, candidate_count, False),
            ):
                destination = output_root / f"round-{round_index:02d}" / name
                started = time.perf_counter()
                run = {"round": round_index, "target": name, "passed": False}
                try:
                    export_preview(source, destination, "powerpoint", args.width, args.height)
                    run["images"] = inspect_images(
                        destination,
                        expected_count,
                        expected_size=(args.width, args.height),
                        require_color_control=require_color,
                    )
                    residual = wait_for_process_baseline(baseline_processes)
                    run["powerpoint_process_count_after"] = residual
                    if baseline_processes is not None and residual is not None and residual > baseline_processes:
                        raise RuntimeError(
                            f"PowerPoint process count did not return to baseline: {baseline_processes} -> {residual}"
                        )
                    run["passed"] = True
                except Exception as exc:
                    run["error"] = " ".join(str(exc).split())[:1200]
                    failed = True
                finally:
                    run["duration_seconds"] = round(time.perf_counter() - started, 3)
                    report["runs"].append(run)
                print(
                    f"round={round_index} target={name} "
                    f"result={'PASS' if run['passed'] else 'FAIL'} duration={run['duration_seconds']}s"
                )
                if failed:
                    break
            if failed:
                break
    finally:
        report["finished_at"] = datetime.now(timezone.utc).isoformat()
        report["rounds_completed"] = sum(
            1
            for round_index in range(1, args.rounds + 1)
            if len([run for run in report["runs"] if run["round"] == round_index]) == 2
            and all(run["passed"] for run in report["runs"] if run["round"] == round_index)
        )
        report["passed"] = not failed and report["rounds_completed"] == args.rounds
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
        print(f"Report: {report_path}")
    raise SystemExit(1 if not report["passed"] else 0)


if __name__ == "__main__":
    main()
