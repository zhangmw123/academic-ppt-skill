"""Render a non-authoritative PPTX preview from actual shape geometry for diagnosis."""

from __future__ import annotations

import argparse
import io
import json
import math
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def color(value, default):
    try:
        rgb = value.rgb
        return tuple(rgb) if rgb is not None else default
    except Exception:
        return default


def font_for(size: int, chinese: bool):
    candidates = (
        (Path("C:/Windows/Fonts/msyh.ttc"), Path("C:/Windows/Fonts/msyhbd.ttc"))
        if chinese else (Path("C:/Windows/Fonts/arial.ttf"), Path("C:/Windows/Fonts/arialbd.ttf"))
    )
    for path in candidates:
        if path.is_file():
            return ImageFont.truetype(str(path), max(8, size))
    return ImageFont.load_default()


def wrap(draw, text: str, font, width: int):
    lines = []
    for paragraph in text.splitlines() or [""]:
        current = ""
        for character in paragraph:
            trial = current + character
            if current and draw.textbbox((0, 0), trial, font=font)[2] > width:
                lines.append(current)
                current = character
            else:
                current = trial
        lines.append(current)
    return lines


def render_text(canvas, shape, box, scale):
    text = shape.text.strip() if getattr(shape, "has_text_frame", False) else ""
    if not text:
        return
    draw = ImageDraw.Draw(canvas)
    paragraph = shape.text_frame.paragraphs[0]
    run = paragraph.runs[0] if paragraph.runs else None
    size_pt = (run.font.size.pt if run and run.font.size else paragraph.font.size.pt if paragraph.font.size else 16)
    chinese = any("\u4e00" <= char <= "\u9fff" for char in text)
    font_size = max(8, round(size_pt * scale))
    font = font_for(font_size, chinese)
    fill = color(run.font.color if run else paragraph.font.color, (40, 45, 52))
    try:
        margin_x = max(0, round(shape.text_frame.margin_left * (box[2] - box[0]) / max(shape.width, 1)))
        margin_y = max(0, round(shape.text_frame.margin_top * (box[3] - box[1]) / max(shape.height, 1)))
    except Exception:
        margin_x = margin_y = max(1, round(scale))
    available_width = max(8, box[2] - box[0] - margin_x * 2)
    available_height = max(8, box[3] - box[1] - margin_y * 2)
    while True:
        lines = wrap(draw, text, font, available_width)
        line_height = max(8, draw.textbbox((0, 0), "Ag国", font=font)[3] + round(scale))
        if line_height * len(lines) <= available_height or font_size <= 8:
            break
        font_size -= 1
        font = font_for(font_size, chinese)
    total_height = line_height * len(lines)
    y = box[1] + margin_y
    try:
        if shape.text_frame.vertical_anchor is not None and int(shape.text_frame.vertical_anchor) == 3:
            y = box[1] + max(margin_y, (box[3] - box[1] - total_height) // 2)
    except Exception:
        pass
    for line in lines:
        line_width = draw.textbbox((0, 0), line, font=font)[2]
        x = box[0] + margin_x
        try:
            alignment = int(paragraph.alignment) if paragraph.alignment is not None else 1
            if alignment == 2:
                x = box[0] + (box[2] - box[0] - line_width) // 2
            elif alignment == 3:
                x = box[2] - margin_x - line_width
        except Exception:
            pass
        if y + line_height > box[3]:
            break
        draw.text((x, y), line, font=font, fill=fill)
        y += line_height


def render_shape(canvas, shape, sx, sy, text_scale):
    left = round(shape.left * sx)
    top = round(shape.top * sy)
    right = round((shape.left + shape.width) * sx)
    bottom = round((shape.top + shape.height) * sy)
    if right <= left or bottom <= top:
        return
    box = (left, top, right, bottom)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            render_shape(canvas, child, sx, sy, text_scale)
        return
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        try:
            with Image.open(io.BytesIO(shape.image.blob)) as picture:
                picture = picture.convert("RGBA")
                picture.thumbnail((right - left, bottom - top), Image.Resampling.LANCZOS)
                x = left + (right - left - picture.width) // 2
                y = top + (bottom - top - picture.height) // 2
                canvas.alpha_composite(picture, (x, y))
        except Exception:
            ImageDraw.Draw(canvas).rectangle(box, outline=(190, 60, 60), width=2)
        return
    draw = ImageDraw.Draw(canvas)
    try:
        if shape.fill.type is not None:
            draw.rectangle(box, fill=color(shape.fill.fore_color, (255, 255, 255)))
    except Exception:
        pass
    try:
        line_color = color(shape.line.color, None)
        if line_color:
            draw.rectangle(box, outline=line_color, width=max(1, round(text_scale)))
    except Exception:
        pass
    render_text(canvas, shape, box, text_scale)


def contact_sheet(images, output):
    with Image.open(images[0]) as sample:
        width, height = sample.size
    thumb_w = 480
    thumb_h = round(height * thumb_w / width)
    columns = 3
    rows = math.ceil(len(images) / columns)
    sheet = Image.new("RGB", (columns * thumb_w, rows * (thumb_h + 28)), "white")
    draw = ImageDraw.Draw(sheet)
    for index, path in enumerate(images):
        with Image.open(path) as image:
            thumb = image.convert("RGB").resize((thumb_w, thumb_h), Image.Resampling.LANCZOS)
        x = (index % columns) * thumb_w
        y = (index // columns) * (thumb_h + 28)
        sheet.paste(thumb, (x, y))
        draw.text((x + 8, y + thumb_h + 4), f"Slide {index + 1}", fill="black", font=font_for(15, False))
    sheet.save(output, quality=92)


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx")
    parser.add_argument("--output-dir", required=True)
    parser.add_argument("--width", type=int, default=1600)
    args = parser.parse_args()
    presentation = Presentation(args.pptx)
    width = args.width
    height = round(width * presentation.slide_height / presentation.slide_width)
    sx, sy = width / presentation.slide_width, height / presentation.slide_height
    text_scale = width / (presentation.slide_width / 12700)
    output = Path(args.output_dir)
    output.mkdir(parents=True, exist_ok=True)
    images = []
    for index, slide in enumerate(presentation.slides, 1):
        canvas = Image.new("RGBA", (width, height), "white")
        try:
            if slide.background.fill.type is not None:
                canvas.paste(color(slide.background.fill.fore_color, (255, 255, 255)), (0, 0, width, height))
        except Exception:
            pass
        for shape in slide.shapes:
            render_shape(canvas, shape, sx, sy, text_scale)
        path = output / f"slide-{index:03d}.png"
        canvas.convert("RGB").save(path)
        images.append(path)
    contact_sheet(images, output / "contact-sheet.jpg")
    (output / "fallback-preview.json").write_text(json.dumps({
        "authoritative": False,
        "renderer": "python-pptx-shape-diagnostic",
        "slide_count": len(images),
        "pptx": str(Path(args.pptx).resolve()),
    }, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Rendered {len(images)} non-authoritative diagnostic slides to {output}")


if __name__ == "__main__":
    main()
