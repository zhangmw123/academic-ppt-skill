"""Inspect actual PPTX slides and expose stable shape IDs for semantic binding."""

from __future__ import annotations

import argparse
import json
import re
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pptx_utils import emu_to_inches, iter_shapes


SAMPLE_PATTERNS = [
    r"这里(填写|输入)", r"输入(相关|正文|标题|意义|应用)", r"xxxx", r"20xx",
    r"小云集市", r"大云集市", r"click to", r"insert (text|title)", r"lorem ipsum",
]


def normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip().lower()


def representative_font_size(shape):
    sizes = []
    if not getattr(shape, "has_text_frame", False):
        return None
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.font.size:
            sizes.append(paragraph.font.size.pt)
        for run in paragraph.runs:
            if run.font.size:
                sizes.append(run.font.size.pt)
    return round(max(sizes), 1) if sizes else None


def infer_role(index: int, total: int, texts: list[str]) -> str:
    combined = " ".join(texts).lower()
    if index == 1:
        return "cover"
    if index == total or any(k in combined for k in ["感谢", "致谢", "thank you", "questions"]):
        return "ending"
    if any(k in combined for k in ["目录", "contents", "outline"]):
        return "section_or_contents"
    if len(texts) <= 3 and any(k in combined for k in ["第", "chapter", "section"]):
        return "section_divider"
    return "content"


def analyze_template(path: Path) -> dict:
    prs = Presentation(path)
    all_texts = []
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                all_texts.append(normalize_text(shape.text))
    repeated = Counter(all_texts)
    repeat_threshold = max(2, round(len(prs.slides) * 0.3))

    slides = []
    for slide_index, slide in enumerate(prs.slides, 1):
        shapes = []
        slide_texts = []
        for order, shape in enumerate(iter_shapes(slide.shapes), 1):
            entry = {
                "order": order,
                "shape_id": shape.shape_id,
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "box": {
                    "left": emu_to_inches(shape.left), "top": emu_to_inches(shape.top),
                    "width": emu_to_inches(shape.width), "height": emu_to_inches(shape.height),
                    "unit": "inch",
                },
            }
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                norm = normalize_text(text)
                entry.update({
                    "kind": "text",
                    "text": text,
                    "font_size_pt": representative_font_size(shape),
                    "sample_content": bool(text and any(re.search(p, norm, re.I) for p in SAMPLE_PATTERNS)),
                    "repeated_text": bool(norm and repeated[norm] >= repeat_threshold),
                })
                if text:
                    slide_texts.append(text)
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                entry["kind"] = "picture"
            else:
                entry["kind"] = "decoration"
            shapes.append(entry)
        semantic_texts = [
            entry["text"] for entry in shapes
            if entry.get("kind") == "text" and entry.get("text") and not entry.get("repeated_text")
        ]
        slides.append({
            "slide_index": slide_index,
            "role_hint": infer_role(slide_index, len(prs.slides), semantic_texts),
            "layout_name": slide.slide_layout.name,
            "text_summary": slide_texts,
            "shapes": shapes,
        })

    return {
        "schema_version": 2,
        "template_path": str(path.resolve()),
        "template_name": path.name,
        "slide_width_inches": emu_to_inches(prs.slide_width),
        "slide_height_inches": emu_to_inches(prs.slide_height),
        "slide_count": len(prs.slides),
        "master_count": len(prs.slide_masters),
        "slides": slides,
    }


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", nargs="?", help="PPTX template")
    parser.add_argument("--all", action="store_true", help="Inspect all PPTX files in --skill-dir")
    parser.add_argument("--skill-dir", default=str(Path(__file__).resolve().parents[1]))
    parser.add_argument("--output", "-o", required=True, help="Output JSON")
    args = parser.parse_args()

    if args.all:
        root = Path(args.skill_dir)
        files = sorted(p for p in root.rglob("*.pptx") if "ppt_output" not in p.parts and not p.name.startswith("~$"))
    elif args.pptx:
        files = [Path(args.pptx)]
    else:
        parser.error("provide a PPTX path or --all")
    if not files:
        raise SystemExit("No PPTX templates found")

    payload = [analyze_template(path) for path in files]
    output = Path(args.output)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(json.dumps(payload[0] if len(payload) == 1 else payload, ensure_ascii=False, indent=2), encoding="utf-8")
    for item in payload:
        print(f"{item['template_name']}: {item['slide_count']} actual slides inspected")
    print(f"Wrote {output}")


if __name__ == "__main__":
    main()
