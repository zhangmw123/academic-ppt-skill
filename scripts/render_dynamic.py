"""Render an academic deck from a learned visual system and content-driven page archetypes."""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor.from_string(value.upper())


def font_pair(font) -> tuple[str, str]:
    if isinstance(font, dict):
        zh = str(font.get("zh") or font.get("name") or "Microsoft YaHei")
        return zh, str(font.get("latin") or zh)
    value = str(font or "Microsoft YaHei")
    return value, value


def apply_run_fonts(run, font):
    zh_font, latin_font = font_pair(font)
    run.font.name = latin_font
    properties = run._r.get_or_add_rPr()
    for tag, font_name in (("a:ea", zh_font), ("a:latin", latin_font)):
        element = properties.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            properties.append(element)
        element.set("typeface", font_name)


def add_box(slide, x, y, w, h, fill, line=None, radius=True, dashed=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    if dashed:
        shape.line.dash_style = 4
    if radius:
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
    return shape


def add_text(slide, text, x, y, w, h, size, color, font, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.04):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = rgb(color)
    if paragraph.runs:
        run = paragraph.runs[0]
        apply_run_fonts(run, font)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_bullets(slide, bullets, x, y, w, h, style, size=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.06)
    frame.margin_top = frame.margin_bottom = Inches(0.04)
    font_size = size or style["typography"]["body"]
    for index, item in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = rgb(style["colors"]["text"])
        paragraph.space_after = Pt(7)
        paragraph.line_spacing = 1.15
        if paragraph.runs:
            apply_run_fonts(paragraph.runs[0], style["fonts"]["body"])
    return box


def template_mode(style):
    return style.get("_template_identity", {}).get("panel_mode", "native_shapes")


def add_panel(slide, x, y, w, h, style, *, fill=None):
    colors = style["colors"]
    image_scaffold = template_mode(style) == "image_scaffold"
    return add_box(
        slide, x, y, w, h, fill or colors["surface"], colors["line"],
        radius=not image_scaffold, dashed=image_scaffold,
    )


def template_panel_boxes(scaffold, expected):
    """Rebuild combined template frames as separately editable panel geometry."""
    if not scaffold or expected <= 0:
        return []
    candidates = []
    seen = set()
    for item in (*scaffold.get("picture_slots", ()), *scaffold.get("decorative_assets", ())):
        if item.get("role") == "logo":
            continue
        box = item.get("box", {})
        area = float(item.get("area", float(box.get("width", 0)) * float(box.get("height", 0))))
        if float(box.get("top", 0)) < 0.18 or not 0.025 <= area <= 0.30:
            continue
        key = tuple(round(float(box.get(name, 0)), 3) for name in ("left", "top", "width", "height"))
        if key in seen:
            continue
        seen.add(key)
        candidates.append(box)
    if len(candidates) < expected:
        return []
    candidates.sort(key=lambda box: (float(box.get("top", 0)), float(box.get("left", 0))))
    if len(candidates) > expected:
        candidates = sorted(
            candidates,
            key=lambda box: float(box.get("width", 0)) * float(box.get("height", 0)),
            reverse=True,
        )[:expected]
        candidates.sort(key=lambda box: (float(box.get("top", 0)), float(box.get("left", 0))))
    return [
        {
            "x": float(box["left"]) * 13.333,
            "y": float(box["top"]) * 7.5,
            "w": float(box["width"]) * 13.333,
            "h": float(box["height"]) * 7.5,
        }
        for box in candidates
    ]


def template_figure_box(scaffold):
    """Select a half-width source-image slot that leaves room for editable explanation."""
    if not scaffold:
        return None
    candidates = []
    for item in scaffold.get("picture_slots", ()):
        if item.get("role") != "content_image":
            continue
        box = item.get("box", {})
        width = float(box.get("width", 0))
        height = float(box.get("height", 0))
        area = float(item.get("area", width * height))
        if float(box.get("top", 0)) >= 0.18 and 0.08 <= area <= 0.35 and width <= 0.58:
            candidates.append((area, box))
    if not candidates:
        return None
    _, box = max(candidates, key=lambda item: item[0])
    return {
        "x": float(box["left"]) * 13.333,
        "y": float(box["top"]) * 7.5,
        "w": float(box["width"]) * 13.333,
        "h": float(box["height"]) * 7.5,
    }


def add_key_points(slide, items, x, y, w, h, style):
    """Render heading/explanation pairs without exposing internal slot labels."""
    colors = style["colors"]
    gap = 0.12
    item_h = (h - gap * max(0, len(items) - 1)) / max(1, len(items))
    for index, item in enumerate(items):
        if isinstance(item, str):
            title, body = "要点", item
        else:
            title, body = item.get("title", "要点"), item.get("body", "")
        top = y + index * (item_h + gap)
        add_text(slide, title, x, top, w, min(0.3, item_h * 0.34), 12.5,
                 colors["primary"], style["fonts"]["title"], bold=True, margin=0)
        add_text(slide, body, x, top + min(0.34, item_h * 0.38), w,
                 max(0.24, item_h - min(0.34, item_h * 0.38)), 11,
                 colors["text"], style["fonts"]["body"], margin=0)


def add_picture_contain(slide, path: Path, x, y, w, h):
    with Image.open(path) as image:
        ratio = image.width / image.height
    box_ratio = w / h
    if ratio > box_ratio:
        final_w, final_h = w, w / ratio
    else:
        final_h, final_w = h, h * ratio
    return slide.shapes.add_picture(str(path), Inches(x + (w - final_w) / 2),
                                    Inches(y + (h - final_h) / 2), Inches(final_w), Inches(final_h))


def add_picture_cover(slide, path: Path, x, y, w, h):
    with Image.open(path) as image:
        ratio = image.width / image.height
    box_ratio = w / h
    picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    if ratio > box_ratio:
        crop = (1 - box_ratio / ratio) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    elif ratio < box_ratio:
        crop = (1 - ratio / box_ratio) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    return picture


def runtime_geometry(style, grammar=None):
    nav_style = (grammar or {}).get("identity", {}).get("navigation", "top")
    if nav_style == "sidebar":
        return {
            "nav_style": "sidebar", "title_x": 3.02, "title_y": 0.62,
            "title_w": 9.65, "content_x": 3.02, "content_y": 1.42,
            "content_w": 9.65, "content_bottom": 6.9,
        }
    if nav_style == "minimal":
        return {
            "nav_style": "minimal", "title_x": 0.78, "title_y": 0.62,
            "title_w": 11.75, "content_x": 0.78, "content_y": 1.5,
            "content_w": 11.75, "content_bottom": 6.88,
        }
    title = (grammar or {}).get("geometry", {}).get("title_box", {})
    return {
        "nav_style": "top",
        "title_x": max(0.52, title.get("left", 0.05) * 13.333),
        "title_y": max(0.82, title.get("top", 0.13) * 7.5),
        "title_w": min(12.25, title.get("width", 0.9) * 13.333),
        "content_x": 0.65, "content_y": 1.68,
        "content_w": 12.03, "content_bottom": 6.18,
    }


def scaffold_for_page(grammar, page):
    if not grammar:
        return None
    desired = page.get("scaffold_role")
    if not desired:
        desired = {
            "cover": "cover", "agenda": "agenda", "section": "section",
            "ending": "ending", "text_figure": "text_figure_right",
            "comparison": "comparison", "architecture": "title_body",
            "process": "gallery", "points": "grid",
        }.get(page.get("layout"), "title_body")
    candidates = grammar.get("archetypes", [])
    if page.get("scaffold_slide"):
        return next((item for item in candidates if item["slide_index"] == page["scaffold_slide"]), None)
    scaffold = next((item for item in candidates if item["role"] == desired), None)
    if scaffold:
        return scaffold
    return next((item for item in candidates if item["role"] == "title_body"), None)


def add_template_scaffold(slide, scaffold, base_dir=None, *, mode="full", allow_logo=False):
    if not scaffold:
        return False
    added = False
    for asset in scaffold.get("decorative_assets", []):
        path = Path(asset["path"])
        if not path.is_absolute() and base_dir:
            path = Path(base_dir) / path
        if not path.exists():
            continue
        box = asset["box"]
        area = float(asset.get("area", float(box.get("width", 0)) * float(box.get("height", 0))))
        is_top_right_logo = (
            float(box.get("left", 0)) > 0.75
            and float(box.get("top", 1)) < 0.16
            and area < 0.05
        )
        if is_top_right_logo and not allow_logo:
            continue
        matching_slot = next((
            slot for slot in scaffold.get("picture_slots", ())
            if all(abs(float(slot.get("box", {}).get(name, 0)) - float(box.get(name, 0))) < 0.002
                   for name in ("left", "top", "width", "height"))
        ), None)
        asset_role = (matching_slot or {}).get("role")
        if mode == "identity":
            touches_edge = (
                float(box.get("top", 1)) < 0.12
                or float(box.get("top", 0)) + float(box.get("height", 0)) > 0.94
                or float(box.get("left", 1)) < 0.04
            )
            if area > 0.08 or not touches_edge:
                continue
        elif mode == "structure":
            touches_edge = (
                float(box.get("top", 1)) < 0.12
                or float(box.get("top", 0)) + float(box.get("height", 0)) > 0.94
                or float(box.get("left", 1)) < 0.04
            )
            # Content frames baked into one raster are reconstructed below as
            # independent editable shapes. Keep only small edge decoration.
            if area > 0.08 or not touches_edge:
                continue
        slide.shapes.add_picture(
            str(path), Inches(box["left"] * 13.333), Inches(box["top"] * 7.5),
            Inches(box["width"] * 13.333), Inches(box["height"] * 7.5),
        )
        added = True
    return added


def add_navigation(slide, sections, active, style, logo_path=None):
    if not sections:
        return
    geometry = style.get("_runtime", {})
    if geometry.get("nav_style") == "minimal":
        add_text(slide, active or "", geometry.get("title_x", 0.78), 0.28, 2.6, 0.28,
                 style["typography"]["section_nav"], style["colors"]["primary"],
                 style["fonts"]["body"], bold=True, margin=0)
        return
    if geometry.get("nav_style") == "sidebar":
        x, y, w = 0.48, 1.78, 2.05
        gap = 0.13
        h = min(0.62, (4.7 - gap * (len(sections) - 1)) / len(sections))
        for index, section in enumerate(sections):
            active_tab = section == active
            fill = style["colors"]["primary"] if active_tab else style["colors"]["primary_pale"]
            text_color = "#FFFFFF" if active_tab else style["colors"]["primary"]
            add_box(slide, x, y + index * (h + gap), w, h, fill, fill, radius=False)
            add_text(slide, section, x + 0.08, y + index * (h + gap), w - 0.16, h,
                     style["typography"]["section_nav"], text_color, style["fonts"]["body"],
                     bold=active_tab, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
        if logo_path:
            add_picture_contain(slide, Path(logo_path), 0.7, 0.45, 1.55, 0.72)
        return
    x, y, h = 0.52, 0.28, 0.42
    logo_w = 1.15 if logo_path else 0
    available = 12.09 - logo_w
    gap = 0.08
    tab_w = (available - gap * (len(sections) - 1)) / len(sections)
    image_scaffold = template_mode(style) == "image_scaffold"
    for index, section in enumerate(sections):
        active_tab = section == active
        if image_scaffold:
            fill = style["colors"]["primary"] if active_tab else style["colors"]["background"]
            text_color = "#FFFFFF" if active_tab else style["colors"]["text"]
            if active_tab:
                add_box(slide, x + index * (tab_w + gap), y, tab_w, h, fill, fill, radius=False)
        else:
            fill = style["colors"]["primary"] if active_tab else style["colors"]["primary_soft"]
            text_color = "#FFFFFF" if active_tab else style["colors"]["primary"]
            add_box(slide, x + index * (tab_w + gap), y, tab_w, h, fill, fill, radius=False)
        add_text(slide, section, x + index * (tab_w + gap), y, tab_w, h,
                 style["typography"]["section_nav"], text_color, style["fonts"]["body"],
                 bold=active_tab, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    if image_scaffold:
        add_box(slide, x, y + h + 0.04, available, 0.025,
                style["colors"]["primary"], style["colors"]["primary"], radius=False)
    if logo_path:
        add_picture_contain(slide, Path(logo_path), 12.18, 0.23, 0.55, 0.52)


def add_page_conclusion(slide, text, style, y=6.36):
    if not text:
        return
    colors = style["colors"]
    geometry = style.get("_runtime", {})
    x = geometry.get("content_x", 0.68)
    w = geometry.get("content_w", 12.0)
    add_box(slide, x, y, w, 0.46, colors["primary_pale"], colors["primary_soft"], radius=False)
    add_box(slide, x, y, 0.12, 0.46, colors["primary"], colors["primary"], radius=False)
    add_text(slide, text, x + 0.22, y + 0.02, w - 0.45, 0.4, 11, colors["primary"],
             style["fonts"]["body"], bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)


def page_conclusion(page):
    """Support old plans while making the conclusion optional and label-free."""
    return page.get("page_conclusion", page.get("takeaway"))


def add_arrow(slide, x1, y1, x2, y2, color, dashed=False):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(1.6)
    if dashed:
        line.line.dash_style = 4
    line_properties = line._element.spPr.get_or_add_ln()
    tail = line_properties.find(qn("a:tailEnd"))
    if tail is None:
        tail = OxmlElement("a:tailEnd")
        line_properties.append(tail)
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    return line


def add_outline(slide, x, y, w, h, color, *, dashed=True, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.background()
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(1.0)
    if dashed:
        shape.line.dash_style = 4
    if radius:
        try:
            shape.adjustments[0] = 0.04
        except Exception:
            pass
    return shape


def add_chrome(slide, page, page_number, total, plan, style):
    add_navigation(slide, plan.get("sections", []), page.get("section"), style, plan.get("logo_path"))
    geometry = style.get("_runtime", {})
    title_x, title_y = geometry.get("title_x", 0.72), geometry.get("title_y", 0.92)
    add_text(slide, page["title"], title_x, title_y, geometry.get("title_w", 11.9), 0.52,
             style["typography"]["page_title"], style["colors"]["text"],
             style["fonts"]["title"], bold=True, valign=MSO_ANCHOR.MIDDLE)
    marker = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(title_x - 0.16), Inches(title_y + 0.14), Inches(0.12), Inches(0.18))
    marker.fill.solid(); marker.fill.fore_color.rgb = rgb(style["colors"]["primary"])
    marker.line.fill.background()
    footer = plan.get("footer", "")
    if footer:
        add_text(slide, footer, geometry.get("content_x", 0.68), 7.17, 6.5, 0.17,
                 style["typography"]["footer"], style["colors"]["muted"], style["fonts"]["body"])
    add_text(slide, f"{page_number:02d} / {total:02d}", 11.75, 7.15, 0.9, 0.18,
             style["typography"]["footer"], style["colors"]["muted"], style["fonts"]["body"],
             align=PP_ALIGN.RIGHT)


def render_cover(slide, page, plan, style, grammar=None, has_scaffold=False):
    colors, fonts, type_ = style["colors"], style["fonts"], style["typography"]
    if not has_scaffold:
        add_box(slide, 0.0, 0.0, 13.333, 7.5, colors["background"], colors["background"], radius=False)
        add_box(slide, 0.75, 1.25, 0.12, 4.8, colors["primary"], colors["primary"], radius=False)
    title_box = (grammar or {}).get("geometry", {}).get("cover_title_box")
    if title_box:
        tx, ty, tw, th = title_box["left"] * 13.333, title_box["top"] * 7.5, title_box["width"] * 13.333, title_box["height"] * 7.5
    else:
        tx, ty, tw, th = 1.15, 1.65, 10.9, 1.25
    title_color = colors["text"] if has_scaffold else colors["primary"]
    rendered_title_h = max(1.05, th)
    add_text(slide, page["title"], tx, ty, tw, rendered_title_h, type_["cover_title"], title_color, fonts["title"], bold=True,
             align=PP_ALIGN.CENTER if title_box else PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, page.get("subtitle", ""), max(1.0, tx), ty + rendered_title_h + 0.18, min(11.2, tw), 0.7,
             type_["cover_subtitle"], colors["text"], fonts["body"], align=PP_ALIGN.CENTER if title_box else PP_ALIGN.LEFT)
    metadata = page.get("metadata", [])
    meta_y = max(4.15, ty + rendered_title_h + 1.1) if title_box else 4.15
    add_text(slide, "\n".join(metadata), 2.4 if title_box else 1.18, meta_y, 8.6, min(1.45, 7.15 - meta_y), 12, colors["muted"], fonts["body"],
             align=PP_ALIGN.CENTER if title_box else PP_ALIGN.LEFT)
    if plan.get("logo_path"):
        add_picture_contain(slide, Path(plan["logo_path"]), 10.9, 0.55, 1.35, 0.85)


def render_text_figure(slide, page, style, base_dir, scaffold=None):
    colors = style["colors"]
    g = style.get("_runtime", {})
    x0, y0, total_w = g.get("content_x", 0.65), g.get("content_y", 1.68), g.get("content_w", 12.03)
    bottom = g.get("content_bottom", 6.18)
    h = bottom - y0
    gap = 0.24
    image_box = template_figure_box(scaffold)
    if image_box:
        if image_box["x"] >= x0 + total_w / 2:
            text_x, text_w = x0, max(3.4, image_box["x"] - x0 - gap)
        else:
            text_x = image_box["x"] + image_box["w"] + gap
            text_w = max(3.4, x0 + total_w - text_x)
        text_y, text_h = image_box["y"], image_box["h"]
        image_x, image_y, image_w, image_h = image_box["x"], image_box["y"], image_box["w"], image_box["h"]
    else:
        text_x, text_y, text_w, text_h = x0, y0, total_w * 0.42, h
        image_x = x0 + text_w + gap
        image_y, image_w, image_h = y0, total_w - text_w - gap, h
    if template_mode(style) == "image_scaffold":
        frame_x = min(text_x, image_x) - 0.12
        frame_y = min(text_y, image_y) - 0.12
        frame_right = max(text_x + text_w, image_x + image_w) + 0.12
        frame_bottom = max(text_y + text_h, image_y + image_h) + 0.12
        add_outline(slide, frame_x, frame_y, frame_right - frame_x, frame_bottom - frame_y,
                    colors["primary_soft"], dashed=True, radius=True)
    add_panel(slide, text_x, text_y, text_w, text_h, style)
    add_text(slide, page.get("lead", "核心要点"), text_x + 0.27, text_y + 0.24, text_w - 0.54, 0.36, 15,
             colors["primary"], style["fonts"]["title"], bold=True)
    add_key_points(slide, page.get("bullets", []), text_x + 0.27, text_y + 0.72,
                   text_w - 0.54, text_h - 0.96, style)
    add_panel(slide, image_x, image_y, image_w, image_h, style)
    image_path = (base_dir / page["image"]).resolve()
    if page.get("image_fit") == "cover":
        add_picture_cover(slide, image_path, image_x + 0.22, image_y + 0.24, image_w - 0.44, image_h - 0.72)
    else:
        add_picture_contain(slide, image_path, image_x + 0.22, image_y + 0.24, image_w - 0.44, image_h - 0.72)
    if page.get("caption"):
        add_text(slide, page["caption"], image_x + 0.22, image_y + image_h - 0.34,
                 image_w - 0.44, 0.25, style["typography"]["caption"],
                 colors["muted"], style["fonts"]["body"], align=PP_ALIGN.CENTER)
    if template_mode(style) != "image_scaffold" and not image_box:
        add_page_conclusion(slide, page_conclusion(page), style, y=min(6.38, bottom + 0.12))


def media_grid_boxes(count, x, y, w, h, layout=None):
    gap = 0.18
    if layout == "primary_plus_supporting" and count >= 3:
        primary_w = w * 0.58
        support_x = x + primary_w + gap
        support_w = w - primary_w - gap
        support_count = count - 1
        support_columns = 2 if support_count > 1 else 1
        support_rows = (support_count + support_columns - 1) // support_columns
        support_h = (h - gap * (support_rows - 1)) / support_rows
        support_cell_w = (support_w - gap * (support_columns - 1)) / support_columns
        boxes = [{"x": x, "y": y, "w": primary_w, "h": h, "primary": True}]
        for index in range(support_count):
            row, column = divmod(index, support_columns)
            boxes.append({
                "x": support_x + column * (support_cell_w + gap),
                "y": y + row * (support_h + gap),
                "w": support_cell_w,
                "h": support_h,
                "primary": False,
            })
        return boxes
    columns, rows = {
        1: (1, 1), 2: (2, 1), 3: (3, 1), 4: (2, 2), 5: (3, 2), 6: (3, 2),
    }[count]
    cell_w = (w - gap * (columns - 1)) / columns
    cell_h = (h - gap * (rows - 1)) / rows
    boxes = []
    for index in range(count):
        row, column = divmod(index, columns)
        row_offset = (cell_w + gap) / 2 if count == 5 and row == 1 else 0
        boxes.append({
            "x": x + row_offset + column * (cell_w + gap),
            "y": y + row * (cell_h + gap),
            "w": cell_w,
            "h": cell_h,
            "primary": False,
        })
    return boxes


def render_media_gallery(slide, page, style, base_dir):
    colors = style["colors"]
    geometry = style.get("_runtime", {})
    x = geometry.get("content_x", 0.65)
    y = geometry.get("content_y", 1.68)
    w = geometry.get("content_w", 12.03)
    bottom = geometry.get("content_bottom", 6.18)
    items = page.get("media_items", ())
    boxes = media_grid_boxes(len(items), x, y, w, bottom - y, page.get("media_layout"))
    for item, box in zip(items, boxes):
        add_panel(slide, box["x"], box["y"], box["w"], box["h"], style)
        caption_h = 0.28
        padding = 0.16 if not box.get("primary") else 0.2
        image_path = (base_dir / item["path"]).resolve()
        add_picture_contain(
            slide,
            image_path,
            box["x"] + padding,
            box["y"] + padding,
            box["w"] - padding * 2,
            max(0.35, box["h"] - caption_h - padding * 2),
        )
        add_text(
            slide,
            item.get("caption", ""),
            box["x"] + padding,
            box["y"] + box["h"] - caption_h - 0.04,
            box["w"] - padding * 2,
            caption_h,
            9,
            colors["muted"],
            style["fonts"]["body"],
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
            margin=0,
        )
    if page_conclusion(page):
        add_page_conclusion(slide, page_conclusion(page), style, y=min(6.38, bottom + 0.12))


def render_module_media(slide, page, style, base_dir):
    colors = style["colors"]
    geometry = style.get("_runtime", {})
    x = geometry.get("content_x", 0.65)
    y = geometry.get("content_y", 1.68)
    w = geometry.get("content_w", 12.03)
    bottom = geometry.get("content_bottom", 6.18)
    modules = page.get("modules", ())
    boxes = media_grid_boxes(len(modules), x, y, w, bottom - y)
    for module, box in zip(modules, boxes):
        add_panel(slide, box["x"], box["y"], box["w"], box["h"], style)
        header_h = 0.38
        body_h = 0.64 if box["h"] >= 3.0 else 0.48
        caption_h = 0.25
        add_text(
            slide, module["title"], box["x"] + 0.18, box["y"] + 0.14,
            box["w"] - 0.36, header_h, 14, colors["primary"], style["fonts"]["title"],
            bold=True, margin=0,
        )
        add_text(
            slide, module["body"], box["x"] + 0.18, box["y"] + 0.54,
            box["w"] - 0.36, body_h, 11, colors["text"], style["fonts"]["body"], margin=0,
        )
        image_y = box["y"] + 0.58 + body_h
        image_h = box["h"] - (image_y - box["y"]) - caption_h - 0.28
        add_picture_contain(
            slide,
            (base_dir / module["image"]).resolve(),
            box["x"] + 0.18,
            image_y,
            box["w"] - 0.36,
            max(0.3, image_h),
        )
        add_text(
            slide, module.get("caption", ""), box["x"] + 0.18,
            box["y"] + box["h"] - caption_h - 0.05, box["w"] - 0.36, caption_h,
            9, colors["muted"], style["fonts"]["body"], align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE, margin=0,
        )
    if page_conclusion(page):
        add_page_conclusion(slide, page_conclusion(page), style, y=min(6.38, bottom + 0.12))


def render_comparison(slide, page, style):
    colors = style["colors"]
    g = style.get("_runtime", {})
    x0, y0, total_w = g.get("content_x", 0.64), g.get("content_y", 1.68), g.get("content_w", 12.05)
    bottom = g.get("content_bottom", 6.18)
    columns = page["columns"]
    gap = 0.28
    width = (total_w - gap * (len(columns) - 1)) / len(columns)
    for index, column in enumerate(columns):
        x = x0 + index * (width + gap)
        add_box(slide, x, y0, width, bottom - y0, colors["surface"], colors["line"])
        add_box(slide, x, y0, width, 0.48, colors["primary"] if index == 0 else colors["primary_soft"],
                colors["primary"] if index == 0 else colors["primary_soft"], radius=False)
        add_text(slide, column["title"], x + 0.12, y0 + 0.02, width - 0.24, 0.42, 15,
                 "#FFFFFF" if index == 0 else colors["primary"], style["fonts"]["title"],
                 bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if column.get("lead"):
            add_text(slide, column["lead"], x + 0.22, y0 + 0.7, width - 0.44, 0.65, 14,
                     colors["text"], style["fonts"]["title"], bold=True)
            y = y0 + 1.42
        else:
            y = y0 + 0.74
        add_bullets(slide, column.get("bullets", []), x + 0.22, y, width - 0.44, bottom - y - 0.18, style)
    add_page_conclusion(slide, page_conclusion(page), style, y=min(6.38, bottom + 0.12))


def render_process(slide, page, style, base_dir, scaffold=None):
    colors = style["colors"]
    g = style.get("_runtime", {})
    x0, y0, total_w = g.get("content_x", 0.64), g.get("content_y", 1.68), g.get("content_w", 12.05)
    bottom = g.get("content_bottom", 6.18)
    if page.get("image"):
        add_box(slide, x0 + 0.5, y0, total_w - 1.0, 2.55, colors["surface"], colors["line"])
        add_picture_contain(slide, (base_dir / page["image"]).resolve(), x0 + 0.7, y0 + 0.16, total_w - 1.4, 2.2)
    steps = page["steps"]
    native_boxes = template_panel_boxes(scaffold, len(steps))
    gap = 0.13
    width = (total_w - gap * (len(steps) - 1)) / len(steps)
    if page.get("image"):
        step_y = y0 + 2.82
        step_h = max(1.15, bottom - step_y)
    else:
        step_h = 2.35 if len(steps) <= 4 else 2.15
        step_y = y0 + max(0.35, (bottom - y0 - step_h) / 2)
    for index, step in enumerate(steps):
        if native_boxes:
            box = native_boxes[index]
            x, current_y, current_w, current_h = box["x"], box["y"], box["w"], box["h"]
        else:
            x, current_y, current_w, current_h = x0 + index * (width + gap), step_y, width, step_h
        module = step if isinstance(step, dict) else {"title": f"步骤 {index + 1}", "body": str(step)}
        add_panel(slide, x, current_y, current_w, current_h, style)
        add_box(slide, x, current_y, current_w, 0.44, colors["primary"], colors["primary"], radius=False)
        add_text(slide, f"{index + 1:02d}  {module.get('title', '')}", x, current_y, current_w, 0.42, 13, "#FFFFFF",
                 style["fonts"]["body"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
        add_text(slide, module.get("body", ""), x + 0.1, current_y + 0.52,
                 current_w - 0.2, current_h - 0.65, 11, colors["text"],
                 style["fonts"]["body"], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    if native_boxes and template_mode(style) == "image_scaffold":
        centers = [box["x"] + box["w"] / 2 for box in native_boxes]
        top = min(box["y"] for box in native_boxes)
        circle_y, circle_size = max(y0, top - 0.88), 0.58
        frame_x = min(box["x"] for box in native_boxes) - 0.12
        frame_right = max(box["x"] + box["w"] for box in native_boxes) + 0.12
        frame_bottom = max(box["y"] + box["h"] for box in native_boxes) + 0.12
        add_outline(slide, frame_x, circle_y - 0.1, frame_right - frame_x,
                    frame_bottom - circle_y + 0.1, colors["primary_soft"], dashed=True, radius=True)
        for index, center_x in enumerate(centers):
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(center_x - circle_size / 2), Inches(circle_y),
                Inches(circle_size), Inches(circle_size),
            )
            circle.fill.solid(); circle.fill.fore_color.rgb = rgb("#FFFFFF")
            circle.line.color.rgb = rgb(colors["primary"]); circle.line.width = Pt(1.5)
            add_text(slide, str(index + 1), center_x - circle_size / 2, circle_y,
                     circle_size, circle_size, 12, colors["primary"], style["fonts"]["title"],
                     bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
            if index < len(centers) - 1:
                add_arrow(slide, center_x + circle_size / 2 + 0.08, circle_y + circle_size / 2,
                          centers[index + 1] - circle_size / 2 - 0.08, circle_y + circle_size / 2,
                          colors["primary"])
    if template_mode(style) != "image_scaffold" and not native_boxes:
        add_page_conclusion(slide, page_conclusion(page), style, y=min(6.38, bottom + 0.12))


def render_points(slide, page, style, scaffold=None):
    colors = style["colors"]
    g = style.get("_runtime", {})
    x0, y0, total_w = g.get("content_x", 0.68), g.get("content_y", 1.7), g.get("content_w", 12.02)
    bottom = g.get("content_bottom", 6.18)
    points = page["points"]
    native_boxes = template_panel_boxes(scaffold, len(points))
    gap_x, gap_y = 0.3, 0.24
    box_w = (total_w - gap_x) / 2
    box_h = (bottom - y0 - gap_y) / 2
    for index, point in enumerate(points):
        if native_boxes:
            box = native_boxes[index]
            x, y, current_w, current_h = box["x"], box["y"], box["w"], box["h"]
        else:
            row, column = divmod(index, 2)
            x, y = x0 + column * (box_w + gap_x), y0 + row * (box_h + gap_y)
            current_w, current_h = box_w, box_h
        add_panel(slide, x, y, current_w, current_h, style)
        add_box(slide, x, y, 0.11, current_h, colors["primary"], colors["primary"], radius=False)
        add_text(slide, point["title"], x + 0.3, y + 0.24, current_w - 0.55, 0.38, 14,
                 colors["primary"], style["fonts"]["title"], bold=True)
        add_text(slide, point["body"], x + 0.3, y + 0.72, current_w - 0.55, current_h - 0.9, 11.5,
                 colors["text"], style["fonts"]["body"])
    if template_mode(style) != "image_scaffold" and not native_boxes:
        add_page_conclusion(slide, page_conclusion(page), style, y=min(6.38, bottom + 0.12))


def render_architecture(slide, page, style, scaffold=None):
    """Render a left-to-right editable architecture from structured columns and edges."""
    colors = style["colors"]
    architecture = page["architecture"]
    columns = architecture["columns"]
    g = style.get("_runtime", {})
    gap = 0.28
    total_w = g.get("content_w", 12.05)
    col_w = (total_w - gap * (len(columns) - 1)) / len(columns)
    x0, top = g.get("content_x", 0.64), g.get("content_y", 1.68)
    bottom = g.get("content_bottom", 6.02)
    positions = {}
    node_records = []
    native_boxes = template_panel_boxes(scaffold, len(columns))

    for col_index, column in enumerate(columns):
        if native_boxes:
            box = native_boxes[col_index]
            x, column_top, current_w, current_h = box["x"], box["y"], box["w"], box["h"]
            column_bottom = column_top + current_h
        else:
            x, column_top, current_w, current_h = x0 + col_index * (col_w + gap), top, col_w, bottom - top
            column_bottom = bottom
        add_panel(slide, x, column_top, current_w, current_h, style)
        add_box(slide, x, column_top, current_w, 0.46, colors["primary_soft"], colors["primary_soft"], radius=False)
        add_text(slide, column["title"], x + 0.08, column_top + 0.01, current_w - 0.16, 0.42, 13,
                 colors["primary"], style["fonts"]["title"], bold=True,
                 align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
        nodes = column.get("nodes", [])
        node_gap = 0.22
        available_h = column_bottom - column_top - 0.78
        node_h = min(0.9, (available_h - node_gap * max(0, len(nodes) - 1)) / max(1, len(nodes)))
        stack_h = len(nodes) * node_h + max(0, len(nodes) - 1) * node_gap
        node_y = column_top + 0.62 + max(0, (available_h - stack_h) / 2)
        for node_index, node in enumerate(nodes):
            y = node_y + node_index * (node_h + node_gap)
            positions[node["id"]] = {
                "left": x + 0.16, "right": x + current_w - 0.16,
                "top": y, "bottom": y + node_h,
                "center_x": x + current_w / 2, "center_y": y + node_h / 2,
            }
            node_records.append((node, x, y, node_h, current_w))

    for edge in architecture.get("edges", []):
        source, target = positions.get(edge["source"]), positions.get(edge["target"])
        if not source or not target:
            continue
        same_column = abs(source["center_x"] - target["center_x"]) < 0.01
        if same_column:
            downward = source["center_y"] <= target["center_y"]
            x1 = x2 = source["center_x"]
            y1 = source["bottom"] if downward else source["top"]
            y2 = target["top"] if downward else target["bottom"]
        else:
            forward = source["center_x"] < target["center_x"]
            x1 = source["right"] if forward else source["left"]
            x2 = target["left"] if forward else target["right"]
            y1, y2 = source["center_y"], target["center_y"]
        add_arrow(slide, x1, y1, x2, y2,
                  colors["primary"], edge.get("style") == "feedback")
        if edge.get("label"):
            mid_x = (x1 + x2) / 2
            mid_y = (y1 + y2) / 2
            add_text(slide, edge["label"], mid_x - 0.42, mid_y - 0.22, 0.84, 0.22, 8.5,
                     colors["muted"], style["fonts"]["body"], align=PP_ALIGN.CENTER,
                     valign=MSO_ANCHOR.MIDDLE, margin=0)

    # Draw nodes after connectors so lines never obscure labels or run through boxes.
    for node, x, y, node_h, current_w in node_records:
        add_panel(slide, x + 0.16, y, current_w - 0.32, node_h, style, fill="#FFFFFF")
        add_text(slide, node["label"], x + 0.26, y + 0.08, current_w - 0.52,
                 0.32 if node.get("detail") else node_h - 0.16, 12, colors["text"],
                 style["fonts"]["body"], bold=True, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE, margin=0)
        if node.get("detail"):
            add_text(slide, node["detail"], x + 0.26, y + 0.43, current_w - 0.52,
                     max(0.22, node_h - 0.49), 10.5, colors["muted"], style["fonts"]["body"],
                     align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
    if template_mode(style) != "image_scaffold" and not native_boxes:
        add_page_conclusion(slide, page_conclusion(page), style, y=min(6.38, bottom + 0.12))


def render_agenda(slide, page, plan, style):
    colors = style["colors"]
    g = style.get("_runtime", {})
    sections = page.get("items", plan.get("sections", []))
    x0, y0, total_w = g.get("content_x", 0.68), g.get("content_y", 1.72), g.get("content_w", 12.0)
    rows = len(sections)
    gap = 0.16
    row_h = min(0.78, (4.65 - gap * max(0, rows - 1)) / max(1, rows))
    for index, item in enumerate(sections):
        label = item["title"] if isinstance(item, dict) else item
        note = item.get("note", "") if isinstance(item, dict) else ""
        y = y0 + index * (row_h + gap)
        add_text(slide, f"{index + 1:02d}", x0, y, 0.72, row_h, 18, colors["primary"],
                 style["fonts"]["title"], bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)
        add_box(slide, x0 + 0.82, y, total_w - 0.82, row_h, colors["surface"], colors["line"], radius=False)
        add_text(slide, label, x0 + 1.08, y, 4.1, row_h, 15, colors["text"], style["fonts"]["title"],
                 bold=True, valign=MSO_ANCHOR.MIDDLE, margin=0)
        if note:
            add_text(slide, note, x0 + 5.15, y, total_w - 5.45, row_h, 11, colors["muted"],
                     style["fonts"]["body"], align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE, margin=0)


def render_section(slide, page, style):
    colors = style["colors"]
    add_text(slide, page.get("number", ""), 0.9, 1.45, 2.0, 1.0, 34, colors["primary_soft"],
             style["fonts"]["title"], bold=True)
    add_text(slide, page["title"], 1.1, 2.55, 11.0, 0.9, 30, colors["primary"],
             style["fonts"]["title"], bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_box(slide, 4.75, 3.62, 3.85, 0.04, colors["primary"], colors["primary"], radius=False)
    add_text(slide, page.get("subtitle", ""), 2.0, 4.05, 9.35, 0.72, 15, colors["text"],
             style["fonts"]["body"], align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def render_full_figure(slide, page, style, base_dir):
    colors = style["colors"]
    g = style.get("_runtime", {})
    x0, y0, total_w = g.get("content_x", 0.68), g.get("content_y", 1.68), g.get("content_w", 12.0)
    bottom = g.get("content_bottom", 6.18)
    add_box(slide, x0, y0, total_w, bottom - y0, colors["surface"], colors["line"])
    image_path = (base_dir / page["image"]).resolve()
    add_picture_contain(slide, image_path, x0 + 0.22, y0 + 0.2, total_w - 0.44, bottom - y0 - 0.68)
    if page.get("caption"):
        add_text(slide, page["caption"], x0 + 0.22, bottom - 0.36, total_w - 0.44, 0.24,
                 style["typography"]["caption"], colors["muted"], style["fonts"]["body"],
                 align=PP_ALIGN.CENTER)
    add_page_conclusion(slide, page_conclusion(page), style, y=min(6.38, bottom + 0.12))


def render_ending(slide, page, style, grammar=None, has_scaffold=False):
    colors = style["colors"]
    if not has_scaffold:
        add_box(slide, 0, 0, 13.333, 7.5, colors["background"], colors["background"], radius=False)
        add_box(slide, 1.0, 1.15, 11.3, 5.2, colors["surface"], colors["line"])
    title_box = (grammar or {}).get("geometry", {}).get("ending_title_box")
    if title_box:
        x, y, w, h = title_box["left"] * 13.333, title_box["top"] * 7.5, title_box["width"] * 13.333, title_box["height"] * 7.5
    else:
        x, y, w, h = 1.65, 2.25, 10.0, 1.0
    add_text(slide, page["title"], x, y, w, max(0.9, h), 28, colors["primary"],
             style["fonts"]["title"], bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, page.get("subtitle", ""), 2.05, y + max(1.2, h), 9.2, 1.0, 16, colors["text"],
             style["fonts"]["body"], align=PP_ALIGN.CENTER)


def render(plan_path: Path, style_path: Path, output: Path, selected_page_ids: set[str] | None = None):
    plan = json.loads(plan_path.read_text(encoding="utf-8"))
    style = json.loads(style_path.read_text(encoding="utf-8"))
    latin_font = style.get("fonts", {}).get("latin")
    for role in ("title", "body"):
        current = style.get("fonts", {}).get(role)
        if isinstance(current, str):
            style["fonts"][role] = {"zh": current, "latin": latin_font or current}
    grammar = None
    if plan.get("template_grammar"):
        grammar_path = Path(plan["template_grammar"])
        if not grammar_path.is_absolute():
            grammar_path = Path.cwd() / grammar_path
        grammar = json.loads(grammar_path.read_text(encoding="utf-8"))
        grammar["_base_dir"] = str(grammar_path.parent.resolve())
        style["_template_identity"] = dict(grammar.get("identity", {}))
    style["_runtime"] = runtime_geometry(style, grammar)
    if not plan.get("confirmed"):
        raise ValueError("dynamic plan is not confirmed")
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = min(prs.slide_layouts, key=lambda layout: len(layout.placeholders))
    pages = plan["pages"]
    if selected_page_ids:
        pages = [page for page in pages if page.get("page_id") in selected_page_ids]
    base_dir = Path(plan.get("asset_base_dir", plan_path.resolve().parent.parent)).resolve()
    for index, page in enumerate(pages, 1):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid(); slide.background.fill.fore_color.rgb = rgb(style["colors"]["background"])
        layout = page["layout"]
        scaffold_mode = page.get("use_template_scaffold")
        scaffold = scaffold_for_page(grammar, page) if scaffold_mode else None
        has_scaffold = add_template_scaffold(
            slide,
            scaffold,
            grammar.get("_base_dir") if grammar else None,
            mode="identity" if scaffold_mode == "identity" else "full",
            allow_logo=bool(plan.get("logo_path")),
        )
        if layout == "cover":
            render_cover(slide, page, plan, style, grammar, has_scaffold)
        elif layout == "ending":
            render_ending(slide, page, style, grammar, has_scaffold)
        elif layout == "section":
            render_section(slide, page, style)
        else:
            add_chrome(slide, page, index, len(pages), plan, style)
            if layout == "agenda":
                render_agenda(slide, page, plan, style)
            elif layout == "text_figure":
                render_text_figure(slide, page, style, base_dir, scaffold)
            elif layout == "media_gallery":
                render_media_gallery(slide, page, style, base_dir)
            elif layout == "module_media":
                render_module_media(slide, page, style, base_dir)
            elif layout == "full_figure":
                render_full_figure(slide, page, style, base_dir)
            elif layout == "comparison":
                render_comparison(slide, page, style)
            elif layout == "process":
                render_process(slide, page, style, base_dir, scaffold)
            elif layout == "points":
                render_points(slide, page, style, scaffold)
            elif layout == "architecture":
                render_architecture(slide, page, style, scaffold)
            else:
                raise ValueError(f"unsupported dynamic layout: {layout}")
        if page.get("speaker_notes"):
            slide.notes_slide.notes_text_frame.text = page["speaker_notes"]
    # Remove the default first slide only if Presentation() supplied one.
    while len(prs.slides) > len(pages):
        from pptx_utils import remove_slide
        remove_slide(prs, 0)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"Rendered {len(pages)} dynamic slides to {output}")


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--plan", required=True)
    parser.add_argument("--visual-system", required=True)
    parser.add_argument("--output", required=True)
    parser.add_argument("--pages", help="Comma-separated page IDs for representative samples")
    args = parser.parse_args()
    selected = {value.strip() for value in args.pages.split(",")} if args.pages else None
    render(Path(args.plan), Path(args.visual_system), Path(args.output), selected)


if __name__ == "__main__":
    main()
