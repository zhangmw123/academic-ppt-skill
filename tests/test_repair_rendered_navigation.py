import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


SCRIPTS = Path(__file__).resolve().parents[1] / "scripts"
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

from repair_rendered_navigation import repair


def test_repairs_only_legacy_navigation_backings_and_updates_manifest(tmp_path: Path):
    image = tmp_path / "old-active-tab.png"
    Image.new("RGB", (320, 80), "darkgreen").save(image)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    old_backing = slide.shapes.add_picture(
        str(image), Inches(0.5), Inches(0.375), Inches(2.0), Inches(0.525)
    )
    content = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(4), Inches(2))
    source = tmp_path / "curated.pptx"
    presentation.save(source)
    grammar = tmp_path / "grammar.json"
    grammar.write_text(json.dumps({
        "archetypes": [{
            "decorative_assets": [{
                "role": "navigation",
                "box": {"left": 0.05, "top": 0.05, "width": 0.2, "height": 0.07},
            }],
        }],
    }), encoding="utf-8")
    manifest = tmp_path / "manifest.json"
    manifest.write_text(json.dumps({
        "pages": [{
            "slide_number": 1,
            "additional_identity_shape_ids": [old_backing.shape_id, content.shape_id],
            "identity_feature_bindings": [{
                "feature": "top navigation",
                "shape_ids": [old_backing.shape_id, content.shape_id],
            }],
            "regions": [{"owned_shape_ids": [content.shape_id]}],
        }],
    }), encoding="utf-8")
    output = tmp_path / "repaired.pptx"
    output_manifest = tmp_path / "repaired-manifest.json"

    report = repair(
        source,
        grammar,
        manifest,
        output,
        output_manifest,
        tmp_path / "repair-report.json",
    )

    repaired = Presentation(output)
    assert not any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in repaired.slides[0].shapes)
    assert any(shape.shape_id == content.shape_id for shape in repaired.slides[0].shapes)
    repaired_manifest = json.loads(output_manifest.read_text(encoding="utf-8"))
    page = repaired_manifest["pages"][0]
    assert old_backing.shape_id not in page["additional_identity_shape_ids"]
    assert old_backing.shape_id not in page["identity_feature_bindings"][0]["shape_ids"]
    assert content.shape_id in page["additional_identity_shape_ids"]
    assert report["removed_count"] == 1
