import hashlib
import json
from pathlib import Path
import subprocess
import sys

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from academic_ppt.object_qa import ObjectLevelQualityGate, TemplateIdentityDifferenceGate
from academic_ppt.template_semantics import (
    MEDIA_PROVENANCE_FIELDS,
    StandardTemplateSpecValidator,
)
from academic_ppt.templates import TemplateCatalog


ROOT = Path(__file__).resolve().parents[1]
SPEC_PATHS = {
    "T01": ROOT / "assets" / "template_specs" / "T01_green_research.semantic.json",
    "T03": ROOT / "assets" / "template_specs" / "T03_blue_defense.semantic.json",
}


def _load(template_id: str) -> dict:
    return json.loads(SPEC_PATHS[template_id].read_text(encoding="utf-8"))


def _page(specification: dict, page_id: str) -> dict:
    return next(page for page in specification["pages"] if page["page_id"] == page_id)


def test_t01_t03_standard_specs_record_powerpoint_review_without_claiming_release():
    validator = StandardTemplateSpecValidator()

    for template_id, path in SPEC_PATHS.items():
        specification = _load(template_id)
        result = validator.validate(specification)

        assert result.passed, result.errors
        assert specification["template"]["id"] == template_id
        assert specification["acceptance"]["semantic_compile_passed"]
        assert specification["acceptance"]["product_accepted"] is False
        assert specification["acceptance"]["powerpoint_visual_review"] == "passed"
        assert path.is_file()


def test_semantic_specs_encode_page_and_module_media_as_independent_contracts():
    t01 = _load("T01")
    t03 = _load("T03")

    assert _page(t01, "T01_P04")["media_layout"] == {
        "scope": "module",
        "kind": "one_per_module",
        "slot_count": 3,
        "native_reuse_supported": False,
        "missing_asset_behavior": "remove_or_reflow",
    }
    assert _page(t01, "T01_P07")["media_layout"]["slot_count"] == 6
    assert _page(t01, "T01_P07")["media_layout"]["scope"] == "page"
    assert _page(t03, "T03_P09")["media_layout"]["kind"] == "six_image"
    assert _page(t03, "T03_P09")["media_layout"]["slot_count"] == 6

    for specification in (t01, t03):
        contract = specification["media_contract"]
        assert {"one_image", "two_image", "three_image", "four_image", "six_image"} <= set(
            contract["page_level_layouts"]
        )
        assert contract["module_level_media"]
        for page in specification["pages"]:
            for module in page["semantic_modules"]:
                for slot in module["child_slots"]:
                    if slot["role"] == "image_or_chart":
                        assert tuple(slot["source_binding"]["required_fields"]) == MEDIA_PROVENANCE_FIELDS
                        assert slot["empty_behavior"] == "remove_or_reflow"


def test_semantic_specs_enforce_exclusive_regions_complete_removal_and_font_bounds():
    for specification in (_load("T01"), _load("T03")):
        policy = specification["typography_capacity_policy"]
        assert policy["cover_title"]["min_pt"] == 28
        assert policy["cover_title"]["max_pt"] == 32
        assert policy["page_title"]["min_pt"] == 20
        assert policy["module_heading"]["min_pt"] == 13
        assert policy["body"]["min_pt"] == 11
        assert policy["caption"]["min_pt"] == 8.5
        assert policy["normal_content_min_pt"] == 10.5
        assert policy["unbounded_shrink_forbidden"]
        for page in specification["pages"]:
            assert not page["ownership_audit"]["unowned_component_ids"]
            assert not page["ownership_audit"]["duplicate_component_ids"]
            for module in page["semantic_modules"]:
                render = module["render_contract"]
                assert render["mutually_exclusive"]
                assert render["masking_forbidden"]
                ownership = module["ownership_group"]
                assert ownership["shape_ids"] == ownership["complete_delete_shape_ids"]
                assert ownership["removal_policy"] == "remove_complete_ownership_group"


def test_semantic_specs_keep_content_modules_out_of_navigation_with_usable_geometry():
    for specification in (_load("T01"), _load("T03")):
        for page in specification["pages"]:
            modules = [
                module for module in page["semantic_modules"]
                if module["semantic_region"] == "content"
            ]
            if page["page_role"] not in {"cover", "ending", "section"}:
                for module in modules:
                    box = module["box"]
                    assert box["top"] >= 0.20, (page["page_id"], module["module_id"], box)
                    assert box["width"] >= 0.12, (page["page_id"], module["module_id"], box)
                    assert box["height"] >= 0.12, (page["page_id"], module["module_id"], box)
                    assert box["left"] >= 0
                    assert box["left"] + box["width"] <= 1.001
                    assert box["top"] + box["height"] <= 1.001
                for index, first in enumerate(modules):
                    for second in modules[index + 1:]:
                        left = max(first["box"]["left"], second["box"]["left"])
                        top = max(first["box"]["top"], second["box"]["top"])
                        right = min(
                            first["box"]["left"] + first["box"]["width"],
                            second["box"]["left"] + second["box"]["width"],
                        )
                        bottom = min(
                            first["box"]["top"] + first["box"]["height"],
                            second["box"]["top"] + second["box"]["height"],
                        )
                        assert right <= left or bottom <= top, (
                            page["page_id"], first["module_id"], second["module_id"]
                        )


def test_t01_t03_standard_pptx_objects_match_semantic_ownership():
    catalog = TemplateCatalog.load()
    by_id = {item["id"]: item for item in catalog.templates}

    for template_id in ("T01", "T03"):
        result = ObjectLevelQualityGate().inspect(
            catalog.root / by_id[template_id]["path"],
            SPEC_PATHS[template_id],
        )

        assert result.passed, result.errors
        assert result.inspected_slide_count == _load(template_id)["template"]["slide_count"]


def test_t01_t03_identity_signatures_are_not_color_only_duplicates():
    result = TemplateIdentityDifferenceGate().compare([_load("T01"), _load("T03")])

    assert result.passed, result.errors
    assert result.signatures["T01"] != result.signatures["T03"]


def test_object_gate_rejects_duplicate_empty_bad_source_small_font_and_sample_residue(tmp_path: Path):
    image_path = tmp_path / "figure.png"
    Image.new("RGB", (160, 90), "white").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    first.text = "Template sample"
    first.text_frame.paragraphs[0].font.size = Pt(8)
    second = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    second.text = "Template sample"
    second.text_frame.paragraphs[0].font.size = Pt(8)
    picture = slide.shapes.add_picture(str(image_path), Inches(6), Inches(1), Inches(3), Inches(2))
    pptx = tmp_path / "bad.pptx"
    presentation.save(pptx)
    sample_hash = hashlib.sha256("Template sample".encode("utf-8")).hexdigest()
    module_id = "T01_P01_M01"
    specification = {
        "pages": [{
            "page_id": "T01_P01",
            "semantic_modules": [{
                "module_id": module_id,
                "ownership_group": {
                    "shape_ids": [first.shape_id, second.shape_id, picture.shape_id],
                    "complete_delete_shape_ids": [first.shape_id, second.shape_id, picture.shape_id],
                },
                "render_contract": {"allowed_modes": ["native_reuse", "full_reconstruction"]},
                "child_slots": [
                    {
                        "slot_id": "TEXT",
                        "role": "explanation",
                        "typography": {"min_pt": 11, "max_pt": 13},
                        "source_sample_fingerprints": [{"shape_id": first.shape_id, "sha256": sample_hash}],
                    },
                    {"slot_id": "MEDIA_EMPTY", "role": "image_or_chart", "required": True},
                    {"slot_id": "MEDIA_BAD_SOURCE", "role": "image_or_chart", "required": False},
                ],
            }],
            "identity_regions": [],
        }],
    }
    manifest = {
        "pages": [{
            "slide_number": 1,
            "template_page_id": "T01_P01",
            "regions": [{
                "module_id": module_id,
                "render_mode": "native_reuse",
                "slot_bindings": [
                    {"slot_id": "TEXT", "shape_ids": [first.shape_id]},
                    {
                        "slot_id": "MEDIA_BAD_SOURCE",
                        "shape_ids": [picture.shape_id],
                        "source": {"source_id": "SRC", "source_path": "missing.pdf"},
                    },
                ],
            }],
        }],
    }

    result = ObjectLevelQualityGate().inspect(
        pptx,
        specification,
        render_manifest=manifest,
        asset_root=tmp_path,
    )

    assert not result.passed
    assert result.categories["duplicate_objects"]
    assert result.categories["empty_media_slots"]
    assert result.categories["media_provenance"]
    assert result.categories["font_bounds"]
    assert result.categories["template_residue"]


def test_object_gate_rejects_picture_that_does_not_match_declared_asset(tmp_path: Path):
    actual_asset = tmp_path / "actual.png"
    declared_asset = tmp_path / "declared.png"
    Image.new("RGB", (160, 90), "white").save(actual_asset)
    Image.new("RGB", (160, 90), "black").save(declared_asset)
    source_pdf = tmp_path / "paper.pdf"
    source_pdf.write_bytes(b"source")
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    picture = slide.shapes.add_picture(str(actual_asset), Inches(1), Inches(1), Inches(4), Inches(2.25))
    pptx = tmp_path / "wrong-picture.pptx"
    presentation.save(pptx)
    module_id = "T01_P01_M01"
    slot_id = "T01_P01_M01_MEDIA"
    specification = {
        "pages": [{
            "page_id": "T01_P01",
            "semantic_modules": [{
                "module_id": module_id,
                "ownership_group": {
                    "shape_ids": [picture.shape_id],
                    "complete_delete_shape_ids": [picture.shape_id],
                },
                "render_contract": {"allowed_modes": ["native_reuse", "full_reconstruction"]},
                "child_slots": [{"slot_id": slot_id, "role": "image_or_chart", "required": True}],
            }],
            "identity_regions": [],
        }],
    }
    manifest = {
        "pages": [{
            "slide_number": 1,
            "template_page_id": "T01_P01",
            "base_mode": "cloned_template",
            "source_components_present": True,
            "regions": [{
                "module_id": module_id,
                "render_mode": "native_reuse",
                "reflowed_or_removed_slots": [],
                "slot_bindings": [{
                    "slot_id": slot_id,
                    "shape_ids": [picture.shape_id],
                    "source": {
                        "source_id": "FIG-1",
                        "source_path": str(source_pdf),
                        "pdf_page": 3,
                        "semantic_use": "验证处理前后的形态差异",
                        "asset_path": str(declared_asset),
                    },
                }],
            }],
        }],
    }

    result = ObjectLevelQualityGate().inspect(
        pptx,
        specification,
        render_manifest=manifest,
        asset_root=tmp_path,
    )

    assert not result.passed
    assert any("does not match" in value for value in result.categories["media_provenance"])


def test_object_gate_rejects_nested_frames_and_incomplete_native_removal(tmp_path: Path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    outer = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(8), Inches(4.5)
    )
    inner = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.3), Inches(1.3), Inches(7.4), Inches(3.9)
    )
    pptx = tmp_path / "nested-frames.pptx"
    presentation.save(pptx)
    module_id = "T01_P01_M01"
    specification = {
        "template_identity": {"minimum_preserved_features": 4},
        "pages": [{
            "page_id": "T01_P01",
            "page_role": "content",
            "semantic_modules": [{
                "module_id": module_id,
                "ownership_group": {
                    "shape_ids": [outer.shape_id, inner.shape_id],
                    "complete_delete_shape_ids": [outer.shape_id, inner.shape_id],
                },
                "render_contract": {"allowed_modes": ["full_reconstruction"]},
                "child_slots": [],
            }],
            "identity_regions": [],
        }],
    }
    manifest = {
        "pages": [{
            "slide_number": 1,
            "template_page_id": "T01_P01",
            "base_mode": "blank_reconstruction",
            "source_components_present": False,
            "additional_identity_shape_ids": [],
            "identity_feature_bindings": [
                {"feature": f"feature-{index}", "shape_ids": [outer.shape_id]}
                for index in range(4)
            ],
            "regions": [{
                "module_id": module_id,
                "render_mode": "full_reconstruction",
                "owned_shape_ids": [outer.shape_id, inner.shape_id],
                "removed_native_shape_ids": [outer.shape_id],
                "slot_bindings": [],
                "reflowed_or_removed_slots": [],
            }],
        }],
    }

    result = ObjectLevelQualityGate().inspect(
        pptx,
        specification,
        render_manifest=manifest,
        asset_root=tmp_path,
    )

    assert not result.passed
    assert any("nested frame" in value for value in result.categories["duplicate_objects"])
    assert any("incomplete native removal" in value for value in result.categories["template_residue"])


def test_blank_reconstruction_keeps_native_and_rendered_shape_id_namespaces_separate(tmp_path: Path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    identity_shapes = []
    for index in range(7):
        shape = slide.shapes.add_textbox(
            Inches(0.1 + index * 0.2), Inches(0.1), Inches(0.15), Inches(0.15)
        )
        shape.text = f"N{index + 1}"
        identity_shapes.append(shape)
    second = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(0.6))
    second.text = "Rendered module two"
    first = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(3), Inches(0.6))
    first.text = "Rendered module one"
    assert first.shape_id == 10
    pptx = tmp_path / "blank-id-collision.pptx"
    presentation.save(pptx)

    def module(module_id: str, native_shape_id: int, rendered_shape_id: int) -> dict:
        return {
            "module_id": module_id,
            "ownership_group": {
                "shape_ids": [native_shape_id],
                "complete_delete_shape_ids": [native_shape_id],
            },
            "render_contract": {"allowed_modes": ["full_reconstruction"]},
            "child_slots": [{
                "slot_id": f"{module_id}_TEXT",
                "role": "explanation",
                "required": True,
                "typography": {"min_pt": 11, "max_pt": 13},
            }],
            "rendered_shape_id": rendered_shape_id,
        }

    first_module = module("T02_P01_M01", 101, first.shape_id)
    second_module = module("T02_P01_M02", 10, second.shape_id)
    specification = {
        "template_identity": {"minimum_preserved_features": 4},
        "pages": [{
            "page_id": "T02_P01",
            "page_role": "content",
            "semantic_modules": [first_module, second_module],
            "identity_regions": [],
        }],
    }
    manifest = {
        "pages": [{
            "slide_number": 1,
            "template_page_id": "T02_P01",
            "base_mode": "blank_reconstruction",
            "source_components_present": False,
            "additional_identity_shape_ids": [shape.shape_id for shape in identity_shapes],
            "identity_feature_bindings": [
                {"feature": f"feature-{index}", "shape_ids": [identity_shapes[index].shape_id]}
                for index in range(4)
            ],
            "regions": [{
                "module_id": item["module_id"],
                "render_mode": "full_reconstruction",
                "owned_shape_ids": [item["rendered_shape_id"]],
                "removed_native_shape_ids": item["ownership_group"]["shape_ids"],
                "slot_bindings": [{
                    "slot_id": f"{item['module_id']}_TEXT",
                    "shape_ids": [item["rendered_shape_id"]],
                }],
                "reflowed_or_removed_slots": [],
            } for item in (first_module, second_module)],
        }],
    }

    result = ObjectLevelQualityGate().inspect(
        pptx,
        specification,
        render_manifest=manifest,
        asset_root=tmp_path,
    )

    assert result.passed, result.errors


def test_catalog_discloses_semantic_compile_status_without_claiming_all_templates_complete():
    catalog = TemplateCatalog.load()
    t01 = catalog.select("组会-文献精读", "T01")
    t02 = catalog.select("学术会议报告", "T02")

    assert Path(t01.semantic_spec_path).is_file()
    assert t01.standardization_status == "semantic_compiled_powerpoint_review_passed"
    assert Path(t02.semantic_spec_path).is_file()
    assert t02.standardization_status == "semantic_compilation_in_progress"


def test_validate_pptx_cli_runs_object_qa_when_semantic_spec_is_supplied(tmp_path: Path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    shape.text = "Bound content"
    pptx = tmp_path / "valid.pptx"
    presentation.save(pptx)
    component_id = "S01_C2"
    specification = {
        "pages": [{
            "page_id": "T01_P01",
            "semantic_modules": [{
                "module_id": "T01_P01_M01",
                "ownership_group": {
                    "component_ids": [component_id],
                    "shape_ids": [shape.shape_id],
                    "complete_delete_shape_ids": [shape.shape_id],
                },
                "child_slots": [],
            }],
            "identity_regions": [],
        }],
    }
    spec_path = tmp_path / "semantic.json"
    spec_path.write_text(json.dumps(specification), encoding="utf-8")
    report_path = tmp_path / "report.json"

    subprocess.run(
        [
            sys.executable,
            str(ROOT / "scripts" / "validate_pptx.py"),
            str(pptx),
            "--render-check", "off",
            "--semantic-spec", str(spec_path),
            "--output", str(report_path),
        ],
        cwd=ROOT,
        check=True,
        capture_output=True,
        text=True,
        encoding="utf-8",
    )
    report = json.loads(report_path.read_text(encoding="utf-8"))

    assert report["object_qa"]["passed"]
    assert report["failed_error"] == 0


def test_unbound_standard_template_defers_delivery_text_geometry_checks(tmp_path: Path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    first = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.2))
    first.text = "Native layered heading with intentionally compact geometry"
    second = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.2))
    second.text = "Native supporting label"
    pptx = tmp_path / "unbound-standard-template.pptx"
    presentation.save(pptx)
    specification = {
        "pages": [{
            "page_id": "T01_P01",
            "semantic_modules": [{
                "module_id": "T01_P01_M01",
                "ownership_group": {
                    "component_ids": ["S01_C2", "S01_C3"],
                    "shape_ids": [first.shape_id, second.shape_id],
                    "complete_delete_shape_ids": [first.shape_id, second.shape_id],
                },
                "child_slots": [],
            }],
            "identity_regions": [],
        }],
    }
    spec_path = tmp_path / "semantic.json"
    spec_path.write_text(json.dumps(specification), encoding="utf-8")
    report_path = tmp_path / "report.json"

    subprocess.run(
        [
            sys.executable,
            str(ROOT / "scripts" / "validate_pptx.py"),
            str(pptx),
            "--render-check", "off",
            "--semantic-spec", str(spec_path),
            "--output", str(report_path),
        ],
        cwd=ROOT,
        check=True,
        capture_output=True,
        text=True,
        encoding="utf-8",
    )
    report = json.loads(report_path.read_text(encoding="utf-8"))

    assert report["failed_error"] == 0
    assert any("object binding manifest" in value for value in report["observations"])
    assert "No severe estimated text overflow" not in {
        check["name"] for check in report["checks"]
    }


def test_native_template_validator_rejects_placeholder_masks_even_when_marked_temporary(tmp_path: Path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    shape.text = "Editable"
    template = tmp_path / "template.pptx"
    presentation.save(template)
    plan = {
        "template_mode": "template_native",
        "pages": [{
            "page_id": "P001",
            "render_mode": "template_native",
            "source_slide_index": 1,
            "text_bindings": [{"shape_id": shape.shape_id, "content": "Bound"}],
            "placeholder_masks": [{
                "container_shape_id": shape.shape_id,
                "temporary_approved": True,
                "reason": "legacy workaround",
            }],
        }],
    }
    plan_path = tmp_path / "plan.json"
    plan_path.write_text(json.dumps(plan), encoding="utf-8")

    completed = subprocess.run(
        [
            sys.executable,
            str(ROOT / "scripts" / "validate_template_native.py"),
            str(plan_path),
            str(template),
        ],
        cwd=ROOT,
        check=False,
        capture_output=True,
        text=True,
        encoding="utf-8",
    )

    assert completed.returncode == 1
    assert "placeholder masks are forbidden" in completed.stdout
