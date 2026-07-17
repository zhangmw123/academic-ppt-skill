import base64
import re
import sys
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from PIL import Image
from pptx import Presentation


SCRIPTS = Path(__file__).resolve().parents[1] / "scripts"
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

from export_preview import export_preview, export_with_powerpoint, list_slide_images


class PowerPointExportTests(unittest.TestCase):
    def test_stages_unicode_source_path_before_invoking_powerpoint_com(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "材料.pptx"
            source.write_bytes(b"test")
            output = root / "preview"
            output.mkdir()
            commands = []

            def fake_run(command, **_kwargs):
                commands.append(command)
                power_shell = base64.b64decode(command[-1]).decode("utf-16-le")
                match = re.search(r"Export\('([^']+)'", power_shell)
                self.assertIsNotNone(match)
                Path(match.group(1)).joinpath("Slide1.PNG").write_bytes(b"png")

            with patch("export_preview.subprocess.run", side_effect=fake_run):
                export_with_powerpoint(source, output, 1600, 900)

            command_text = base64.b64decode(commands[0][-1]).decode("utf-16-le")
            self.assertIn("deck.pptx", command_text)
            self.assertNotIn("材料", command_text)
            self.assertTrue((output / "Slide1.PNG").is_file())

    def test_lists_localized_slide_images_without_assuming_english_names(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            output = Path(temp_dir)
            for filename in ("幻灯片1.PNG", "幻灯片2.PNG", "Slide10.png", "contact-sheet.png"):
                (output / filename).write_bytes(b"png")

            self.assertEqual(
                [path.name for path in list_slide_images(output)],
                ["Slide10.png", "幻灯片1.PNG", "幻灯片2.PNG"],
            )

    def test_rejects_partial_powerpoint_exports(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.slides.add_slide(presentation.slide_layouts[6])
            source = root / "two-slides.pptx"
            presentation.save(source)
            output = root / "preview"

            def fake_export(_pptx, destination, _width, _height):
                Image.new("RGB", (1600, 900), "white").save(destination / "幻灯片1.PNG")

            with (
                patch("export_preview.find_powerpoint", return_value=True),
                patch("export_preview.export_with_powerpoint", side_effect=fake_export),
            ):
                with self.assertRaisesRegex(RuntimeError, "expected 2 slide images; rendered 1"):
                    export_preview(source, output, "powerpoint", 1600, 900)


if __name__ == "__main__":
    unittest.main()
