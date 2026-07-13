import tempfile
import unittest
from pathlib import Path

from openpyxl import Workbook
from docx import Document
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from academic_ppt.ingest import SourceIngestor


class SourceIngestorTests(unittest.TestCase):
    def test_text_source_is_normalized_with_role_hash_and_line_locators(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "实验说明.md"
            path.write_text("# 实验结果\n模型 F1 为 92.9%。\n边界：仅一个数据集。\n", encoding="utf-8")

            source = SourceIngestor().ingest(path, role="primary_evidence")

            self.assertEqual(source.format, "markdown")
            self.assertEqual(source.role, "primary_evidence")
            self.assertEqual(source.blocks[1].text, "模型 F1 为 92.9%。")
            self.assertEqual(source.blocks[1].locator, {"line": 2})
            self.assertTrue(source.sha256)
            self.assertTrue(source.source_id.startswith("SRC_"))

    def test_xlsx_source_preserves_sheet_range_values_and_formulas(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "最新实验.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "Metrics"
            sheet.append(["模型", "F1", "增益"])
            sheet.append(["Baseline", 89.7, None])
            sheet.append(["Proposed", 92.9, "=B3-B2"])
            workbook.save(path)

            source = SourceIngestor().ingest(path)

            table = source.blocks[0]
            self.assertEqual(source.role, "incremental_evidence")
            self.assertEqual(table.kind, "table")
            self.assertEqual(table.locator, {"sheet": "Metrics", "range": "A1:C3"})
            self.assertEqual(table.data["rows"][2][1], 92.9)
            self.assertEqual(table.data["formulas"], {"C3": "=B3-B2"})

    def test_csv_source_preserves_rows_and_row_range(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "实验增量.csv"
            path.write_text("模型,F1\nBaseline,89.7\nProposed,92.9\n", encoding="utf-8")

            source = SourceIngestor().ingest(path)

            self.assertEqual(source.format, "csv")
            self.assertEqual(source.role, "incremental_evidence")
            self.assertEqual(source.blocks[0].kind, "table")
            self.assertEqual(source.blocks[0].locator, {"row_start": 1, "row_end": 3})
            self.assertEqual(source.blocks[0].data["rows"][2], ["Proposed", "92.9"])

    def test_docx_source_preserves_heading_paragraph_and_table_locations(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "论文.docx"
            document = Document()
            document.add_heading("实验结果", level=1)
            document.add_paragraph("模型在测试集上取得最优结果。")
            table = document.add_table(rows=2, cols=2)
            table.cell(0, 0).text = "模型"
            table.cell(0, 1).text = "F1"
            table.cell(1, 0).text = "Proposed"
            table.cell(1, 1).text = "92.9"
            document.save(path)

            source = SourceIngestor().ingest(path)

            self.assertEqual(source.format, "docx")
            self.assertEqual(source.blocks[0].kind, "heading")
            self.assertEqual(source.blocks[1].locator, {"paragraph": 2, "heading_path": ["实验结果"]})
            self.assertEqual(source.blocks[2].kind, "table")
            self.assertEqual(source.blocks[2].locator, {"table": 1, "heading_path": ["实验结果"]})

    def test_pptx_source_preserves_shapes_tables_pictures_and_notes(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            image_path = root / "figure.png"
            Image.new("RGB", (320, 180), "white").save(image_path)
            path = root / "历史汇报.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            text_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            text_shape.text = "实验结果提升 3.2 个百分点"
            table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(2), Inches(4), Inches(1.5))
            table_shape.table.cell(0, 0).text = "模型"
            table_shape.table.cell(0, 1).text = "F1"
            table_shape.table.cell(1, 0).text = "Proposed"
            table_shape.table.cell(1, 1).text = "92.9"
            slide.shapes.add_picture(str(image_path), Inches(6), Inches(1), Inches(3), Inches(2))
            slide.notes_slide.notes_text_frame.text = "先讲完整模型，再解释消融。"
            presentation.save(path)

            source = SourceIngestor().ingest(path, role="historical_expression")

            kinds = [block.kind for block in source.blocks]
            self.assertIn("text", kinds)
            self.assertIn("table", kinds)
            self.assertIn("picture", kinds)
            self.assertIn("speaker_notes", kinds)
            text_block = next(block for block in source.blocks if block.kind == "text")
            self.assertEqual(text_block.locator["slide"], 1)
            self.assertIn("shape_id", text_block.locator)

    def test_pdf_source_preserves_page_numbered_text_and_embedded_image_metadata(self):
        import fitz

        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            image_path = root / "result.png"
            Image.new("RGB", (320, 180), "white").save(image_path)
            path = root / "论文结果.pdf"
            document = fitz.open()
            first_page = document.new_page()
            first_page.insert_text((72, 72), "Model F1 reached 92.9%.")
            first_page.insert_image(fitz.Rect(72, 100, 232, 190), filename=image_path)
            second_page = document.new_page()
            second_page.insert_text((72, 72), "Boundary: one evaluation dataset.")
            document.save(path)
            document.close()

            source = SourceIngestor().ingest(path)

            text_blocks = [block for block in source.blocks if block.kind == "text"]
            picture_block = next(block for block in source.blocks if block.kind == "picture")
            self.assertEqual(source.format, "pdf")
            self.assertEqual([block.locator["page"] for block in text_blocks], [1, 2])
            self.assertEqual(text_blocks[0].text, "Model F1 reached 92.9%.")
            self.assertEqual(picture_block.locator["page"], 1)
            self.assertEqual(picture_block.data["image_size"], [320, 180])
            self.assertEqual(picture_block.data["image_format"], "png")
            self.assertTrue(picture_block.data["image_sha256"])

    def test_image_source_preserves_safe_visual_metadata_without_exif(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "观察图.png"
            Image.new("RGB", (640, 480), "white").save(path)

            source = SourceIngestor().ingest(path)

            self.assertEqual(source.format, "image")
            self.assertEqual(source.role, "visual_asset")
            self.assertEqual(source.metadata["image_size"], [640, 480])
            self.assertEqual(source.metadata["image_mode"], "RGB")
            self.assertEqual(source.metadata["image_format"], "png")
            self.assertNotIn("exif", source.metadata)
            self.assertEqual(source.blocks[0].kind, "picture")
            self.assertTrue(source.blocks[0].data["image_sha256"])


if __name__ == "__main__":
    unittest.main()
