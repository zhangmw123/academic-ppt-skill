import sys
from pathlib import Path

import pytest
from PIL import Image, ImageDraw
from pptx import Presentation


SCRIPTS = Path(__file__).resolve().parents[1] / "scripts"
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

from check_powerpoint_runtime import build_minimal_image_control, inspect_images


def test_builds_a_colored_minimal_image_control(tmp_path: Path):
    output = build_minimal_image_control(tmp_path / "minimal.pptx")

    presentation = Presentation(output)
    assert len(presentation.slides) == 1
    assert len(presentation.slides[0].shapes) == 2
    assert (tmp_path / "minimal.png").is_file()


def test_inspects_localized_color_control_images(tmp_path: Path):
    image = Image.new("RGB", (1600, 900), "white")
    draw = ImageDraw.Draw(image)
    for index, color in enumerate(("red", "green", "blue", "white")):
        draw.rectangle((index * 400, 0, (index + 1) * 400 - 1, 899), fill=color)
    image.save(tmp_path / "幻灯片1.PNG")

    details = inspect_images(
        tmp_path,
        1,
        expected_size=(1600, 900),
        require_color_control=True,
    )

    assert details[0]["sample_color_count"] >= 4


def test_rejects_incomplete_runtime_control_exports(tmp_path: Path):
    Image.new("RGB", (1600, 900), "white").save(tmp_path / "Slide1.PNG")

    with pytest.raises(RuntimeError, match="expected 2 slide images; rendered 1"):
        inspect_images(
            tmp_path,
            2,
            expected_size=(1600, 900),
            require_color_control=False,
        )
