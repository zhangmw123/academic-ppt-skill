import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

SCRIPTS = Path(__file__).resolve().parents[1] / "scripts"
if str(SCRIPTS) not in sys.path:
    sys.path.insert(0, str(SCRIPTS))

from extract_template_grammar import extract
from academic_ppt.template_semantics import StandardTemplateCompiler


def test_navigation_picture_backings_are_not_classified_as_logos(tmp_path: Path):
    backing = tmp_path / "active-tab.png"
    Image.new("RGB", (320, 80), "#176B4D").save(backing)
    inactive = tmp_path / "inactive-tab.png"
    Image.new("RGB", (320, 80), "#D9E8E2").save(inactive)
    presentation = Presentation()
    for _ in range(3):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(str(backing), Inches(0.5), Inches(0.25), Inches(2.0), Inches(0.5))
        slide.shapes.add_picture(str(inactive), Inches(2.7), Inches(0.25), Inches(2.0), Inches(0.5))
        for index, label in enumerate(("Section A", "Section B", "Section C")):
            textbox = slide.shapes.add_textbox(
                Inches(0.5 + index * 2.2),
                Inches(0.34),
                Inches(2.0),
                Inches(0.24),
            )
            textbox.text = label
    source = tmp_path / "navigation-template.pptx"
    presentation.save(source)
    output = tmp_path / "grammar.json"

    extract(source, output, tmp_path / "assets")

    grammar = json.loads(output.read_text(encoding="utf-8"))
    first = grammar["archetypes"][0]
    active_asset = next(asset for asset in first["decorative_assets"] if asset["shape_id"] == 2)
    active_slot = next(slot for slot in first["picture_slots"] if slot["shape_id"] == 2)
    inactive_asset = next(asset for asset in first["decorative_assets"] if asset["shape_id"] == 3)
    assert active_asset["role"] == "navigation"
    assert active_slot["role"] == "navigation"
    assert not active_slot["logo_candidate"]
    assert inactive_asset["role"] == "navigation"

    identity_shape_ids = StandardTemplateCompiler._identity_shape_ids(
        first,
        first["text_slots"],
    )
    assert {2, 3} <= identity_shape_ids
    identity_components = [
        component for component in first["components"]
        if identity_shape_ids.intersection(component["shape_ids"])
    ]
    navigation_region = next(
        region for region in StandardTemplateCompiler._identity_regions(identity_components, first)
        if region["role"] == "navigation"
    )
    assert {2, 3} <= set(navigation_region["ownership_group"]["shape_ids"])
