import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation


class AcademicPptCliTests(unittest.TestCase):
    def test_cli_persists_an_edit_baseline_and_plans_a_targeted_revision(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir) / "presentation"
            root.mkdir()
            generated = root / "generated.pptx"
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(generated)

            adopted = subprocess.run(
                [sys.executable, "-m", "academic_ppt", "adopt-baseline", str(root), str(generated), "--page-ids", "P001"],
                check=True, capture_output=True, text=True, encoding="utf-8",
            )
            planned = subprocess.run(
                [sys.executable, "-m", "academic_ppt", "plan-revision", str(root), str(generated), "--targets", "P001"],
                check=True, capture_output=True, text=True, encoding="utf-8",
            )

            self.assertTrue(Path(json.loads(adopted.stdout)["baseline"]).is_file())
            self.assertEqual(json.loads(planned.stdout)["target_page_ids"], ["P001"])
            self.assertTrue((root / "audit" / "revision_plan.json").is_file())

    def test_cli_prepares_a_confirmable_task_summary_from_sources(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir) / "presentation"
            source_path = Path(temp_dir) / "model.md"
            source_path.write_text("Model algorithm training dataset baseline ablation.\nModel F1 reached 92.9%.", encoding="utf-8")

            result = subprocess.run(
                [
                    sys.executable,
                    "-m",
                    "academic_ppt",
                    "prepare",
                    str(root),
                    str(source_path),
                    "--scene",
                    "毕业答辩",
                ],
                check=True,
                capture_output=True,
                text=True,
                encoding="utf-8",
            )

            report = json.loads(result.stdout)
            self.assertEqual(report["project"], str(root.resolve()))
            self.assertEqual(report["phase"], "awaiting_confirmation")
            self.assertEqual(report["task_summary"]["source_count"], 1)
            self.assertTrue((root / "audit" / "task_summary.json").is_file())

    def test_cli_initializes_project_and_reports_status(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir) / "demo"
            init_result = subprocess.run(
                [
                    sys.executable,
                    "-m",
                    "academic_ppt",
                    "init",
                    str(root),
                    "--scene",
                    "毕业答辩",
                    "--rigor",
                    "standard",
                ],
                check=True,
                capture_output=True,
                text=True,
                encoding="utf-8",
            )
            status_result = subprocess.run(
                [sys.executable, "-m", "academic_ppt", "status", str(root)],
                check=True,
                capture_output=True,
                text=True,
                encoding="utf-8",
            )

            initialized = json.loads(init_result.stdout)
            status = json.loads(status_result.stdout)
            self.assertEqual(initialized["scene"], "毕业答辩")
            self.assertEqual(status["phases"]["brief"], "draft")
            self.assertEqual(status["phases"]["evidence"], "not_started")

    def test_cli_analyzes_sources_into_profiles_evidence_and_conflicts(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source_path = root / "模型结果.md"
            source_path.write_text(
                "模型 算法 训练 数据集 基线 消融 泛化。\n模型 F1 为 92.9%。\n",
                encoding="utf-8",
            )

            result = subprocess.run(
                [sys.executable, "-m", "academic_ppt", "analyze", str(source_path)],
                check=True,
                capture_output=True,
                text=True,
                encoding="utf-8",
            )

            report = json.loads(result.stdout)
            self.assertEqual(report["source_count"], 1)
            self.assertEqual(report["method_profiles"]["primary"], "computational_modeling")
            self.assertEqual(report["evidence"][1]["evidence_type"], "result")
            self.assertEqual(report["conflicts"], [])


if __name__ == "__main__":
    unittest.main()
