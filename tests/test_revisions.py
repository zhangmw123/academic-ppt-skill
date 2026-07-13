import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from academic_ppt.revisions import AuthoritativeEditBaseline


class AuthoritativeEditBaselineTests(unittest.TestCase):
    def test_detects_manual_changes_and_protects_untargeted_pages_from_rebuild(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            generated = root / "generated.pptx"
            presentation = Presentation()
            for text in ("Original first page", "Original second page"):
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1)).text = text
            presentation.save(generated)
            baseline = AuthoritativeEditBaseline.adopt(generated, ("P001", "P002"))

            edited = root / "edited.pptx"
            presentation = Presentation(generated)
            presentation.slides[0].shapes[0].text = "User-edited first page"
            presentation.slides[1].shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1)).text = "Manual note"
            presentation.save(edited)

            revision = baseline.plan_revision(edited, target_page_ids=("P001",))

            self.assertEqual(revision.target_page_ids, ("P001",))
            self.assertEqual(revision.protected_page_ids, ("P002",))
            self.assertEqual(revision.changed_page_ids, ("P001", "P002"))
            self.assertEqual(revision.manual_shape_ids["P002"], (3,))
            self.assertEqual(revision.replacement_disclosures, ("P001: user edits may be replaced by the requested rebuild",))

    def test_rejects_a_revision_target_that_does_not_exist_in_the_baseline(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            path = Path(temp_dir) / "generated.pptx"
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(path)
            baseline = AuthoritativeEditBaseline.adopt(path, ("P001",))

            with self.assertRaisesRegex(ValueError, "unknown revision target"):
                baseline.plan_revision(path, target_page_ids=("P404",))

    def test_persists_and_reloads_the_authoritative_baseline(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            path = root / "generated.pptx"
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(path)

            baseline_path = AuthoritativeEditBaseline.adopt(path, ("P001",)).write(root / "audit")
            restored = AuthoritativeEditBaseline.load(baseline_path)

            self.assertEqual(restored.page_ids, ("P001",))
            self.assertEqual(restored.source_path, path.resolve())


if __name__ == "__main__":
    unittest.main()
