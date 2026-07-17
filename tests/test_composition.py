from dataclasses import replace
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

from academic_ppt.autobuild import CompleteContentDraft, CompleteContentPackage
from academic_ppt.composition import CompositionQualityGate, DynamicCompositionCompiler
from academic_ppt.layout import ScientificPageContract
from academic_ppt.manifest import SlideManifestBuilder
from academic_ppt.object_qa import ObjectLevelQualityGate
from academic_ppt.planning import PagePlan, PlannedPage
from academic_ppt.scenes import ScenePlanContract
from scripts.render_dynamic import render as render_dynamic


def _page(page_id: str, title: str, strategy: str, text_count: int) -> PlannedPage:
    return PlannedPage(
        page_id=page_id,
        section="研究内容",
        title=title,
        question_answered="该页证据说明什么？",
        claim_id=f"CLM_{page_id}",
        claim_text="多源证据共同支持当前判断，并给出可以复核的实现依据。",
        evidence_ids=(f"EVD_{page_id}_1", f"EVD_{page_id}_2"),
        interpretation="该结果说明方法链路有效，但仍需结合边界条件解释。",
        next_link="下一页进入进一步验证。",
        time_seconds=60,
        visual_strategy=strategy,
        contract=ScientificPageContract(page_id, f"CLM_{page_id}", {"text": text_count}),
    )


def _package(pages: tuple[PlannedPage, ...], image_path: Path) -> CompleteContentPackage:
    drafts = tuple(
        CompleteContentDraft(
            page_id=page.page_id,
            section=page.section,
            title=page.title,
            question_answered=page.question_answered,
            claim_text=page.claim_text,
            evidence_ids=page.evidence_ids,
            interpretation=page.interpretation,
            next_link=page.next_link,
            time_seconds=page.time_seconds,
            visual_strategy=page.visual_strategy,
            component_requirements=dict(page.contract.component_requirements),
        )
        for page in pages
    )
    text_content = {
        "P001": [pages[0].title, "完整论文汇报"],
        "P002": [pages[1].title, "原文图展示关键结构与比较对象", "读图时需要同时关注机制、对照和适用边界"],
        "P003": [pages[2].title, "输入数据清洗与结构化", "模型完成语义表示与关系推断", "输出结果接受定量和案例验证", "解释方法为何有效", "转入结果验证"],
        "P004": [pages[3].title, "核心结论与下一步行动"],
    }
    return CompleteContentPackage(
        sections=("研究内容",),
        drafts=drafts,
        scene_contract=ScenePlanContract("complete", "published"),
        text_content=text_content,
        image_content={"P002": [str(image_path)]},
        figure_manifests=(),
    )


def test_dynamic_composition_routes_figures_and_missing_figures_to_editable_visuals(tmp_path: Path):
    pages = (
        _page("P001", "论文题目", "text_only", 2),
        _page("P002", "结果图同时支持机制与比较结论", "source_figure", 3),
        _page("P003", "系统流程形成完整验证链路", "native_diagram", 6),
        _page("P004", "结论与行动", "text_only", 2),
    )
    plan = PagePlan("组会-文献精读", ("研究内容",), pages, deck_scope="complete")
    image = tmp_path / "figure.png"
    image.write_bytes(b"placeholder")
    grammar = tmp_path / "grammar.json"
    grammar.write_text(json.dumps({"archetypes": [
        {"slide_index": 1, "role": "cover", "layout_signature": "cover"},
        {"slide_index": 2, "role": "gallery", "layout_signature": "text_figure_right"},
        {"slide_index": 3, "role": "process", "layout_signature": "four_columns"},
        {"slide_index": 4, "role": "ending", "layout_signature": "ending"},
    ]}), encoding="utf-8")

    payload = DynamicCompositionCompiler().compile(
        plan,
        _package(pages, image),
        visual_system_path=tmp_path / "visual.json",
        template_grammar_path=grammar,
        asset_base_dir=tmp_path,
        confirmed=True,
    )

    assert payload["composition_mode"] == "template_hybrid_editable"
    assert [page["layout"] for page in payload["pages"]] == ["cover", "text_figure", "architecture", "ending"]
    assert len(payload["pages"][1]["bullets"]) >= 3
    assert all(set(item) == {"title", "body"} for item in payload["pages"][1]["bullets"])
    assert sum(len(column["nodes"]) for column in payload["pages"][2]["architecture"]["columns"]) >= 4
    assert all(page.get("use_template_scaffold") == "identity" for page in payload["pages"][1:-1])


def test_composition_uses_template_panel_geometry_and_never_leaks_internal_labels(tmp_path: Path):
    pages = (
        _page("P001", "Title", "text_only", 2),
        _page("P002", "Four findings", "text_only", 4),
        _page("P003", "Method details", "text_only", 4),
        _page("P004", "End", "text_only", 2),
    )
    package = _package(pages, tmp_path / "unused.png")
    package = replace(package, text_content={
        "P001": ["Title", "Defense"],
        "P002": [
            "Four findings",
            "Dataset scale: 6,200 annotated sentences and 20,688 entities",
            "Model result: precision reaches 91.58%",
            "Query result: average accuracy reaches 91.14%",
        ],
        "P003": ["Method details", "Input: annotated text", "Model: contextual encoding", "Output: entity triples"],
        "P004": ["End", "Thanks"],
    }, image_content={})
    grammar = tmp_path / "grammar.json"
    grammar.write_text(json.dumps({
        "identity": {"panel_mode": "image_scaffold"},
        "archetypes": [
            {"slide_index": 1, "role": "cover", "layout_signature": "cover"},
            {
                "slide_index": 6, "role": "gallery", "layout_signature": "three_columns",
                "picture_slots": [
                    {"role": "content_image", "area": 0.13, "box": {"left": x, "top": y, "width": 0.4, "height": 0.3}}
                    for x, y in ((0.06, 0.24), (0.52, 0.24), (0.06, 0.62), (0.52, 0.62))
                ],
            },
            {"slide_index": 3, "role": "gallery", "layout_signature": "three_columns"},
            {"slide_index": 11, "role": "ending", "layout_signature": "ending"},
        ],
    }), encoding="utf-8")
    plan = PagePlan("毕业答辩", ("Research",), pages, deck_scope="complete")

    payload = DynamicCompositionCompiler().compile(
        plan, package, visual_system_path=tmp_path / "visual.json",
        template_grammar_path=grammar, asset_base_dir=tmp_path, confirmed=True,
    )

    page = payload["pages"][1]
    assert page["scaffold_slide"] == 6
    assert page["use_template_scaffold"] == "structure"
    assert all(item["title"] not in {"核心判断", "证据依据", "机制解释", "边界与行动"} for item in page["points"])
    assert all(item["body"] for item in page["points"])
    assert any("91.14" in item["body"] for item in page["points"])


def test_composition_gate_rejects_sparse_placeholder_pages():
    result = CompositionQualityGate().inspect({
        "pages": [{
            "page_id": "P002",
            "layout": "text_figure",
            "title": "结果",
            "lead": "结果",
            "bullets": ["一句话"],
            "image": "figure.png",
        }],
    })

    assert not result.passed
    assert any("visible modules" in error for error in result.errors)
    assert any("at least three" in error for error in result.errors)


def test_semantic_selector_prefers_exact_module_count_over_reflowable_media_scope():
    root = Path(__file__).resolve().parents[1]
    specification = json.loads(
        (root / "assets" / "template_specs" / "T03_blue_defense.semantic.json").read_text(
            encoding="utf-8"
        )
    )
    reference = DynamicCompositionCompiler._semantic_template_reference(
        {
            "layout": "points",
            "points": [
                {"title": f"模块 {index}", "body": "来源支持的解释"}
                for index in range(1, 5)
            ],
        },
        specification,
    )

    assert reference["semantic_spec_page_id"] == "T03_P06"
    assert reference["content_module_count"] == 4


def test_slide_manifest_requires_editable_content_and_template_identity(tmp_path: Path):
    image = tmp_path / "figure.png"
    Image.new("RGB", (320, 180), "white").save(image)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    for index, text in enumerate(("标题", "证据一", "证据二", "解释", "结论")):
        slide.shapes.add_textbox(Inches(0.5), Inches(0.5 + index * 0.5), Inches(4), Inches(0.4)).text = text
    for index in range(3):
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5 + index), Inches(4.5), Inches(0.8), Inches(0.5))
    slide.shapes.add_picture(str(image), Inches(6), Inches(1), Inches(5), Inches(3))
    pptx = tmp_path / "deck.pptx"
    presentation.save(pptx)
    visual = tmp_path / "visual.json"
    visual.write_text(json.dumps({
        "fonts": {"title": "Microsoft YaHei", "body": "Microsoft YaHei", "latin": "Arial"},
        "typography": {"page_title": 24, "body": 14},
        "colors": {"primary": "#123456"},
    }), encoding="utf-8")
    grammar = tmp_path / "grammar.json"
    grammar.write_text(json.dumps({"identity": {"navigation": "top"}}), encoding="utf-8")
    plan = {
        "composition_mode": "template_hybrid_editable",
        "pages": [{
            "page_id": "P001",
            "layout": "text_figure",
            "lead": "证据、读图与边界",
            "bullets": [
                "第一条证据说明方法能够覆盖核心问题并形成可复核结果",
                "第二条证据提供对照关系并说明改进来自关键模块",
                "解释部分同时保留适用边界和后续验证要求",
            ],
            "page_conclusion": "综合证据支持当前结论，但仍需在外部数据上继续验证。",
            "image": str(image),
            "template_reference": {"source_slide_index": 3, "layout_signature": "text_figure_right"},
        }],
    }

    manifest = SlideManifestBuilder().build(
        pptx_path=pptx,
        dynamic_plan=plan,
        visual_system_path=visual,
        template_grammar_path=grammar,
        template={"id": "T01"},
    )

    assert manifest["passed"]
    assert manifest["slides"][0]["editable_information_layer_pass"]
    assert manifest["slides"][0]["template_identity_bound"]


def test_dynamic_composition_and_renderer_support_page_gallery_and_module_media(tmp_path: Path):
    pages = (
        _page("P001", "研究主题", "text_only", 2),
        replace(
            _page("P002", "四组图像共同证明处理效果", "source_figure", 4),
            media_scope="page",
            media_layout="four_image",
        ),
        replace(
            _page("P003", "四个实验模块分别提供图像证据", "module_media", 4),
            media_scope="module",
            media_layout="one_per_module",
        ),
        _page("P004", "结论与行动", "text_only", 2),
    )
    images = []
    for index in range(8):
        path = tmp_path / f"evidence_{index + 1}.png"
        Image.new("RGB", (320, 180), (240 - index * 8, 245, 250)).save(path)
        images.append(str(path))
    package = _package(pages, Path(images[0]))
    package = replace(
        package,
        text_content={
            "P001": ["研究主题", "完整证据汇报"],
            "P002": [
                "四组图像共同证明处理效果",
                "对照组显示原始结构，处理组显示稳定改善",
                "四个视野来自同一实验批次并采用一致尺度",
                "跨面板比较支持结论，同时保留样本边界",
            ],
            "P003": [
                "四个实验模块分别提供图像证据",
                "细胞形态：处理后边界更清晰且异常区域减少",
                "表达结果：目标蛋白变化与表型方向一致",
                "结构验证：结合位点预测支持所提出的机制",
                "剂量响应：效应随处理浓度增加并趋于稳定",
            ],
            "P004": ["结论与行动", "完成外部数据验证"],
        },
        image_content={"P002": images[:4], "P003": images[4:8]},
    )
    grammar_path = tmp_path / "grammar.json"
    grammar_path.write_text(json.dumps({
        "identity": {"panel_mode": "native_shapes", "navigation": "top"},
        "geometry": {},
        "archetypes": [
            {"slide_index": 1, "role": "cover", "layout_signature": "cover"},
            {"slide_index": 4, "role": "gallery", "layout_signature": "three_columns"},
            {"slide_index": 7, "role": "gallery", "layout_signature": "image_grid"},
            {"slide_index": 10, "role": "ending", "layout_signature": "ending"},
        ],
    }), encoding="utf-8")
    semantic_path = Path(__file__).resolve().parents[1] / "assets" / "template_specs" / "T01_green_research.semantic.json"
    plan = PagePlan("组会-文献精读", ("研究内容",), pages, deck_scope="complete")

    payload = DynamicCompositionCompiler().compile(
        plan,
        package,
        visual_system_path=tmp_path / "visual.json",
        template_grammar_path=grammar_path,
        template_semantic_spec_path=semantic_path,
        asset_base_dir=tmp_path,
        confirmed=True,
    )

    assert payload["pages"][1]["layout"] == "media_gallery"
    assert len(payload["pages"][1]["media_items"]) == 4
    assert payload["pages"][1]["template_reference"]["selection_source"] == "standard_template_specification"
    assert payload["pages"][2]["layout"] == "module_media"
    assert len(payload["pages"][2]["modules"]) == 4
    assert all(module["media_source"]["semantic_use"] for module in payload["pages"][2]["modules"])

    visual_path = tmp_path / "visual.json"
    visual_path.write_text(json.dumps({
        "colors": {
            "background": "#F7FAF8", "surface": "#FFFFFF", "primary": "#176B4D",
            "primary_soft": "#CDE8DB", "primary_pale": "#EDF7F2", "text": "#1F2933",
            "muted": "#667085", "line": "#D9E2DD",
        },
        "fonts": {"title": "Microsoft YaHei", "body": "Microsoft YaHei", "latin": "Arial"},
        "typography": {
            "cover_title": 30, "cover_subtitle": 17, "page_title": 22, "section_nav": 10,
            "panel_title": 14, "body": 12, "caption": 9, "footer": 8,
        },
    }), encoding="utf-8")
    plan_path = tmp_path / "dynamic_plan.json"
    plan_path.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")
    output = tmp_path / "media_deck.pptx"
    object_manifest = tmp_path / "object_bindings.json"

    render_dynamic(plan_path, visual_path, output, object_manifest_path=object_manifest)

    rendered = Presentation(output)
    assert len(rendered.slides) == 4
    assert sum(1 for shape in rendered.slides[1].shapes if shape.shape_type == 13) == 4
    assert sum(1 for shape in rendered.slides[2].shapes if shape.shape_type == 13) == 4
    binding_result = ObjectLevelQualityGate().inspect(
        output,
        semantic_path,
        render_manifest=object_manifest,
        asset_root=tmp_path,
    )
    assert binding_result.passed, binding_result.errors


def test_image_scaffold_reconstruction_has_no_nested_old_and_new_frames(tmp_path: Path):
    image = tmp_path / "evidence.png"
    Image.new("RGB", (640, 360), "white").save(image)
    old_frame = tmp_path / "old-content-frame.png"
    Image.new("RGB", (1280, 720), "red").save(old_frame)
    identity_strip = tmp_path / "identity-strip.png"
    Image.new("RGB", (1280, 32), "green").save(identity_strip)
    old_navigation_state = tmp_path / "old-navigation-state.png"
    Image.new("RGB", (240, 70), "darkgreen").save(old_navigation_state)
    grammar_path = tmp_path / "grammar.json"
    grammar_path.write_text(json.dumps({
        "identity": {"panel_mode": "image_scaffold", "navigation": "top"},
        "geometry": {
            "title_box": {"left": 0.05, "top": 0.14, "width": 0.9, "height": 0.05},
            "content_region": {"left": 0.05, "top": 0.22, "width": 0.9, "height": 0.64},
        },
        "archetypes": [{
            "slide_index": 8,
            "role": "text_figure",
            "layout_signature": "figure_text_right",
            "picture_slots": [{
                "shape_id": 42,
                "role": "content_image",
                "area": 0.22,
                "box": {"left": 0.55, "top": 0.25, "width": 0.38, "height": 0.5},
            }],
            "decorative_assets": [
                {
                    "path": str(old_frame),
                    "area": 0.65,
                    "box": {"left": 0.04, "top": 0.22, "width": 0.92, "height": 0.71},
                },
                {
                    "path": str(identity_strip),
                    "area": 0.02,
                    "box": {"left": 0.05, "top": 0.03, "width": 0.9, "height": 0.025},
                },
                {
                    "shape_id": 7,
                    "role": "navigation",
                    "path": str(old_navigation_state),
                    "area": 0.009,
                    "box": {"left": 0.3, "top": 0.05, "width": 0.12, "height": 0.075},
                },
            ],
        }],
    }), encoding="utf-8")
    visual_path = tmp_path / "visual.json"
    visual_path.write_text(json.dumps({
        "colors": {
            "background": "#F7FAF8", "surface": "#FFFFFF", "primary": "#176B4D",
            "primary_soft": "#CDE8DB", "primary_pale": "#EDF7F2", "text": "#1F2933",
            "muted": "#667085", "line": "#D9E2DD",
        },
        "fonts": {"title": "Microsoft YaHei", "body": "Microsoft YaHei", "latin": "Arial"},
        "typography": {
            "cover_title": 30, "cover_subtitle": 17, "page_title": 22, "section_nav": 10,
            "panel_title": 14, "body": 12, "caption": 9, "footer": 8,
        },
    }), encoding="utf-8")
    plan_path = tmp_path / "plan.json"
    plan_path.write_text(json.dumps({
        "schema_version": 2,
        "confirmed": True,
        "sections": ["研究结果"],
        "template_grammar": str(grammar_path),
        "asset_base_dir": str(tmp_path),
        "pages": [{
            "page_id": "P001",
            "section": "研究结果",
            "layout": "text_figure",
            "title": "图像证据支持核心结论",
            "lead": "处理组呈现稳定改善",
            "bullets": [
                {"title": "形态变化", "body": "处理后边界更清晰"},
                {"title": "对照关系", "body": "同批对照排除批次差异"},
                {"title": "适用边界", "body": "仍需外部样本验证"},
            ],
            "image": str(image),
            "caption": "来源图像 | PDF 第 8 页",
            "media_source": {
                "source_id": "FIG-8", "source_path": str(image),
                "pdf_page": "not_applicable", "semantic_use": "支持处理效果判断",
                "asset_path": str(image),
            },
            "scaffold_slide": 8,
            "use_template_scaffold": "structure",
        }],
    }, ensure_ascii=False), encoding="utf-8")
    output = tmp_path / "single-layer.pptx"

    render_dynamic(plan_path, visual_path, output)

    slide = Presentation(output).slides[0]
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    # Keep the evidence and identity strip; exclude both the old content raster and
    # the source page's active navigation backing.
    assert len(pictures) == 2
    frames = [
        shape for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and not getattr(shape, "text", "").strip()
        and shape.width >= Inches(1)
        and shape.height >= Inches(0.55)
    ]
    nested = []
    for outer in frames:
        for inner in frames:
            if outer.shape_id == inner.shape_id:
                continue
            contains = (
                outer.left <= inner.left
                and outer.top <= inner.top
                and outer.left + outer.width >= inner.left + inner.width
                and outer.top + outer.height >= inner.top + inner.height
            )
            if contains and outer.width * outer.height > inner.width * inner.height * 1.1:
                nested.append((outer.shape_id, inner.shape_id))
    assert not nested


def test_semantic_process_uses_module_headers_without_legacy_step_circles(tmp_path: Path):
    grammar_path = tmp_path / "grammar.json"
    grammar_path.write_text(json.dumps({
        "identity": {"panel_mode": "image_scaffold", "navigation": "top"},
        "geometry": {},
        "archetypes": [{
            "slide_index": 7,
            "role": "process",
            "layout_signature": "five_columns",
            "decorative_assets": [],
        }],
    }), encoding="utf-8")
    semantic_path = tmp_path / "semantic.json"
    semantic_path.write_text(json.dumps({
        "template_id": "T03",
        "pages": [{
            "page_id": "T03_P07",
            "source_slide_index": 7,
            "semantic_modules": [
                {
                    "module_id": f"STEP_{index}",
                    "role": "process_step",
                    "semantic_region": "content",
                    "box": {
                        "left": 0.05 + (index - 1) * 0.19,
                        "top": 0.25,
                        "width": 0.17,
                        "height": 0.6,
                    },
                    "child_slots": [],
                }
                for index in range(1, 6)
            ],
        }],
    }), encoding="utf-8")
    visual_path = tmp_path / "visual.json"
    visual_path.write_text(json.dumps({
        "colors": {
            "background": "#F7FAF8", "surface": "#FFFFFF", "primary": "#176B4D",
            "primary_soft": "#CDE8DB", "primary_pale": "#EDF7F2", "text": "#1F2933",
            "muted": "#667085", "line": "#D9E2DD",
        },
        "fonts": {"title": "Arial", "body": "Arial", "latin": "Arial"},
        "typography": {
            "cover_title": 30, "cover_subtitle": 17, "page_title": 22, "section_nav": 10,
            "panel_title": 14, "body": 12, "caption": 9, "footer": 8,
        },
    }), encoding="utf-8")
    plan_path = tmp_path / "plan.json"
    plan_path.write_text(json.dumps({
        "schema_version": 2,
        "confirmed": True,
        "sections": ["Results"],
        "template_grammar": str(grammar_path),
        "template_semantic_spec": str(semantic_path),
        "asset_base_dir": str(tmp_path),
        "pages": [{
            "page_id": "P008",
            "section": "Results",
            "layout": "process",
            "title": "A five-stage validation process",
            "steps": [
                {"title": f"Stage {index}", "body": "Source-grounded explanation"}
                for index in range(1, 6)
            ],
            "template_reference": {"semantic_spec_page_id": "T03_P07"},
            "scaffold_slide": 7,
            "use_template_scaffold": "structure",
        }],
    }), encoding="utf-8")
    output = tmp_path / "semantic-process.pptx"

    render_dynamic(plan_path, visual_path, output)

    slide = Presentation(output).slides[0]
    legacy_circles = [
        shape for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and shape.auto_shape_type == MSO_SHAPE.OVAL
    ]
    cards = [
        shape for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and shape.height > Inches(1)
        and shape.width > Inches(1.5)
    ]
    assert not legacy_circles
    assert len(cards) == 5
    assert all(shape.height <= Inches(3.2) for shape in cards)
    assert sum("Stage " in getattr(shape, "text", "") for shape in slide.shapes) == 5
