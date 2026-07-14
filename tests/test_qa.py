import tempfile
import unittest
import json
from unittest.mock import patch
from pathlib import Path

from pptx import Presentation

from academic_ppt.qa import RenderQualityGate
from academic_ppt.qa import ScientificSemanticGate
from academic_ppt.evidence import EvidenceGraph
from academic_ppt.ingest import SourceIngestor
from academic_ppt.planning import PageDraft, PagePlanner
from academic_ppt.visuals import VisualTask
from academic_ppt.qa import VisualCompositionGate
from academic_ppt.qa import QualitySummaryBuilder, RenderQualityResult, SemanticQualityResult, VisualQualityResult


class RenderQualityGateTests(unittest.TestCase):
    def test_structural_validation_without_real_render_is_not_formal_acceptance(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            pptx = root / "deck.pptx"
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(pptx)

            result = RenderQualityGate().inspect(pptx, root / "audit", require_runtime=False)

            self.assertTrue(result.structural_passed)
            self.assertFalse(result.real_render_passed)
            self.assertFalse(result.formal_accepted)
            self.assertTrue(result.report_path.is_file())

    def test_required_powerpoint_runtime_reports_unavailable_as_non_acceptance(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            pptx = root / "deck.pptx"
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(pptx)

            commands = []

            def write_unavailable_report(command, **_kwargs):
                commands.append(command)
                report_path = Path(command[command.index("--output") + 1])
                report_path.write_text(json.dumps({
                    "checks": [{
                        "name": "Real application render succeeds",
                        "level": "ERROR",
                        "passed": False,
                        "detail": "PowerPoint is not installed",
                    }],
                    "failed_error": 1,
                }), encoding="utf-8")

            with patch("academic_ppt.qa.subprocess.run", side_effect=write_unavailable_report):
                result = RenderQualityGate().inspect(pptx, root / "audit", require_runtime=True)

            self.assertIn("--render-engine", commands[0])
            self.assertEqual(commands[0][commands[0].index("--render-engine") + 1], "powerpoint")
            self.assertIn("--preview-dir", commands[0])
            self.assertFalse(result.real_render_passed)
            self.assertFalse(result.formal_accepted)
            self.assertEqual(result.runtime_status, "unavailable")

    def test_semantic_gate_rejects_page_plan_while_blocking_evidence_conflict_exists(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "result.md"
            source.write_text("Model F1 reached 92.9%.", encoding="utf-8")
            graph = EvidenceGraph.from_sources([SourceIngestor().ingest(source)])
            evidence = graph.all()[0]
            graph.add_derived_evidence(
                "System-derived F1 is 90.1%.", [evidence.evidence_id], metric="f1", value=90.1,
                authorization_note="Authorized recomputation.",
            )
            claim = graph.add_claim("Model F1 reached 92.9%.", [evidence.evidence_id])
            plan = PagePlanner().build("毕业答辩", ["实验与结果"], [PageDraft(
                page_id="P001", section="实验与结果", title="F1 result",
                question_answered="Does the result support the model?", claim_id=claim.claim_id,
                interpretation="Interpret with the reported evidence.", next_link="Discuss the boundary.",
                time_seconds=60, visual_strategy="native_chart", component_requirements={"text": 1, "chart": 1},
            )], graph)

            result = ScientificSemanticGate().inspect(plan, graph)

            self.assertFalse(result.passed)
            self.assertIn("blocking evidence conflicts", result.errors[0])

    def test_semantic_gate_keeps_single_evidence_density_as_an_advisory_observation(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "result.md"
            source.write_text("Model F1 reached 92.9%.", encoding="utf-8")
            graph = EvidenceGraph.from_sources([SourceIngestor().ingest(source)])
            evidence = graph.all()[0]
            claim = graph.add_claim("Model F1 reached 92.9%.", [evidence.evidence_id])
            plan = PagePlanner().build("Graduation defense", ["Results"], [PageDraft(
                page_id="P001", section="Results", title="F1 result",
                question_answered="Does the result support the model?", claim_id=claim.claim_id,
                interpretation="Interpret with the reported evidence.", next_link="Discuss the boundary.",
                time_seconds=60, visual_strategy="native_chart", component_requirements={"text": 1, "chart": 1},
            )], graph)

            result = ScientificSemanticGate().inspect(plan, graph)

            self.assertTrue(result.passed)
            self.assertEqual(result.errors, ())
            self.assertEqual(result.observations, ("P001: only one evidence unit supports the page claim",))

    def test_visual_composition_gate_rejects_unaccepted_required_visual_task(self):
        task = VisualTask.create(
            task_id="VIS_001", page_id="P001", evidence_ids=["EVD_001"], task_type="chart"
        )

        result = VisualCompositionGate().inspect([task])

        self.assertFalse(result.passed)
        self.assertEqual(result.errors, ("VIS_001: visual task is planned, not accepted",))

    def test_visual_composition_gate_rejects_evidence_not_bound_to_the_page_contract(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source = root / "result.md"
            source.write_text("Model F1 reached 92.9%.", encoding="utf-8")
            graph = EvidenceGraph.from_sources([SourceIngestor().ingest(source)])
            evidence = graph.all()[0]
            claim = graph.add_claim("Model F1 reached 92.9%.", [evidence.evidence_id])
            plan = PagePlanner().build("Graduation defense", ["Results"], [PageDraft(
                page_id="P001", section="Results", title="F1 result",
                question_answered="Does the result support the model?", claim_id=claim.claim_id,
                interpretation="Interpret with the reported evidence.", next_link="Discuss the boundary.",
                time_seconds=60, visual_strategy="native_chart", component_requirements={"text": 1, "chart": 1},
            )], graph)
            task = VisualTask(
                task_id="VIS_001", page_id="P001", evidence_ids=("EVD_UNBOUND",), task_type="chart", status="accepted"
            )

            result = VisualCompositionGate().inspect([task], plan)

            self.assertFalse(result.passed)
            self.assertEqual(result.errors, ("VIS_001: visual evidence is not bound to the page contract: EVD_UNBOUND",))

    def test_quality_summary_requires_all_hard_gates_for_formal_acceptance(self):
        summary = QualitySummaryBuilder().build(
            RenderQualityResult(Path("report.json"), True, False, False),
            SemanticQualityResult(True, ()),
            VisualQualityResult(True, ()),
        )

        self.assertEqual(summary["structural"], True)
        self.assertEqual(summary["semantic"], True)
        self.assertEqual(summary["visual"], True)
        self.assertEqual(summary["real_render"], False)
        self.assertFalse(summary["formal_accepted"])

    def test_quality_summary_rejects_a_non_powerpoint_render_for_formal_acceptance(self):
        summary = QualitySummaryBuilder().build(
            RenderQualityResult(Path("report.json"), True, True, False, "passed"),
            SemanticQualityResult(True, ()),
            VisualQualityResult(True, ()),
        )

        self.assertTrue(summary["real_render"])
        self.assertFalse(summary["formal_accepted"])


if __name__ == "__main__":
    unittest.main()
