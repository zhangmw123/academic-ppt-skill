import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from PIL import Image

from academic_ppt.evidence import EvidenceGraph
from academic_ppt.ingest import SourceIngestor
from academic_ppt.layout import LayoutCompiler
from academic_ppt.planning import PageDraft, PagePlanner
from academic_ppt.rendering import NativeRenderAdapter
from academic_ppt.templates import TemplateCapabilityGraph


class NativeRenderAdapterTests(unittest.TestCase):
    def test_renders_page_plan_into_cloned_template_text_shape_without_overlay(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source_path = root / "result.md"
            source_path.write_text("Model F1 reached 92.9%.", encoding="utf-8")
            graph = EvidenceGraph.from_sources([SourceIngestor().ingest(source_path)])
            evidence = graph.all()[0]
            claim = graph.add_claim("Model F1 reached 92.9%.", [evidence.evidence_id])
            page_plan = PagePlanner().build(
                "毕业答辩",
                ["实验与结果"],
                [PageDraft(
                    page_id="P001",
                    section="实验与结果",
                    title="Model F1 reached 92.9%.",
                    question_answered="Does the experiment support effectiveness?",
                    claim_id=claim.claim_id,
                    interpretation="The result supports effectiveness within the dataset boundary.",
                    next_link="Next, explain the evaluation boundary.",
                    time_seconds=75,
                    visual_strategy="text_only",
                    component_requirements={"text": 1},
                )],
                graph,
            )
            template_path = root / "template.pptx"
            template = Presentation()
            source_slide = template.slides.add_slide(template.slide_layouts[6])
            text_shape = source_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1))
            text_shape.text = "SAMPLE TITLE"
            template.save(template_path)
            decision = LayoutCompiler(TemplateCapabilityGraph.from_presentation(template_path)).compile(
                page_plan.pages[0].contract
            )
            adapter = NativeRenderAdapter()
            layout_path = adapter.write_layout_plan(
                page_plan,
                {"P001": decision},
                {"P001": ["Model F1 reached 92.9%."]},
                root / "layout_plan.json",
                confirmed=True,
            )
            output = root / "deck.pptx"

            adapter.render(layout_path, template_path, output)

            rendered = Presentation(output)
            self.assertEqual(len(rendered.slides), 1)
            self.assertEqual(len(rendered.slides[0].shapes), len(source_slide.shapes))
            self.assertEqual(rendered.slides[0].shapes[0].shape_id, text_shape.shape_id)
            self.assertEqual(rendered.slides[0].shapes[0].text, "Model F1 reached 92.9%.")

    def test_replaces_source_picture_through_original_template_picture_shape_id(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source_path = root / "result.md"
            source_path.write_text("Model F1 reached 92.9%.", encoding="utf-8")
            graph = EvidenceGraph.from_sources([SourceIngestor().ingest(source_path)])
            evidence = graph.all()[0]
            claim = graph.add_claim("Model F1 reached 92.9%.", [evidence.evidence_id])
            page_plan = PagePlanner().build(
                "毕业答辩",
                ["实验与结果"],
                [PageDraft(
                    page_id="P001", section="实验与结果", title="Model F1 reached 92.9%.",
                    question_answered="Does the experiment support effectiveness?", claim_id=claim.claim_id,
                    interpretation="The result supports effectiveness within the dataset boundary.",
                    next_link="Next, explain the evaluation boundary.", time_seconds=75,
                    visual_strategy="source_figure", component_requirements={"text": 1, "picture": 1},
                )],
                graph,
            )
            sample_image = root / "sample.png"
            replacement_image = root / "replacement.png"
            Image.new("RGB", (120, 80), "red").save(sample_image)
            Image.new("RGB", (120, 80), "blue").save(replacement_image)
            template_path = root / "template.pptx"
            template = Presentation()
            source_slide = template.slides.add_slide(template.slide_layouts[6])
            source_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(7), Inches(1)).text = "SAMPLE TITLE"
            picture = source_slide.shapes.add_picture(str(sample_image), Inches(1), Inches(2), Inches(4), Inches(2))
            template.save(template_path)
            decision = LayoutCompiler(TemplateCapabilityGraph.from_presentation(template_path)).compile(
                page_plan.pages[0].contract
            )
            adapter = NativeRenderAdapter()
            layout_path = adapter.write_layout_plan(
                page_plan,
                {"P001": decision},
                {"P001": ["Model F1 reached 92.9%."]},
                root / "layout_plan.json",
                confirmed=True,
                image_content={"P001": [replacement_image]},
            )
            output = adapter.render(layout_path, template_path, root / "deck.pptx")

            rendered = Presentation(output)
            rendered_picture = next(shape for shape in rendered.slides[0].shapes if shape.shape_id == picture.shape_id)
            self.assertEqual(rendered_picture.shape_id, picture.shape_id)
            self.assertEqual(rendered_picture.image.blob, replacement_image.read_bytes())


if __name__ == "__main__":
    unittest.main()
