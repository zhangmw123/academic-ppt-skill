import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from academic_ppt.layout import LayoutCompiler, ScientificPageContract
from academic_ppt.templates import TemplateAdmissionGate, TemplateCapabilityGraph, TemplateCatalog


class TemplateCapabilityGraphTests(unittest.TestCase):
    def test_resolves_original_bundled_names_to_disclosed_powerpoint_safe_templates(self):
        catalog = TemplateCatalog.load()

        green = catalog.select("组会-文献精读", "绿色-科研风格PPT模版.pptx")
        blue_defense = catalog.select(
            "毕业答辩",
            "蓝色-学术答辩多版式通用模板 (Academic Defense Multi-Layout Template).pptx",
        )

        self.assertEqual(green.template_id, "T01")
        self.assertEqual(blue_defense.template_id, "T03")
        self.assertTrue(Path(green.path).is_file())
        self.assertTrue(Path(blue_defense.path).is_file())
        self.assertEqual(green.support_level, "bundled_recompiled_source")
        self.assertIn("complete slide structure", green.substitution_reason)
        self.assertTrue(green.source_path.endswith("绿色-科研风格PPT模版.pptx"))
        self.assertEqual(green.source_fidelity, "complete_structure_recompile")
        self.assertIn("all 10 source slides", green.source_limitations)
        self.assertIn("all 11 source slides", blue_defense.source_limitations)

    def test_accepts_an_existing_unregistered_template_as_conditional(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            template_path = Path(temp_dir) / "custom.pptx"
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(template_path)

            selection = TemplateCatalog.load().select("毕业答辩", template_path)

            self.assertIsNone(selection.template_id)
            self.assertEqual(selection.support_level, "conditional_user")
            self.assertEqual(Path(selection.path), template_path.resolve())

    def test_user_template_admission_requires_editable_components_and_grammar(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            valid_path = root / "valid.pptx"
            valid = Presentation()
            valid_slide = valid.slides.add_slide(valid.slide_layouts[6])
            valid_slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "Editable"
            valid.save(valid_path)
            invalid_path = root / "invalid.pptx"
            invalid = Presentation()
            invalid.slides.add_slide(invalid.slide_layouts[6])
            invalid.save(invalid_path)

            accepted = TemplateAdmissionGate().inspect(valid_path, require_runtime=False)
            rejected = TemplateAdmissionGate().inspect(invalid_path, require_runtime=False)

            self.assertTrue(accepted.passed)
            self.assertTrue(accepted.grammar_extracted)
            self.assertGreater(accepted.editable_component_count, 0)
            self.assertFalse(rejected.passed)
            self.assertTrue(any("editable" in error for error in rejected.errors))

    def test_selects_a_native_slide_with_required_editable_component_capacity(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            image_path = root / "figure.png"
            Image.new("RGB", (320, 180), "white").save(image_path)
            template_path = root / "template.pptx"
            presentation = Presentation()
            text_only = presentation.slides.add_slide(presentation.slide_layouts[6])
            text_only.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "标题"
            figure_page = presentation.slides.add_slide(presentation.slide_layouts[6])
            figure_page.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "图表标题"
            figure_page.shapes.add_picture(str(image_path), Inches(1), Inches(2), Inches(4), Inches(2))
            presentation.save(template_path)

            graph = TemplateCapabilityGraph.from_presentation(template_path)
            candidates = graph.find_compatible_slides({"text": 1, "picture": 1})

            self.assertEqual(graph.slides[1].slide_index, 2)
            self.assertEqual(candidates, (graph.slides[1],))
            self.assertEqual(graph.slides[1].component_counts, {"picture": 1, "text": 1})

    def test_layout_compiler_uses_compatible_native_template_slide_for_page_contract(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            image_path = root / "figure.png"
            Image.new("RGB", (320, 180), "white").save(image_path)
            template_path = root / "template.pptx"
            presentation = Presentation()
            text_only = presentation.slides.add_slide(presentation.slide_layouts[6])
            text_only.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "标题"
            figure_page = presentation.slides.add_slide(presentation.slide_layouts[6])
            figure_page.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "图表标题"
            figure_page.shapes.add_picture(str(image_path), Inches(1), Inches(2), Inches(4), Inches(2))
            presentation.save(template_path)
            graph = TemplateCapabilityGraph.from_presentation(template_path)
            contract = ScientificPageContract(
                page_id="P001",
                claim_id="CLM_001",
                component_requirements={"text": 1, "picture": 1},
            )

            decision = LayoutCompiler(graph).compile(contract)

            self.assertEqual(decision.render_mode, "template_native")
            self.assertEqual(decision.source_slide_index, 2)
            self.assertIsNone(decision.fallback_reason)

    def test_layout_decision_binds_required_content_to_original_template_shape_ids(self):
        with tempfile.TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            image_path = root / "figure.png"
            Image.new("RGB", (320, 180), "white").save(image_path)
            template_path = root / "template.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "标题"
            slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), Inches(4), Inches(2))
            presentation.save(template_path)
            graph = TemplateCapabilityGraph.from_presentation(template_path)
            contract = ScientificPageContract(
                page_id="P001",
                claim_id="CLM_001",
                component_requirements={"text": 1, "picture": 1},
            )

            decision = LayoutCompiler(graph).compile(contract)

            expected = {
                kind: tuple(component.shape_id for component in graph.slides[0].components if component.kind == kind)
                for kind in ("text", "picture")
            }
            self.assertEqual(decision.component_bindings, expected)


if __name__ == "__main__":
    unittest.main()
